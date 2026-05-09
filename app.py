
import os
import re
import io
import base64
import sqlite3
import inspect
import hashlib
from datetime import date, datetime, timedelta
from pathlib import Path

import pandas as pd
import numpy as np
import plotly.express as px
import streamlit as st

try:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import landscape, letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, Image, PageBreak
    REPORTLAB_AVAILABLE = True
except Exception:
    REPORTLAB_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
    PPTX_AVAILABLE = True
except Exception:
    PPTX_AVAILABLE = False

APP_TITLE = "Inventory App"
APP_VERSION = "1.2.3"  # Removed default account note from login
APP_DIR = Path(__file__).resolve().parent
DB_PATH = APP_DIR / "inventory_app.db"
UPLOAD_DIR = APP_DIR / "uploads"
UPLOAD_DIR.mkdir(exist_ok=True)
ASSET_DIR = APP_DIR / "assets"
LOGO_PATH = ASSET_DIR / "farmfix_logo.png"
APP_LOGO_PATH = ASSET_DIR / "app_logo.png"


def image_to_data_uri(path):
    try:
        p = Path(path)
        if p.exists():
            encoded = base64.b64encode(p.read_bytes()).decode("utf-8")
            suffix = p.suffix.lower().replace(".", "") or "png"
            if suffix == "jpg":
                suffix = "jpeg"
            return f"data:image/{suffix};base64,{encoded}"
    except Exception:
        pass
    return ""

st.set_page_config(
    page_title=APP_TITLE,
    page_icon="📦",
    layout="wide",
    initial_sidebar_state="expanded",
)

# -----------------------------
# CSS / UI DESIGN
# -----------------------------
st.markdown(
    """
    <style>
    :root {
        --green: #1f7a4d;
        --green-dark: #155f3a;
        --green-light: #eaf7f0;
        --text: #17312a;
        --muted: #667085;
        --bg: #f7faf8;
        --card: #ffffff;
        --border: #e5e7eb;
        --red: #dc2626;
        --orange: #f59e0b;
        --blue: #2563eb;
    }

    .main .block-container,
    div[data-testid="stAppViewBlockContainer"],
    [data-testid="stAppViewContainer"] .main .block-container {
        padding-top: 0.85rem !important;
        margin-top: 0rem !important;
        padding-bottom: 2rem;
        max-width: 1500px;
    }

    header[data-testid="stHeader"] {
        height: 2.2rem !important;
        min-height: 2.2rem !important;
        max-height: 2.2rem !important;
        background: transparent !important;
        pointer-events: auto !important;
        overflow: visible !important;
    }
    header[data-testid="stHeader"] * {
        pointer-events: auto !important;
    }
    /* Do not hide stToolbar because Streamlit places the sidebar show button there after collapse. */
    div[data-testid="stToolbar"] {
        display: flex !important;
        visibility: visible !important;
        opacity: 1 !important;
        pointer-events: auto !important;
    }
    div[data-testid="stDecoration"],
    #MainMenu,
    footer {
        display: none !important;
        visibility: hidden !important;
        height: 0rem !important;
        min-height: 0rem !important;
        max-height: 0rem !important;
    }

    [data-testid="stAppViewContainer"] {
        padding-top: 0rem !important;
        margin-top: 0rem !important;
    }

    section.main,
    section[data-testid="stAppViewContainer"] {
        padding-top: 0rem !important;
        margin-top: 0rem !important;
    }

    .app-header {
        display: flex;
        align-items: center;
        justify-content: space-between;
        gap: 1rem;
        padding: 1rem 1.25rem;
        background: linear-gradient(135deg, #ffffff 0%, #eefaf2 100%);
        border: 1px solid var(--border);
        border-radius: 22px;
        margin-bottom: 1rem;
        box-shadow: 0 8px 26px rgba(16, 24, 40, 0.06);
    }

    .app-title {
        font-size: 1.8rem;
        line-height: 1.1;
        font-weight: 800;
        color: var(--text);
        margin: 0;
    }

    .app-subtitle {
        color: var(--muted);
        margin-top: .25rem;
        font-size: .92rem;
    }

    .soft-card {
        background: var(--card);
        border: 1px solid var(--border);
        border-radius: 20px;
        padding: 1.1rem;
        box-shadow: 0 8px 24px rgba(16, 24, 40, 0.055);
    }

    .module-hero-card {
        display:flex; align-items:center; justify-content:space-between; gap:1rem;
        padding:1.1rem 1.2rem; margin-bottom:1rem;
        background:linear-gradient(135deg,#ffffff 0%,#eefaf4 58%,#e3f4ea 100%);
        border:1px solid #deebe2; border-radius:22px;
        box-shadow:0 8px 26px rgba(16,24,40,.05);
    }
    .module-hero-main { display:flex; align-items:center; gap:.95rem; }
    .module-hero-icon { width:56px; height:56px; border-radius:18px; display:grid; place-items:center; font-size:1.55rem; background:rgba(255,255,255,.92); border:1px solid #e6efe9; box-shadow:0 8px 22px rgba(16,24,40,.05); }
    .module-hero-title { font-size:1.18rem; font-weight:900; color:#17312a; line-height:1.2; margin-bottom:.15rem; }
    .module-hero-subtitle { font-size:.9rem; color:#667085; line-height:1.45; max-width:760px; }
    .module-hero-tag { display:inline-flex; align-items:center; gap:.35rem; padding:.45rem .7rem; border-radius:999px; background:rgba(255,255,255,.86); border:1px solid #e6efe9; color:#246246; font-size:.8rem; font-weight:800; white-space:nowrap; }
    .section-shell {
        background:#ffffff; border:1px solid #e5e7eb; border-radius:22px; box-shadow:0 8px 24px rgba(16,24,40,.045);
        padding:1rem 1rem .9rem 1rem; margin-bottom:1rem;
    }
    .section-shell-title { font-size:1.02rem; font-weight:900; color:#17312a; margin-bottom:.15rem; }
    .section-shell-subtitle { color:#667085; font-size:.84rem; margin-bottom:.8rem; }
    .chart-note { color:#98a2b3; font-size:.76rem; margin-top:.25rem; }
    .table-toolbar { display:flex; align-items:center; justify-content:space-between; gap:1rem; padding:.35rem 0 .65rem 0; }
    .table-toolbar-note { color:#98a2b3; font-size:.78rem; font-weight:600; }
    .custom-table-help { color:#98a2b3; font-size:.76rem; margin-top:-.25rem; margin-bottom:.45rem; }
    .report-summary-card { background:#ffffff; border:1px solid #e6e9ef; border-radius:18px; padding:.95rem 1rem; box-shadow:0 10px 26px rgba(16,24,40,.04); min-height:110px; margin-bottom:.35rem; }
    .report-summary-label { color:#667085; font-size:.82rem; font-weight:800; margin-bottom:.45rem; }
    .report-summary-value { color:#17312a; font-size:1.9rem; line-height:1; font-weight:900; margin-bottom:.45rem; }
    .report-summary-note { color:#98a2b3; font-size:.78rem; font-weight:600; }

    .status-pill {
        display: inline-flex;
        align-items: center;
        gap: 6px;
        padding: 5px 10px;
        border-radius: 999px;
        font-size: 12px;
        font-weight: 700;
    }

    .pill-available { background:#dcfce7; color:#166534; }
    .pill-expiring { background:#fef3c7; color:#92400e; }
    .pill-expired { background:#fee2e2; color:#991b1b; }
    .pill-out { background:#f3f4f6; color:#374151; }
    .pill-info { background:#dbeafe; color:#1e40af; }

    div[data-testid="stMetric"] {
        background: #ffffff;
        border: 1px solid #e5e7eb;
        padding: 1rem;
        border-radius: 18px;
        box-shadow: 0 8px 22px rgba(16, 24, 40, 0.05);
    }

    div.stButton > button {
        border-radius: 18px;
        border: 1px solid #e5e7eb;
        background: #ffffff;
        color: #17312a;
        box-shadow: 0 8px 22px rgba(16, 24, 40, 0.045);
        min-height: 104px;
        font-weight: 700;
        white-space: pre-line;
        transition: all 0.15s ease;
    }

    div.stButton > button:hover {
        border-color: #22c55e;
        transform: translateY(-1px);
        box-shadow: 0 12px 26px rgba(16, 24, 40, 0.08);
    }

    .small-note {
        color: var(--muted);
        font-size: .86rem;
    }

    .section-title {
        font-size: 1.2rem;
        font-weight: 800;
        color: #17312a;
        margin: .35rem 0 .75rem 0;
    }

    .kpi-anchor {
        text-decoration: none !important;
        display: block;
        color: inherit !important;
    }
    .kpi-card {
        position: relative;
        background:
            radial-gradient(circle at 86% 78%, var(--accent-glow, rgba(34,197,94,.12)) 0 16%, transparent 17%),
            linear-gradient(180deg, #ffffff 0%, #fbfdfc 100%);
        border: 1px solid #e5ebf1;
        border-radius: 24px;
        padding: 1.15rem 1.1rem 1.1rem 1.1rem;
        min-height: 205px;
        box-shadow: 0 12px 32px rgba(16,24,40,.06);
        overflow: hidden;
        transition: transform .18s ease, box-shadow .18s ease, border-color .18s ease, background .18s ease;
    }
    .kpi-anchor:hover .kpi-card {
        transform: translateY(-4px);
        box-shadow: 0 20px 42px rgba(16,24,40,.11);
        border-color: color-mix(in srgb, var(--accent, #22c55e) 34%, #e5ebf1);
    }
    .kpi-card::before {
        content: "";
        position: absolute;
        top: 0;
        left: 0;
        right: 0;
        height: 5px;
        background: linear-gradient(90deg, var(--accent, #22c55e), color-mix(in srgb, var(--accent, #22c55e) 55%, white));
    }
    .kpi-icon {
        width: 54px;
        height: 54px;
        border-radius: 999px;
        display: grid;
        place-items: center;
        font-size: 1.55rem;
        background:
            radial-gradient(circle at 30% 25%, rgba(255,255,255,.95) 0 20%, transparent 21%),
            var(--accent-soft, #ecfdf3);
        color: var(--accent, #22c55e);
        margin-bottom: 1rem;
        border: 1px solid color-mix(in srgb, var(--accent, #22c55e) 18%, white);
        box-shadow: 0 9px 22px var(--accent-shadow, rgba(34,197,94,.16));
    }
    .kpi-title {
        font-size: 1rem;
        font-weight: 850;
        color: #17312a;
        line-height: 1.2;
        min-height: 2.15rem;
        margin-bottom: .65rem;
        letter-spacing: -.015em;
    }
    .kpi-value {
        font-size: 3.55rem;
        line-height: .9;
        font-weight: 950;
        color: var(--accent, #22c55e);
        margin-bottom: .58rem;
        letter-spacing: -0.055em;
    }
    .kpi-note {
        font-size: .72rem;
        line-height: 1.28;
        color: #aeb7c3;
        font-weight: 600;
        max-width: 88%;
    }
    .kpi-chevron {
        position: absolute;
        right: 1rem;
        bottom: .95rem;
        width: 28px;
        height: 28px;
        border-radius: 999px;
        display: grid;
        place-items: center;
        font-size: 1.25rem;
        color: var(--accent, #22c55e);
        background: var(--accent-soft, #ecfdf3);
        font-weight: 800;
        opacity: .92;
    }
    .kpi-available { --accent: #22c55e; --accent-soft: #eaf9f0; --accent-shadow: rgba(34,197,94,.17); --accent-glow: rgba(34,197,94,.11); }
    .kpi-expiring { --accent: #f59e0b; --accent-soft: #fff6e8; --accent-shadow: rgba(245,158,11,.17); --accent-glow: rgba(245,158,11,.11); }
    .kpi-expired { --accent: #ef4444; --accent-soft: #feeeee; --accent-shadow: rgba(239,68,68,.17); --accent-glow: rgba(239,68,68,.10); }
    .kpi-received { --accent: #2563eb; --accent-soft: #edf4ff; --accent-shadow: rgba(37,99,235,.16); --accent-glow: rgba(37,99,235,.10); }
    .kpi-out { --accent: #7c3aed; --accent-soft: #f3ecff; --accent-shadow: rgba(124,58,237,.16); --accent-glow: rgba(124,58,237,.10); }


    /* Native Streamlit KPI button overlay:
       v1.0 keeps original KPI card design and uses isolated wrapper hover.
       The click button remains fully invisible and never draws its own shape. */
    [class*="st-key-kpi_wrap_"] {
        position: relative !important;
    }

    [class*="st-key-dashboard_kpi_btn_"] {
        margin-top: -205px !important;
        height: 205px !important;
        position: relative !important;
        z-index: 20 !important;
        background: transparent !important;
        border: none !important;
        outline: none !important;
        box-shadow: none !important;
    }

    [class*="st-key-dashboard_kpi_btn_"] button {
        width: 100% !important;
        height: 205px !important;
        min-height: 205px !important;
        border-radius: 24px !important;
        opacity: 0 !important;
        cursor: pointer !important;
        padding: 0 !important;
        margin: 0 !important;
        background: transparent !important;
        border: none !important;
        outline: none !important;
        box-shadow: none !important;
    }

    [class*="st-key-dashboard_kpi_btn_"] button:hover,
    [class*="st-key-dashboard_kpi_btn_"] button:active,
    [class*="st-key-dashboard_kpi_btn_"] button:focus,
    [class*="st-key-dashboard_kpi_btn_"]:hover,
    [class*="st-key-dashboard_kpi_btn_"]:active,
    [class*="st-key-dashboard_kpi_btn_"]:focus-within {
        background: transparent !important;
        border: none !important;
        outline: none !important;
        box-shadow: none !important;
    }

    [class*="st-key-dashboard_kpi_btn_"] button p {
        opacity: 0 !important;
        font-size: 0 !important;
        color: transparent !important;
    }

    .kpi-card,
    .kpi-icon,
    .kpi-chevron {
        transition: transform .22s ease, box-shadow .22s ease, border-color .22s ease, background .22s ease, opacity .22s ease !important;
    }

    /* Hover only the KPI card inside the same keyed wrapper. */
    [class*="st-key-kpi_wrap_"]:has([class*="st-key-dashboard_kpi_btn_"]:hover) .kpi-card,
    [class*="st-key-kpi_wrap_"]:has([class*="st-key-dashboard_kpi_btn_"] button:hover) .kpi-card {
        transform: translateY(-6px) !important;
        box-shadow: 0 22px 46px rgba(16,24,40,.16) !important;
        border-color: color-mix(in srgb, var(--accent, #22c55e) 52%, #e5ebf1) !important;
        background:
            radial-gradient(circle at 86% 78%, var(--accent-glow, rgba(34,197,94,.18)) 0 20%, transparent 21%),
            linear-gradient(180deg, #ffffff 0%, #fbfffd 100%) !important;
    }

    [class*="st-key-kpi_wrap_"]:has([class*="st-key-dashboard_kpi_btn_"]:hover) .kpi-icon,
    [class*="st-key-kpi_wrap_"]:has([class*="st-key-dashboard_kpi_btn_"] button:hover) .kpi-icon {
        transform: scale(1.06) !important;
        box-shadow: 0 13px 28px var(--accent-shadow, rgba(34,197,94,.24)) !important;
    }

    [class*="st-key-kpi_wrap_"]:has([class*="st-key-dashboard_kpi_btn_"]:hover) .kpi-chevron,
    [class*="st-key-kpi_wrap_"]:has([class*="st-key-dashboard_kpi_btn_"] button:hover) .kpi-chevron {
        transform: translateX(4px) !important;
        opacity: 1 !important;
    }

    [class*="st-key-kpi_wrap_"]:has([class*="st-key-dashboard_kpi_btn_"]:active) .kpi-card,
    [class*="st-key-kpi_wrap_"]:has([class*="st-key-dashboard_kpi_btn_"] button:active) .kpi-card {
        transform: translateY(-2px) scale(.992) !important;
    }

    [class*="st-key-dashboard_kpi_btn_"] button:focus {
        outline: none !important;
    }

    @media (max-width: 768px) {
        .app-header { padding: .85rem; border-radius: 16px; }
        .app-title { font-size: 1.35rem; }
        .main .block-container { padding-left: .75rem; padding-right: .75rem; padding-top: .55rem !important; padding-bottom: 1.2rem; }

        .module-hero-card {
            flex-direction: column;
            align-items: flex-start;
            gap: .85rem;
            padding: .95rem 1rem;
            border-radius: 18px;
        }
        .module-hero-main {
            align-items: flex-start;
            width: 100%;
        }
        .module-hero-icon {
            width: 48px;
            height: 48px;
            border-radius: 14px;
            font-size: 1.35rem;
        }
        .module-hero-title { font-size: 1.05rem; }
        .module-hero-subtitle { font-size: .84rem; line-height: 1.4; }
        .module-hero-tag { font-size: .74rem; padding: .4rem .62rem; }

        .section-shell {
            padding: .9rem .85rem .8rem .85rem;
            border-radius: 18px;
        }
        .section-shell-title,
        .visual-card-title,
        .section-title {
            font-size: 1rem;
        }
        .section-shell-subtitle,
        .visual-card-subtitle,
        .chart-note,
        .table-toolbar-note {
            font-size: .78rem;
            line-height: 1.35;
        }

        .kpi-card {
            min-height: 170px;
            padding: .9rem .9rem .85rem .9rem;
            border-radius: 18px;
        }
        .kpi-icon-wrap {
            width: 52px;
            height: 52px;
            margin-bottom: .8rem;
        }
        .kpi-icon { font-size: 1.35rem; }
        .kpi-title {
            font-size: .92rem;
            min-height: auto;
            margin-bottom: .6rem;
        }
        .kpi-value {
            font-size: 2.75rem;
            margin-bottom: .5rem;
        }
        .kpi-note {
            font-size: .7rem;
            max-width: 100%;
        }
        .kpi-chevron {
            width: 24px;
            height: 24px;
            right: .8rem;
            bottom: .75rem;
            font-size: 1rem;
        }

        .product-summary-grid {
            grid-template-columns: repeat(2, minmax(0, 1fr));
            gap: .5rem;
        }
        .product-summary-chip {
            padding: .58rem .65rem;
            border-radius: 12px;
        }
        .product-summary-chip-label { font-size: .72rem; }
        .product-summary-chip-value { font-size: .88rem; }

        .report-summary-card,
        .mini-stat-card,
        div[data-testid="stMetric"] {
            border-radius: 16px;
            padding: .85rem .9rem;
            min-height: auto;
        }
        .report-summary-label,
        .mini-stat-label { font-size: .78rem; }
        .report-summary-value,
        .mini-stat-value { font-size: 1.7rem; }
        .report-summary-note,
        .mini-stat-note { font-size: .74rem; }

        .table-toolbar {
            flex-direction: column;
            align-items: flex-start;
            gap: .5rem;
        }

        div[data-testid="stTabs"] button[role="tab"] {
            padding: .38rem .72rem !important;
            font-size: .82rem !important;
            margin-bottom: .25rem !important;
        }

        .stButton > button,
        .stDownloadButton > button {
            min-height: 42px !important;
            font-size: .9rem !important;
        }

        [data-testid="stDataFrame"] {
            border-radius: 14px;
            overflow: hidden;
        }
    }

    @media (max-width: 520px) {
        .main .block-container { padding-left: .6rem; padding-right: .6rem; padding-top: .45rem !important; }
        .product-summary-grid {
            grid-template-columns: 1fr;
        }
        .kpi-value { font-size: 2.45rem; }
        .module-hero-title { font-size: 1rem; }
        .module-hero-subtitle { font-size: .8rem; }
    }
    </style>
    """,
    unsafe_allow_html=True,
)



# -----------------------------
# PREMIUM VISUAL UPGRADE - STEP 1
# -----------------------------
st.markdown(
    """
    <style>
    .stApp {
        background: radial-gradient(circle at top left, #f3fbf6 0, #ffffff 32%, #f8fafc 100%);
        color: #12352b;
    }

    h1, h2, h3, h4, .section-title {
        letter-spacing: -0.02em;
    }

    section[data-testid="stSidebar"] {
        background: linear-gradient(180deg, #063f2a 0%, #075437 42%, #0a2f22 100%) !important;
        border-right: 0 !important;
        box-shadow: 18px 0 45px rgba(6, 63, 42, 0.16);
    }

    section[data-testid="stSidebar"] > div {
        background: transparent !important;
    }

    section[data-testid="stSidebar"] h1,
    section[data-testid="stSidebar"] h2,
    section[data-testid="stSidebar"] h3,
    section[data-testid="stSidebar"] p,
    section[data-testid="stSidebar"] span,
    section[data-testid="stSidebar"] label,
    section[data-testid="stSidebar"] div {
        color: #eefcf3 !important;
    }

    section[data-testid="stSidebar"] [role="radiogroup"] label {
        padding: .65rem .85rem;
        margin: .25rem 0;
        border-radius: 14px;
        transition: all .18s ease;
        border: 1px solid transparent;
    }

    section[data-testid="stSidebar"] [role="radiogroup"] label:hover {
        background: rgba(255,255,255,.10);
        border-color: rgba(255,255,255,.18);
        transform: translateX(2px);
    }

    section[data-testid="stSidebar"] [data-testid="stMarkdownContainer"] hr {
        border-color: rgba(255,255,255,.16) !important;
    }

    .sidebar-brand {
        display: flex;
        align-items: center;
        gap: .72rem;
        padding: .6rem .2rem .85rem .2rem;
    }

    .sidebar-logo-card {
        background: rgba(255,255,255,.96);
        border: 1px solid rgba(255,255,255,.22);
        border-radius: 18px;
        padding: .75rem .8rem .8rem .8rem;
        box-shadow: 0 14px 34px rgba(0,0,0,.14), inset 0 1px 0 rgba(255,255,255,.6);
        margin: .45rem .05rem .95rem .05rem;
    }
    .sidebar-logo-img {
        width: 100%;
        display: block;
        border-radius: 12px;
        object-fit: contain;
        margin-bottom: .55rem;
    }
    .sidebar-logo-app-name {
        color: #0b4f36 !important;
        font-size: .92rem;
        font-weight: 900;
        line-height: 1.15;
        margin-bottom: .18rem;
    }
    .sidebar-logo-app-subtitle {
        color: #667085 !important;
        font-size: .74rem;
        font-weight: 700;
        line-height: 1.35;
    }

    .sidebar-menu-title {
        font-size: .74rem;
        text-transform: uppercase;
        letter-spacing: .12em;
        font-weight: 800;
        color: rgba(238,252,243,.66) !important;
        margin: .10rem 0 .24rem 0;
    }

    section[data-testid="stSidebar"] [role="radiogroup"] label:has(input:checked) {
        background: linear-gradient(135deg, rgba(255,255,255,.18) 0%, rgba(187,247,208,.14) 100%);
        border-color: rgba(255,255,255,.24);
        box-shadow: inset 0 1px 0 rgba(255,255,255,.16), 0 10px 20px rgba(0,0,0,.12);
        transform: translateX(2px);
    }

    /* Sidebar readability fix: prevent general form radio styling from turning menu items into white pills. */
    section[data-testid="stSidebar"] .stRadio [role="radiogroup"] label,
    section[data-testid="stSidebar"] [role="radiogroup"] label {
        background: transparent !important;
        border: 1px solid transparent !important;
        color: rgba(238,252,243,.86) !important;
        box-shadow: none !important;
        padding: .46rem .78rem !important;
        margin: .10rem 0 !important;
        border-radius: 14px !important;
    }
    section[data-testid="stSidebar"] .stRadio [role="radiogroup"] label span,
    section[data-testid="stSidebar"] .stRadio [role="radiogroup"] label p,
    section[data-testid="stSidebar"] [role="radiogroup"] label span,
    section[data-testid="stSidebar"] [role="radiogroup"] label p {
        color: rgba(238,252,243,.90) !important;
        font-weight: 750 !important;
        line-height: 1.15 !important;
    }
    section[data-testid="stSidebar"] .stRadio [role="radiogroup"] label:hover,
    section[data-testid="stSidebar"] [role="radiogroup"] label:hover {
        background: rgba(255,255,255,.10) !important;
        border-color: rgba(255,255,255,.16) !important;
        transform: translateX(2px);
    }
    section[data-testid="stSidebar"] .stRadio [role="radiogroup"] label:has(input:checked),
    section[data-testid="stSidebar"] [role="radiogroup"] label:has(input:checked) {
        background: linear-gradient(135deg, rgba(255,255,255,.20) 0%, rgba(187,247,208,.16) 100%) !important;
        border-color: rgba(255,255,255,.28) !important;
        box-shadow: inset 0 1px 0 rgba(255,255,255,.16), 0 10px 20px rgba(0,0,0,.12) !important;
    }
    section[data-testid="stSidebar"] .stRadio [role="radiogroup"] label:has(input:checked) span,
    section[data-testid="stSidebar"] .stRadio [role="radiogroup"] label:has(input:checked) p,
    section[data-testid="stSidebar"] [role="radiogroup"] label:has(input:checked) span,
    section[data-testid="stSidebar"] [role="radiogroup"] label:has(input:checked) p {
        color: #ffffff !important;
        font-weight: 900 !important;
    }

    .sidebar-info-card {
        margin-top: 1rem;
        padding: .95rem 1rem;
        border-radius: 18px;
        background: linear-gradient(180deg, rgba(255,255,255,.10) 0%, rgba(255,255,255,.06) 100%);
        border: 1px solid rgba(255,255,255,.14);
        box-shadow: inset 0 1px 0 rgba(255,255,255,.12);
    }
    .sidebar-info-title {
        font-size: .8rem;
        font-weight: 800;
        color: #ffffff;
        margin-bottom: .3rem;
    }
    .sidebar-info-text {
        font-size: .76rem;
        line-height: 1.45;
        color: rgba(238,252,243,.78) !important;
    }

    div[data-testid="stTabs"] button[role="tab"] {
        border-radius: 999px !important;
        border: 1px solid #e4e7ec !important;
        background: #ffffff !important;
        color: #475467 !important;
        font-weight: 700 !important;
        padding: .45rem .95rem !important;
        margin-right: .35rem !important;
    }
    div[data-testid="stTabs"] button[role="tab"][aria-selected="true"] {
        background: linear-gradient(135deg, #ecfdf3 0%, #e1f5ea 100%) !important;
        border-color: #cfe6d8 !important;
        color: #155f3a !important;
        box-shadow: 0 8px 20px rgba(16,24,40,.05) !important;
    }

    div[data-testid="stForm"] {
        background: transparent;
    }
    label[data-testid="stWidgetLabel"] p {
        font-size: .86rem !important;
        font-weight: 700 !important;
        color: #344054 !important;
    }
    div[data-baseweb="input"], div[data-baseweb="select"], div[data-baseweb="textarea"] {
        border-radius: 14px !important;
    }
    div[data-baseweb="input"] > div,
    div[data-baseweb="select"] > div,
    div[data-baseweb="textarea"] > div,
    .stDateInput > div > div,
    .stNumberInput > div > div {
        background: #fbfcfd !important;
        border: 1px solid #e4e7ec !important;
        border-radius: 14px !important;
        box-shadow: 0 4px 12px rgba(16,24,40,.03);
    }
    div[data-baseweb="input"]:focus-within > div,
    div[data-baseweb="select"]:focus-within > div,
    div[data-baseweb="textarea"]:focus-within > div,
    .stDateInput > div > div:focus-within,
    .stNumberInput > div > div:focus-within {
        border-color: #94d3ad !important;
        box-shadow: 0 0 0 4px rgba(34,197,94,.08) !important;
    }
    .stTextInput input, .stNumberInput input, .stDateInput input, .stTextArea textarea {
        color: #17312a !important;
        font-size: .94rem !important;
    }
    .stTextInput input::placeholder, .stTextArea textarea::placeholder {
        color: #98a2b3 !important;
    }
    .stSelectbox [data-baseweb="select"] span {
        color: #17312a !important;
    }
    div[data-testid="stFileUploader"] section {
        border: 1px dashed #cbd5e1 !important;
        border-radius: 16px !important;
        background: #fbfcfd !important;
    }
    .stRadio > div {
        gap: .35rem;
    }
    .stRadio [role="radiogroup"] label {
        border-radius: 14px;
        border: 1px solid #e4e7ec;
        padding: .6rem .8rem;
        background: #fbfcfd;
    }
    .stRadio [role="radiogroup"] label:has(input:checked) {
        border-color: #b7dfc8;
        background: #eefaf4;
    }
    .stButton > button, .stDownloadButton > button {
        min-height: 44px !important;
        border-radius: 14px !important;
        font-weight: 800 !important;
    }

    .sidebar-logo {
        width: 42px;
        height: 42px;
        display: grid;
        place-items: center;
        border-radius: 14px;
        background: rgba(255,255,255,.12);
        border: 1px solid rgba(255,255,255,.22);
        font-size: 1.35rem;
        box-shadow: inset 0 1px 0 rgba(255,255,255,.22);
    }

    .sidebar-title {
        font-size: 1.14rem;
        font-weight: 900;
        line-height: 1.08;
        color: white;
    }

    .sidebar-subtitle {
        color: rgba(238, 252, 243, .72) !important;
        font-size: .78rem;
        margin-top: .28rem;
    }

    .sidebar-farm {
        margin-top: 1.4rem;
        min-height: 190px;
        border-radius: 22px;
        overflow: hidden;
        background:
          radial-gradient(circle at 32% 22%, rgba(252,211,77,.70) 0 10%, transparent 11%),
          linear-gradient(180deg, rgba(255,255,255,.08) 0%, rgba(255,255,255,.04) 35%, rgba(0,0,0,.12) 100%),
          linear-gradient(145deg, rgba(34,197,94,.38) 0%, rgba(5,46,34,.28) 100%);
        position: relative;
        border: 1px solid rgba(255,255,255,.14);
    }
    .sidebar-farm:before {
        content:"";
        position:absolute;
        inset:auto -10px -22px -10px;
        height: 110px;
        background:
          repeating-radial-gradient(ellipse at 50% 100%, rgba(187,247,208,.30) 0 2px, transparent 3px 18px),
          linear-gradient(165deg, #166534 0%, #14532d 70%);
        border-radius: 55% 55% 0 0;
        opacity:.9;
    }
    .sidebar-farm:after {
        content:"";
        position:absolute;
        right: 24px;
        bottom: 44px;
        width: 60px;
        height: 48px;
        background: rgba(255,255,255,.15);
        clip-path: polygon(0 45%, 50% 0, 100% 45%, 100% 100%, 0 100%);
        border: 1px solid rgba(255,255,255,.16);
    }

    .dashboard-topbar {
        display:flex;
        align-items:center;
        justify-content:space-between;
        gap:1rem;
        margin-top: -5.4rem !important;
        margin-bottom: .75rem;
    }
    .dashboard-heading {
        font-size: 2rem;
        font-weight: 900;
        color: #12352b;
        line-height:1;
        margin: 0;
    }
    .dashboard-search-fake {
        min-width: 320px;
        max-width: 440px;
        flex: 1;
        background:#fff;
        border:1px solid #e5e7eb;
        color:#667085;
        padding:.72rem 1rem;
        border-radius:999px;
        box-shadow: 0 8px 22px rgba(16,24,40,.045);
        font-size:.92rem;
    }

    div[data-testid="stTextInput"] input[aria-label="Dashboard Search"] {
        border-radius: 999px !important;
    }

    .premium-hero {
        position:relative;
        overflow:hidden;
        display:flex;
        align-items:center;
        gap:1.1rem;
        padding:1.1rem 1.35rem;
        border-radius: 24px;
        background:
          radial-gradient(circle at 85% 20%, rgba(187,247,208,.65) 0 16%, transparent 18%),
          linear-gradient(135deg, #ffffff 0%, #effaf4 52%, #ddf7e8 100%);
        border: 1px solid #d9eee2;
        box-shadow: 0 16px 42px rgba(16, 24, 40, 0.065);
        margin-bottom: 1.1rem;
    }
    .premium-hero:after {
        content:"";
        position:absolute;
        right:0;
        bottom:0;
        width:44%;
        height:100%;
        opacity:.86;
        background:
          linear-gradient(160deg, transparent 10%, rgba(31,122,77,.08) 10% 28%, transparent 29%),
          radial-gradient(ellipse at 60% 80%, rgba(31,122,77,.32) 0 14%, transparent 15%),
          radial-gradient(ellipse at 74% 77%, rgba(34,197,94,.28) 0 10%, transparent 11%),
          linear-gradient(180deg, transparent 0 52%, rgba(22,101,52,.18) 53% 100%);
        clip-path: polygon(10% 0, 100% 0, 100% 100%, 0 100%);
    }
    .hero-icon {
        width:76px;
        height:76px;
        display:grid;
        place-items:center;
        border-radius: 999px;
        background: #e8f8ee;
        border: 8px solid rgba(255,255,255,.72);
        color:#1f7a4d;
        font-size: 2rem;
        box-shadow: 0 10px 28px rgba(31,122,77,.14);
        z-index:1;
    }
    .hero-content { z-index:1; }
    .hero-title {
        font-size: 1.38rem;
        font-weight: 900;
        color: #12352b;
        margin:0 0 .2rem 0;
    }
    .hero-subtitle {
        color:#435467;
        font-size:.96rem;
        max-width: 610px;
        margin:0;
    }

    div.stButton > button {
        min-height: 150px;
        border-radius: 22px !important;
        background: radial-gradient(circle at 86% 84%, rgba(148,163,184,.10) 0%, transparent 34%), linear-gradient(180deg, #ffffff 0%, #fbfffd 100%) !important;
        border: 1px solid #e2e8f0 !important;
        border-top: 5px solid #22a35f !important;
        box-shadow: 0 14px 34px rgba(16,24,40,.075) !important;
        color: #12352b !important;
        font-weight: 850 !important;
        line-height: 1.32 !important;
        transition: all .18s ease !important;
        padding: 1.1rem 1rem !important;
        white-space: pre-line !important;
        position: relative !important;
        overflow: hidden !important;
        text-align: left !important;
        justify-content: flex-start !important;
        align-items: flex-start !important;
    }
    div.stButton > button:hover {
        transform: translateY(-3px);
        border-color: #cbd5e1 !important;
        box-shadow: 0 22px 46px rgba(16,24,40,.13) !important;
        background: radial-gradient(circle at 86% 84%, rgba(34,197,94,.12) 0%, transparent 36%), linear-gradient(180deg, #ffffff 0%, #f8fbfa 100%) !important;
    }
    div.stButton > button p {
        width: 100%;
        margin: 0 !important;
        font-size: 1.02rem !important;
        letter-spacing: -0.01em !important;
        text-align: left !important;
    }
    /* Color-code cards by column position while keeping the full card clickable. */
    div[data-testid="stHorizontalBlock"] > div:nth-child(1) div.stButton > button { border-top-color: #22a35f !important; }
    div[data-testid="stHorizontalBlock"] > div:nth-child(2) div.stButton > button { border-top-color: #f59e0b !important; }
    div[data-testid="stHorizontalBlock"] > div:nth-child(3) div.stButton > button { border-top-color: #ef4444 !important; }
    div[data-testid="stHorizontalBlock"] > div:nth-child(4) div.stButton > button { border-top-color: #3b82f6 !important; }
    div[data-testid="stHorizontalBlock"] > div:nth-child(5) div.stButton > button { border-top-color: #8b5cf6 !important; }
    div.stForm div.stButton > button,
    div[data-testid="stDownloadButton"] > button,
    div[data-testid="stFormSubmitButton"] > button {
        min-height: 2.85rem !important;
        border-radius: 14px !important;
        border-top-width: 1px !important;
        font-weight: 800 !important;
        background: #ffffff !important;
    }

    .visual-card-title {
        font-size: 1.05rem;
        font-weight: 900;
        color:#12352b;
        margin-bottom: .05rem;
    }
    .visual-card-subtitle {
        color:#667085;
        font-size:.84rem;
        margin-bottom:.55rem;
    }

    /* Soft glow effect for dashboard/report chart containers */
    div[data-testid="stVerticalBlockBorderWrapper"]:has(.stPlotlyChart) {
        border-color: rgba(34, 197, 94, .24) !important;
        background:
            radial-gradient(circle at 92% 8%, rgba(34,197,94,.12) 0%, transparent 34%),
            linear-gradient(180deg, #ffffff 0%, #fbfffd 100%) !important;
        box-shadow:
            0 14px 34px rgba(16,24,40,.06),
            0 0 0 1px rgba(34,197,94,.05),
            0 0 32px rgba(34,197,94,.14) !important;
    }
    div[data-testid="stVerticalBlockBorderWrapper"]:has(.stPlotlyChart):hover {
        border-color: rgba(34, 197, 94, .36) !important;
        box-shadow:
            0 18px 42px rgba(16,24,40,.075),
            0 0 0 1px rgba(34,197,94,.10),
            0 0 44px rgba(34,197,94,.20) !important;
        transition: all .2s ease;
    }
    div[data-testid="stPlotlyChart"] {
        border-radius: 18px !important;
        overflow: hidden !important;
        filter: drop-shadow(0 10px 22px rgba(34,197,94,.08));
    }

    div[data-testid="stVerticalBlockBorderWrapper"] {
        border-radius: 22px !important;
        box-shadow: 0 12px 32px rgba(16,24,40,.055) !important;
        border-color: #e6e9ef !important;
        background: #ffffff !important;
    }
    div[data-testid="stDataFrame"] {
        border-radius: 18px;
        overflow: hidden;
        box-shadow: 0 10px 26px rgba(16,24,40,.035);
    }

    .section-title {
        font-size: 1.18rem !important;
        font-weight: 900 !important;
        color:#12352b !important;
        margin: .2rem 0 .75rem 0 !important;
    }

    .module-hero-card {
        background: linear-gradient(135deg, #ffffff 0%, #eefaf4 58%, #e4f7eb 100%);
        border: 1px solid #dbe9df;
        border-radius: 22px;
        padding: 1.05rem 1.2rem;
        box-shadow: 0 12px 30px rgba(16,24,40,.055);
        margin-bottom: 1rem;
    }
    .module-hero-title {
        font-size: 1.18rem;
        font-weight: 900;
        color: #12352b;
        margin-bottom: .25rem;
    }
    .module-hero-subtitle {
        color: #526170;
        font-size: .93rem;
        line-height: 1.5;
    }
    .mini-stat-card {
        background: #ffffff;
        border: 1px solid #e6e9ef;
        border-radius: 18px;
        padding: .95rem 1rem;
        box-shadow: 0 10px 26px rgba(16,24,40,.04);
        min-height: 118px;
        margin-bottom: 1rem;
    }
    .mini-stat-label {
        color: #667085;
        font-size: .82rem;
        font-weight: 700;
        margin-bottom: .45rem;
    }
    .mini-stat-value {
        color: #12352b;
        font-size: 2rem;
        font-weight: 900;
        line-height: 1.0;
        margin-bottom: .3rem;
    }
    .mini-stat-note {
        color: #6b7280;
        font-size: .8rem;
        font-weight: 600;
    }
    .filter-shell, .form-shell {
        background: #ffffff;
        border: 1px solid #e6e9ef;
        border-radius: 22px;
        box-shadow: 0 12px 30px rgba(16,24,40,.045);
        padding: 1rem 1rem .8rem 1rem;
        margin-bottom: 1rem;
    }
    .filter-shell-title, .form-shell-title {
        color: #12352b;
        font-size: 1rem;
        font-weight: 900;
        margin-bottom: .15rem;
    }
    .filter-shell-subtitle, .form-shell-subtitle {
        color: #667085;
        font-size: .84rem;
        margin-bottom: .9rem;
    }
    .form-section-heading {
        color: #12352b;
        font-size: .96rem;
        font-weight: 900;
        margin: .35rem 0 .6rem 0;
        padding: .35rem .7rem;
        background: #f5fbf7;
        border: 1px solid #e2f0e7;
        border-radius: 12px;
        display: inline-block;
    }
    .product-summary-card {
        background: linear-gradient(135deg, #ffffff 0%, #f7fbf9 100%);
        border: 1px solid #e6e9ef;
        border-radius: 20px;
        padding: 1rem 1.1rem;
        box-shadow: 0 10px 28px rgba(16,24,40,.045);
        margin-bottom: 1rem;
    }
    .product-summary-title {
        color: #12352b;
        font-size: 1.1rem;
        font-weight: 900;
        margin-bottom: .2rem;
    }
    .product-summary-subtitle {
        color: #5f6b7a;
        font-size: .86rem;
        margin-bottom: .7rem;
    }
    .product-summary-grid {
        display: grid;
        grid-template-columns: repeat(4, minmax(0, 1fr));
        gap: .55rem;
    }
    .product-summary-chip {
        background: #ffffff;
        border: 1px solid #e8edf2;
        border-radius: 14px;
        padding: .65rem .75rem;
    }
    .product-summary-chip-label {
        color: #6b7280;
        font-size: .76rem;
        font-weight: 700;
        margin-bottom: .18rem;
    }
    .product-summary-chip-value {
        color: #16342c;
        font-size: .96rem;
        font-weight: 800;
        line-height: 1.2;
    }

    @media (max-width: 900px) {
        .dashboard-topbar { flex-direction: column; align-items:flex-start; }
        .dashboard-search-fake { min-width: 100%; }
        .premium-hero { padding: 1rem; }
        .hero-icon { width: 58px; height:58px; font-size:1.45rem; border-width:6px; }
    }


    /* -----------------------------
       FINAL COMPACT SIDEBAR OVERRIDE
       Makes navigation items closer together while preserving Streamlit hide/show behavior.
    ----------------------------- */
    section[data-testid="stSidebar"] {
        background: linear-gradient(180deg, #064d35 0%, #053d2b 100%) !important;
        border-right: 1px solid rgba(255,255,255,0.08) !important;
        box-shadow: 14px 0 34px rgba(6, 63, 42, 0.13) !important;
    }
    section[data-testid="stSidebar"] > div {
        padding-top: 0.20rem !important;
        padding-left: 0.75rem !important;
        padding-right: 0.75rem !important;
        background: transparent !important;
    }
    .sidebar-brand {
        padding: 0rem 0.15rem 0.45rem 0.15rem !important;
        margin-top: 0rem !important;
        gap: 0.55rem !important;
        align-items: center !important;
    }
    .sidebar-logo {
        width: 42px !important;
        height: 42px !important;
        min-width: 42px !important;
        border-radius: 12px !important;
        display: grid !important;
        place-items: center !important;
        background: rgba(255,255,255,0.10) !important;
        border: 1px solid rgba(255,255,255,0.14) !important;
        overflow: hidden !important;
        box-shadow: inset 0 1px 0 rgba(255,255,255,0.10) !important;
    }
    .sidebar-logo-img {
        width: 31px !important;
        height: 31px !important;
        object-fit: contain !important;
        display: block !important;
    }
    .sidebar-title {
        font-size: 1.02rem !important;
        line-height: 1.15 !important;
        font-weight: 850 !important;
        color: #ffffff !important;
    }
    .sidebar-subtitle {
        font-size: 0.74rem !important;
        line-height: 1.25 !important;
        color: rgba(238,252,243,0.72) !important;
        margin-top: 0.1rem !important;
    }
    section[data-testid="stSidebar"] hr {
        margin: 0.72rem 0 0.78rem 0 !important;
        border-color: rgba(255,255,255,0.12) !important;
    }
    .sidebar-menu-title {
        font-size: 0.70rem !important;
        text-transform: uppercase !important;
        letter-spacing: 0.14em !important;
        font-weight: 850 !important;
        color: rgba(238,252,243,0.70) !important;
        margin: 0.1rem 0 0.45rem 0 !important;
    }
    section[data-testid="stSidebar"] .stRadio > label {
        display: none !important;
    }
    section[data-testid="stSidebar"] [role="radiogroup"] {
        gap: 0.16rem !important;
    }
    section[data-testid="stSidebar"] .stRadio [role="radiogroup"] label,
    section[data-testid="stSidebar"] [role="radiogroup"] label {
        min-height: 40px !important;
        padding: 0.38rem 0.55rem !important;
        margin: 0.03rem 0 !important;
        border-radius: 12px !important;
        background: transparent !important;
        border: 1px solid transparent !important;
        box-shadow: none !important;
        color: rgba(238,252,243,0.86) !important;
        display: flex !important;
        align-items: center !important;
        transition: all 0.16s ease-in-out !important;
    }
    section[data-testid="stSidebar"] .stRadio [role="radiogroup"] label:hover,
    section[data-testid="stSidebar"] [role="radiogroup"] label:hover {
        background: rgba(255,255,255,0.08) !important;
        border-color: rgba(255,255,255,0.12) !important;
        transform: translateX(2px) !important;
    }
    section[data-testid="stSidebar"] .stRadio [role="radiogroup"] label:has(input:checked),
    section[data-testid="stSidebar"] [role="radiogroup"] label:has(input:checked) {
        background: linear-gradient(135deg, rgba(255,255,255,0.18), rgba(255,255,255,0.08)) !important;
        border: 1px solid rgba(255,255,255,0.22) !important;
        box-shadow: inset 0 1px 0 rgba(255,255,255,0.14), 0 7px 16px rgba(0,0,0,0.12) !important;
        transform: translateX(1px) !important;
        min-height: 44px !important;
    }
    section[data-testid="stSidebar"] [role="radiogroup"] label span:first-child {
        margin-right: 0.48rem !important;
        transform: scale(0.78) !important;
    }
    section[data-testid="stSidebar"] .stRadio [role="radiogroup"] label p,
    section[data-testid="stSidebar"] [role="radiogroup"] label p,
    section[data-testid="stSidebar"] .stRadio [role="radiogroup"] label span,
    section[data-testid="stSidebar"] [role="radiogroup"] label span {
        font-size: 0.88rem !important;
        font-weight: 760 !important;
        line-height: 1.05 !important;
        color: rgba(238,252,243,0.90) !important;
        margin: 0 !important;
    }
    section[data-testid="stSidebar"] .stRadio [role="radiogroup"] label:has(input:checked) p,
    section[data-testid="stSidebar"] [role="radiogroup"] label:has(input:checked) p,
    section[data-testid="stSidebar"] .stRadio [role="radiogroup"] label:has(input:checked) span,
    section[data-testid="stSidebar"] [role="radiogroup"] label:has(input:checked) span {
        color: #ffffff !important;
        font-weight: 900 !important;
    }
    .sidebar-info-card {
        margin-top: 0.85rem !important;
        padding: 0.78rem 0.85rem !important;
        border-radius: 15px !important;
        background: rgba(255,255,255,0.08) !important;
        border: 1px solid rgba(255,255,255,0.13) !important;
        box-shadow: inset 0 1px 0 rgba(255,255,255,0.10) !important;
    }
    .sidebar-info-title {
        font-size: 0.78rem !important;
        font-weight: 850 !important;
        color: #ffffff !important;
        margin-bottom: 0.25rem !important;
    }
    .sidebar-info-text {
        font-size: 0.70rem !important;
        line-height: 1.32 !important;
        color: rgba(238,252,243,0.72) !important;
    }
    section[data-testid="stSidebar"] [data-testid="stCaptionContainer"],
    section[data-testid="stSidebar"] .stCaptionContainer {
        color: rgba(238,252,243,0.72) !important;
        font-size: 0.72rem !important;
    }

    /* Sidebar logout button: improve contrast and readability */
    section[data-testid="stSidebar"] .stButton > button {
        background: #ffffff !important;
        color: #064d35 !important;
        border: 1px solid rgba(255,255,255,0.35) !important;
        border-radius: 18px !important;
        font-weight: 850 !important;
        font-size: 0.95rem !important;
        min-height: 46px !important;
        box-shadow: 0 8px 20px rgba(0,0,0,0.14) !important;
    }
    section[data-testid="stSidebar"] .stButton > button p,
    section[data-testid="stSidebar"] .stButton > button span {
        color: #064d35 !important;
        font-weight: 850 !important;
    }
    section[data-testid="stSidebar"] .stButton > button:hover {
        background: #e9fff4 !important;
        color: #043d2b !important;
        border-color: #9be7c3 !important;
        transform: translateY(-1px);
    }
    section[data-testid="stSidebar"] .stButton > button:hover p,
    section[data-testid="stSidebar"] .stButton > button:hover span {
        color: #043d2b !important;
    }
    section[data-testid="stSidebar"] .stButton > button:active {
        background: #dff8ec !important;
        color: #043d2b !important;
    }

    </style>
    """,
    unsafe_allow_html=True,
)

# -----------------------------
# DATABASE HELPERS
# -----------------------------

def get_conn():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn


def execute(sql, params=()):
    with get_conn() as conn:
        cur = conn.execute(sql, params)
        conn.commit()
        return cur.lastrowid


def execute_transaction(operations):
    """Run multiple database operations atomically.

    operations must be a list of (sql, params) tuples. If any operation fails,
    all changes are rolled back so product balances do not become inconsistent.
    """
    with get_conn() as conn:
        try:
            conn.execute("BEGIN")
            for sql, params in operations:
                conn.execute(sql, params)
            conn.commit()
        except Exception:
            conn.rollback()
            raise


def query_df(sql, params=()):
    with get_conn() as conn:
        return pd.read_sql_query(sql, conn, params=params)


def table_count(table):
    df = query_df(f"SELECT COUNT(*) AS n FROM {table}")
    return int(df.loc[0, "n"])


def column_exists(table, column):
    with get_conn() as conn:
        rows = conn.execute(f"PRAGMA table_info({table})").fetchall()
    return any(row["name"] == column for row in rows)


def add_column_if_missing(table, column, column_type):
    if not column_exists(table, column):
        execute(f"ALTER TABLE {table} ADD COLUMN {column} {column_type}")


def init_db():
    execute(
        """
        CREATE TABLE IF NOT EXISTS suppliers (
            supplier_id INTEGER PRIMARY KEY AUTOINCREMENT,
            supplier_code TEXT UNIQUE,
            supplier_name TEXT NOT NULL,
            contact_person TEXT,
            contact_number TEXT,
            email TEXT,
            country TEXT,
            status TEXT DEFAULT 'Active',
            remarks TEXT,
            created_at TEXT DEFAULT CURRENT_TIMESTAMP,
            updated_at TEXT DEFAULT CURRENT_TIMESTAMP
        )
        """
    )

    execute(
        """
        CREATE TABLE IF NOT EXISTS products (
            product_id INTEGER PRIMARY KEY AUTOINCREMENT,
            product_code TEXT UNIQUE,
            product_name TEXT NOT NULL,
            active_ingredient TEXT,
            category TEXT,
            supplier_id INTEGER,
            pack_type TEXT,
            pack_size_value REAL,
            pack_size_unit TEXT,
            beginning_inventory REAL DEFAULT 0,
            quantity REAL DEFAULT 0,
            unit_cost REAL DEFAULT 0,
            batch_number TEXT,
            expiry_date TEXT,
            storage_location TEXT,
            remarks TEXT,
            image_path TEXT,
            created_at TEXT DEFAULT CURRENT_TIMESTAMP,
            updated_at TEXT DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY (supplier_id) REFERENCES suppliers(supplier_id)
        )
        """
    )

    execute(
        """
        CREATE TABLE IF NOT EXISTS stock_in (
            stock_in_id INTEGER PRIMARY KEY AUTOINCREMENT,
            reference_no TEXT UNIQUE,
            date_received TEXT,
            product_id INTEGER,
            supplier_id INTEGER,
            quantity REAL,
            unit_cost REAL,
            batch_number TEXT,
            expiry_date TEXT,
            received_by TEXT,
            remarks TEXT,
            attachment_path TEXT,
            created_at TEXT DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY (product_id) REFERENCES products(product_id),
            FOREIGN KEY (supplier_id) REFERENCES suppliers(supplier_id)
        )
        """
    )

    execute(
        """
        CREATE TABLE IF NOT EXISTS stock_out (
            stock_out_id INTEGER PRIMARY KEY AUTOINCREMENT,
            reference_no TEXT UNIQUE,
            date_released TEXT,
            product_id INTEGER,
            quantity REAL,
            purpose TEXT,
            requested_by TEXT,
            released_by TEXT,
            approved_by TEXT,
            status TEXT,
            remarks TEXT,
            created_at TEXT DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY (product_id) REFERENCES products(product_id)
        )
        """
    )

    execute(
        """
        CREATE TABLE IF NOT EXISTS audit_trail (
            audit_id INTEGER PRIMARY KEY AUTOINCREMENT,
            action_time TEXT DEFAULT CURRENT_TIMESTAMP,
            module TEXT,
            action TEXT,
            details TEXT,
            user_name TEXT
        )
        """
    )

    execute(
        """
        CREATE TABLE IF NOT EXISTS app_settings (
            setting_key TEXT PRIMARY KEY,
            setting_value TEXT,
            updated_at TEXT DEFAULT CURRENT_TIMESTAMP
        )
        """
    )

    execute(
        """
        CREATE TABLE IF NOT EXISTS dropdown_options (
            option_id INTEGER PRIMARY KEY AUTOINCREMENT,
            list_name TEXT NOT NULL,
            option_value TEXT NOT NULL,
            is_active INTEGER DEFAULT 1,
            sort_order INTEGER DEFAULT 0,
            created_at TEXT DEFAULT CURRENT_TIMESTAMP,
            UNIQUE(list_name, option_value)
        )
        """
    )

    execute(
        """
        CREATE TABLE IF NOT EXISTS app_users (
            user_id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_name TEXT NOT NULL UNIQUE,
            role TEXT DEFAULT 'Admin',
            status TEXT DEFAULT 'Active',
            created_at TEXT DEFAULT CURRENT_TIMESTAMP
        )
        """
    )

    seed_settings()
    seed_dropdown_options()
    seed_users()

    # Safe migrations for users who already ran an older version of the app.
    add_column_if_missing("stock_in", "document_type", "TEXT")
    add_column_if_missing("stock_in", "document_no", "TEXT")
    add_column_if_missing("stock_in", "currency", "TEXT DEFAULT 'PHP'")
    add_column_if_missing("stock_in", "cost_type", "TEXT DEFAULT 'Sample'")
    add_column_if_missing("stock_in", "receiving_condition", "TEXT DEFAULT 'Good'")
    add_column_if_missing("stock_in", "qc_status", "TEXT DEFAULT 'Accepted'")
    add_column_if_missing("products", "beginning_inventory", "REAL DEFAULT 0")

    initialize_beginning_inventory_once()
    seed_data()


def log_action(module, action, details, user_name="Admin User"):
    execute(
        """
        INSERT INTO audit_trail (module, action, details, user_name)
        VALUES (?, ?, ?, ?)
        """,
        (module, action, details, user_name),
    )


def generate_code(prefix, table, id_col):
    year = date.today().year
    n = table_count(table) + 1
    return f"{prefix}-{year}-{n:04d}" if prefix in ["SIN", "SOUT"] else f"{prefix}-{n:04d}"


def save_upload(uploaded_file, subfolder):
    if uploaded_file is None:
        return None
    folder = UPLOAD_DIR / subfolder
    folder.mkdir(parents=True, exist_ok=True)
    safe_name = uploaded_file.name.replace(" ", "_")
    stamp = datetime.now().strftime("%Y%m%d%H%M%S")
    path = folder / f"{stamp}_{safe_name}"
    with open(path, "wb") as f:
        f.write(uploaded_file.getbuffer())
    return str(path)



def initialize_beginning_inventory_once():
    """One-time migration for existing databases.

    Older app versions stored only the current product quantity. This calculates a
    starting balance so the Products module can display:
    Beginning Inventory + Sample In - Sample Out = Current Inventory.
    """
    try:
        flag = query_df(
            "SELECT setting_value FROM app_settings WHERE setting_key='beginning_inventory_initialized'"
        )
        if not flag.empty and str(flag.loc[0, "setting_value"]) == "1":
            return

        df = query_df(
            """
            SELECT
                p.product_id,
                COALESCE(p.quantity, 0) AS current_qty,
                COALESCE(si.sample_in, 0) AS sample_in,
                COALESCE(so.sample_out, 0) AS sample_out
            FROM products p
            LEFT JOIN (
                SELECT product_id, SUM(quantity) AS sample_in
                FROM stock_in
                GROUP BY product_id
            ) si ON p.product_id = si.product_id
            LEFT JOIN (
                SELECT product_id, SUM(quantity) AS sample_out
                FROM stock_out
                WHERE status='Released'
                GROUP BY product_id
            ) so ON p.product_id = so.product_id
            """
        )
        for _, row in df.iterrows():
            beginning = float(row["current_qty"] or 0) - float(row["sample_in"] or 0) + float(row["sample_out"] or 0)
            if beginning < 0:
                beginning = 0
            execute(
                "UPDATE products SET beginning_inventory=? WHERE product_id=?",
                (beginning, int(row["product_id"])),
            )

        execute(
            """
            INSERT OR REPLACE INTO app_settings (setting_key, setting_value, updated_at)
            VALUES ('beginning_inventory_initialized', '1', CURRENT_TIMESTAMP)
            """
        )
    except Exception:
        # Do not block app startup if an older DB is in a partial state.
        pass


# -----------------------------
# SETTINGS HELPERS
# -----------------------------

DEFAULT_SETTINGS = {
    "app_name": "Sample Inventory App",
    "company_name": "Farmfix",
    "default_currency": "PHP",
    "expiry_alert_days": "90",
    "urgent_expiry_days": "30",
    "current_user": "Admin User",
    # Dashboard chart preferences. These can be changed in Settings > Chart Settings.
    "chart_category_type": "Bar",
    "chart_category_color": "#1f7a4d",
    "chart_category_palette": "Fresh Multi",
    "chart_category_colors": "#1f7a4d,#22c55e,#f59e0b,#2563eb,#a855f7",
    "chart_status_type": "Donut",
    "chart_status_color": "#22c55e",
    "chart_status_palette": "Status Colors",
    "chart_status_colors": "#22c55e,#2563eb,#f59e0b,#ef4444,#6b7280",
    "chart_supplier_type": "Bar",
    "chart_supplier_color": "#2563eb",
    "chart_supplier_palette": "Ocean",
    "chart_supplier_colors": "#2563eb,#06b6d4,#14b8a6,#84cc16,#a855f7",
    "chart_expiry_type": "Line",
    "chart_expiry_color": "#f59e0b",
    "chart_expiry_palette": "Sunset",
    "chart_expiry_colors": "#f59e0b,#fb923c,#ef4444,#ec4899,#8b5cf6",
}

DROPDOWN_DEFAULTS = {
    "category": ["Biostimulant", "Fertilizer", "Fungicide", "Insecticide", "Herbicide", "Biological", "Root Enhancer", "Plant Nutrition", "Adjuvant", "Others"],
    "pack_type": ["Bottle", "Sachet", "Bag", "Drum", "Box", "Can", "Pouch", "Pack", "Vial", "Container", "Others"],
    "pack_size_unit": ["L", "ml", "kg", "g", "pcs"],
    "storage_location": ["PD Office", "Warehouse", "Research Station", "Cold Storage", "Supplier Area", "Others"],
    "currency": ["PHP", "USD"],
    "cost_type": ["Sample", "Purchased", "Free of Charge", "Replacement"],
    "document_type": ["Delivery Receipt", "Invoice", "Sample Form", "Certificate of Analysis", "Others"],
    "receiving_condition": ["Good", "Damaged", "Leaking", "Missing Label", "Expired Upon Receipt", "Others"],
    "qc_status": ["Accepted", "Hold", "Rejected"],
    "sample_out_purpose": ["Sample Release", "Testing", "Demonstration", "Replacement", "Disposal", "Others"],
    "sample_out_status": ["Released", "Pending", "Approved", "Cancelled"],
    "user_role": ["Admin", "Encoder", "Viewer"],
}


def seed_settings():
    for key, value in DEFAULT_SETTINGS.items():
        execute(
            """
            INSERT OR IGNORE INTO app_settings (setting_key, setting_value)
            VALUES (?, ?)
            """,
            (key, value),
        )


def seed_dropdown_options():
    for list_name, values in DROPDOWN_DEFAULTS.items():
        for i, value in enumerate(values, start=1):
            execute(
                """
                INSERT OR IGNORE INTO dropdown_options (list_name, option_value, is_active, sort_order)
                VALUES (?, ?, 1, ?)
                """,
                (list_name, value, i),
            )


def seed_users():
    execute(
        """
        INSERT OR IGNORE INTO app_users (user_name, role, status)
        VALUES ('Admin User', 'Admin', 'Active')
        """
    )


def get_setting(key, default=None):
    try:
        df = query_df("SELECT setting_value FROM app_settings WHERE setting_key=?", (key,))
        if df.empty:
            return DEFAULT_SETTINGS.get(key, default)
        return df.loc[0, "setting_value"]
    except Exception:
        return DEFAULT_SETTINGS.get(key, default)


def set_setting(key, value):
    execute(
        """
        INSERT INTO app_settings (setting_key, setting_value, updated_at)
        VALUES (?, ?, CURRENT_TIMESTAMP)
        ON CONFLICT(setting_key) DO UPDATE SET
            setting_value=excluded.setting_value,
            updated_at=CURRENT_TIMESTAMP
        """,
        (key, str(value)),
    )


def get_int_setting(key, default_value):
    try:
        return int(float(get_setting(key, str(default_value))))
    except Exception:
        return int(default_value)


# -----------------------------
# CHART CUSTOMIZATION HELPERS
# -----------------------------

CHART_CONFIGS = {
    "category": {
        "title": "Product Sample Category",
        "type_key": "chart_category_type",
        "color_key": "chart_category_color",
        "default_type": "Bar",
        "default_color": "#1f7a4d",
        "palette_key": "chart_category_palette",
        "colors_key": "chart_category_colors",
        "default_palette": "Fresh Multi",
        "options": ["Bar", "Horizontal Bar", "Pie", "Donut"],
    },
    "status": {
        "title": "Product Status Summary",
        "type_key": "chart_status_type",
        "color_key": "chart_status_color",
        "default_type": "Donut",
        "default_color": "#22c55e",
        "palette_key": "chart_status_palette",
        "colors_key": "chart_status_colors",
        "default_palette": "Status Colors",
        "options": ["Donut", "Pie", "Bar", "Horizontal Bar"],
    },
    "supplier": {
        "title": "Recently Received Samples — 70 Days",
        "type_key": "chart_supplier_type",
        "color_key": "chart_supplier_color",
        "default_type": "Table",
        "default_color": "#2563eb",
        "palette_key": "chart_supplier_palette",
        "colors_key": "chart_supplier_colors",
        "default_palette": "Ocean",
        "options": ["Table"],
    },
    "expiry": {
        "title": "Samples Received — Last 12 Months",
        "type_key": "chart_expiry_type",
        "color_key": "chart_expiry_color",
        "default_type": "Bar",
        "default_color": "#f59e0b",
        "palette_key": "chart_expiry_palette",
        "colors_key": "chart_expiry_colors",
        "default_palette": "Sunset",
        "options": ["Line", "Bar", "Area"],
    },
}

COLOR_PALETTES = {
    "Farm Green": ["#1f7a4d", "#22c55e", "#84cc16", "#14b8a6", "#a3e635"],
    "Fresh Multi": ["#1f7a4d", "#22c55e", "#f59e0b", "#2563eb", "#a855f7", "#ef4444"],
    "Status Colors": ["#22c55e", "#2563eb", "#f59e0b", "#ef4444", "#6b7280"],
    "Ocean": ["#2563eb", "#06b6d4", "#14b8a6", "#84cc16", "#a855f7"],
    "Sunset": ["#f59e0b", "#fb923c", "#ef4444", "#ec4899", "#8b5cf6"],
    "Earth": ["#92400e", "#a16207", "#65a30d", "#166534", "#0f766e"],
    "Professional": ["#0f172a", "#334155", "#64748b", "#2563eb", "#22c55e"],
    "Custom": [],
}


def normalize_hex_color(value, fallback="#1f7a4d"):
    value = str(value or "").strip()
    if len(value) == 7 and value.startswith("#"):
        try:
            int(value[1:], 16)
            return value.lower()
        except Exception:
            return fallback
    return fallback


def hex_to_rgb(hex_color):
    hex_color = normalize_hex_color(hex_color).lstrip("#")
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def hex_to_rgba(hex_color, alpha=0.22):
    r, g, b = hex_to_rgb(hex_color)
    return f"rgba({r},{g},{b},{alpha})"


def mix_with_white(hex_color, amount):
    """Return a lighter color by mixing a hex color with white."""
    r, g, b = hex_to_rgb(hex_color)
    r = int(r + (255 - r) * amount)
    g = int(g + (255 - g) * amount)
    b = int(b + (255 - b) * amount)
    return f"#{r:02x}{g:02x}{b:02x}"


def make_color_sequence(base_color, n):
    """Build a soft sequential palette from one selected color."""
    base_color = normalize_hex_color(base_color)
    n = max(int(n or 1), 1)
    if n == 1:
        return [base_color]
    amounts = [0.0, 0.18, 0.32, 0.46, 0.58, 0.68, 0.76, 0.84, 0.90, 0.94]
    colors = []
    for i in range(n):
        amount = amounts[i % len(amounts)]
        colors.append(mix_with_white(base_color, amount))
    return colors


def parse_color_list(value, fallback_colors):
    if isinstance(value, (list, tuple)):
        raw_colors = list(value)
    else:
        raw_colors = str(value or "").split(",")
    valid_colors = []
    for color in raw_colors:
        color = str(color or "").strip()
        if len(color) == 7 and color.startswith("#"):
            try:
                int(color[1:], 16)
                valid_colors.append(color.lower())
            except Exception:
                pass
    if valid_colors:
        return valid_colors
    return list(fallback_colors or ["#1f7a4d"])


def expand_palette(colors, n):
    colors = parse_color_list(colors, ["#1f7a4d"])
    n = max(int(n or 1), 1)
    if len(colors) >= n:
        return colors[:n]
    repeated = []
    while len(repeated) < n:
        repeated.extend(colors)
    return repeated[:n]


def palette_swatches(colors):
    colors = parse_color_list(colors, ["#1f7a4d"])
    swatches = "".join(
        f"<span style='display:inline-block;width:22px;height:22px;border-radius:6px;background:{c};margin-right:6px;border:1px solid #e5e7eb;'></span>"
        for c in colors[:8]
    )
    st.markdown(swatches, unsafe_allow_html=True)


def get_chart_pref(chart_key):
    config = CHART_CONFIGS[chart_key]
    chart_type = get_setting(config["type_key"], config["default_type"])
    if chart_type not in config["options"]:
        chart_type = config["default_type"]

    primary_color = normalize_hex_color(get_setting(config["color_key"], config["default_color"]), config["default_color"])
    palette_name = get_setting(config["palette_key"], config["default_palette"])
    if palette_name not in COLOR_PALETTES:
        palette_name = config["default_palette"]

    default_palette_colors = COLOR_PALETTES.get(palette_name) or make_color_sequence(primary_color, 5)
    saved_colors = get_setting(config["colors_key"], ",".join(default_palette_colors))
    if palette_name == "Custom":
        color_sequence = parse_color_list(saved_colors, make_color_sequence(primary_color, 5))
    else:
        color_sequence = list(default_palette_colors)

    if not color_sequence:
        color_sequence = make_color_sequence(primary_color, 5)
    return chart_type, palette_name, color_sequence


def apply_chart_layout(fig):
    # Keep Plotly clean. If title is blank, force a true empty title so
    # Streamlit/Plotly will not render stray text such as "undefined".
    try:
        raw_title = getattr(fig.layout.title, "text", "")
        current_title = "" if raw_title in [None, "None", "undefined"] else str(raw_title)
    except Exception:
        current_title = ""

    top_margin = 58 if current_title.strip() else 18
    fig.update_layout(
        height=360,
        margin=dict(l=8, r=8, t=top_margin, b=8),
        paper_bgcolor="rgba(0,0,0,0)",
        plot_bgcolor="rgba(255,255,255,0.35)",
        font=dict(size=12, color="#12352b", family="Inter, Segoe UI, Arial, sans-serif"),
        title_text=current_title,
        title_font=dict(size=18, color="#12352b"),
        legend_title_text="",
        legend=dict(orientation='h', yanchor='bottom', y=1.02, xanchor='right', x=1),
        hoverlabel=dict(bgcolor="white", font_size=12, font_color="#12352b"),
    )
    fig.update_xaxes(showgrid=False, zeroline=False, linecolor="#edf2f7", tickfont=dict(color="#667085"))
    fig.update_yaxes(gridcolor="#e8edf2", zeroline=False, tickfont=dict(color="#667085"), tickformat=",.0f")
    return fig

def make_custom_chart(chart_df, label_col, value_col, title, chart_type, colors):
    chart_df = chart_df.copy()
    if chart_df.empty:
        return None
    display_title = title if str(title or "").strip() else None
    chart_df[label_col] = chart_df[label_col].fillna("Unspecified").astype(str)
    chart_df[value_col] = pd.to_numeric(chart_df[value_col], errors="coerce").fillna(0)
    chart_df["_display_value"] = chart_df[value_col].apply(clean_number_for_chart)
    palette = expand_palette(colors, len(chart_df))
    primary_color = palette[0]

    if chart_type in ["Pie", "Donut"]:
        fig = px.pie(
            chart_df,
            names=label_col,
            values=value_col,
            title=display_title,
            hole=0.52 if chart_type == "Donut" else 0,
            color_discrete_sequence=palette,
        )
        fig.update_traces(
            textposition="inside",
            textinfo="percent+label",
            marker=dict(line=dict(color="rgba(255,255,255,.95)", width=3)),
            hovertemplate="%{label}<br>%{customdata}<br>%{percent}<extra></extra>",
            customdata=chart_df["_display_value"],
        )
    elif chart_type == "Horizontal Bar":
        fig = px.bar(
            chart_df,
            y=label_col,
            x=value_col,
            orientation="h",
            title=display_title,
            text="_display_value",
            color=label_col,
            color_discrete_sequence=palette,
        )
        fig.update_traces(marker_line_color="rgba(255,255,255,.9)", marker_line_width=1.5, opacity=.92, texttemplate="%{text}", hovertemplate="%{y}<br>%{text}<extra></extra>")
        fig.update_layout(yaxis={"categoryorder": "total ascending"}, showlegend=False)
    elif chart_type == "Line":
        fig = px.line(chart_df, x=label_col, y=value_col, markers=True, title=display_title)
        fig.update_traces(
            line=dict(color=primary_color, width=4),
            marker=dict(size=9, color=primary_color, line=dict(color="white", width=2)),
            customdata=chart_df["_display_value"],
            hovertemplate="%{x}<br>%{customdata}<extra></extra>",
        )
    elif chart_type == "Area":
        fig = px.area(chart_df, x=label_col, y=value_col, title=display_title, color_discrete_sequence=[primary_color])
        fig.update_traces(
            line=dict(color=primary_color, width=4, shape="spline", smoothing=1.15),
            fillcolor=hex_to_rgba(primary_color, 0.22),
            mode="lines+markers",
            marker=dict(size=7, color=primary_color, line=dict(color="white", width=2)),
            customdata=chart_df["_display_value"],
            hovertemplate="%{x}<br>%{customdata}<extra></extra>",
        )
    else:
        fig = px.bar(
            chart_df,
            x=label_col,
            y=value_col,
            title=display_title,
            text="_display_value",
            color=label_col,
            color_discrete_sequence=palette,
        )
        fig.update_traces(marker_line_color="rgba(255,255,255,.9)", marker_line_width=1.5, opacity=.92, texttemplate="%{text}", hovertemplate="%{x}<br>%{text}<extra></extra>")
        fig.update_layout(showlegend=False)
    return apply_chart_layout(fig)


def chart_preferences_form(form_key_prefix="settings"):
    st.caption("Customize the chart type and select either a preset palette or your own colors.")
    with st.form(f"{form_key_prefix}_chart_preferences_form"):
        chart_values = {}
        row1 = st.columns(2)
        row2 = st.columns(2)
        chart_order = ["category", "status", "supplier", "expiry"]
        palette_options = list(COLOR_PALETTES.keys())
        for col, chart_key in zip(row1 + row2, chart_order):
            config = CHART_CONFIGS[chart_key]
            chart_type, palette_name, color_sequence = get_chart_pref(chart_key)
            with col:
                st.markdown(f"**{config['title']}**")
                selected_type = st.selectbox(
                    "Chart Type",
                    config["options"],
                    index=config["options"].index(chart_type),
                    key=f"{form_key_prefix}_{chart_key}_chart_type",
                )
                selected_palette = st.selectbox(
                    "Color Palette",
                    palette_options,
                    index=palette_options.index(palette_name) if palette_name in palette_options else 0,
                    key=f"{form_key_prefix}_{chart_key}_chart_palette",
                    help="Choose a ready-made palette or select Custom to choose each color.",
                )

                if selected_palette == "Custom":
                    base_custom_colors = expand_palette(color_sequence, 5)
                    picker_cols = st.columns(5)
                    selected_colors = []
                    for i, picker_col in enumerate(picker_cols):
                        selected_colors.append(
                            picker_col.color_picker(
                                f"Color {i + 1}",
                                value=base_custom_colors[i],
                                key=f"{form_key_prefix}_{chart_key}_custom_color_{i}",
                            )
                        )
                else:
                    selected_colors = COLOR_PALETTES[selected_palette]
                    palette_swatches(selected_colors)

                chart_values[chart_key] = (selected_type, selected_palette, selected_colors)
        save_charts = st.form_submit_button("Save Chart Settings", use_container_width=True)
    if save_charts:
        for chart_key, (selected_type, selected_palette, selected_colors) in chart_values.items():
            config = CHART_CONFIGS[chart_key]
            selected_colors = parse_color_list(selected_colors, COLOR_PALETTES.get(config["default_palette"], [config["default_color"]]))
            set_setting(config["type_key"], selected_type)
            set_setting(config["color_key"], selected_colors[0])
            set_setting(config["palette_key"], selected_palette)
            set_setting(config["colors_key"], ",".join(selected_colors))
        log_action("Settings", "Update Chart Settings", "Updated dashboard chart types and color palettes", get_setting("current_user", "Admin User"))
        st.success("Chart settings saved.")
        rerun_app()


def get_dropdown_options(list_name, fallback=None, include_inactive=False):
    fallback = fallback if fallback is not None else DROPDOWN_DEFAULTS.get(list_name, [])
    try:
        if include_inactive:
            df = query_df(
                """
                SELECT option_value FROM dropdown_options
                WHERE list_name=?
                ORDER BY is_active DESC, sort_order, option_value
                """,
                (list_name,),
            )
        else:
            df = query_df(
                """
                SELECT option_value FROM dropdown_options
                WHERE list_name=? AND is_active=1
                ORDER BY sort_order, option_value
                """,
                (list_name,),
            )
        values = [str(x) for x in df["option_value"].tolist()] if not df.empty else []
        return values or fallback
    except Exception:
        return fallback


def add_dropdown_option(list_name, option_value):
    clean_value = str(option_value).strip()
    if not clean_value:
        return
    max_order_df = query_df("SELECT COALESCE(MAX(sort_order), 0) + 1 AS next_order FROM dropdown_options WHERE list_name=?", (list_name,))
    next_order = int(max_order_df.loc[0, "next_order"])
    execute(
        """
        INSERT INTO dropdown_options (list_name, option_value, is_active, sort_order)
        VALUES (?, ?, 1, ?)
        ON CONFLICT(list_name, option_value) DO UPDATE SET is_active=1
        """,
        (list_name, clean_value, next_order),
    )


def set_dropdown_active(list_name, option_value, is_active):
    execute(
        """
        UPDATE dropdown_options
        SET is_active=?
        WHERE list_name=? AND option_value=?
        """,
        (1 if is_active else 0, list_name, option_value),
    )


def image_to_base64(image_path):
    try:
        image_path = Path(image_path)
        if image_path.exists():
            return base64.b64encode(image_path.read_bytes()).decode("utf-8")
    except Exception:
        pass
    return ""


def export_all_data_bytes():
    sheets = {
        "Products": products_base_df(),
        "Suppliers": query_df("SELECT * FROM suppliers ORDER BY supplier_name"),
        "Sample In": query_df("SELECT * FROM stock_in ORDER BY created_at DESC"),
        "Sample Out": query_df("SELECT * FROM stock_out ORDER BY created_at DESC"),
        "Audit Trail": query_df("SELECT * FROM audit_trail ORDER BY action_time DESC"),
        "Settings": query_df("SELECT * FROM app_settings ORDER BY setting_key"),
        "Dropdowns": query_df("SELECT * FROM dropdown_options ORDER BY list_name, sort_order, option_value"),
    }
    return to_excel_bytes(sheets)


def seed_data():
    if table_count("suppliers") == 0:
        suppliers = [
            ("SUP-0001", "Tradecorp", "Juan Dela Cruz", "0912-345-6789", "sales@tradecorp.example", "Spain", "Active", "Main biostimulant supplier"),
            ("SUP-0002", "LIDA Plant Research", "Maria Santos", "0917-111-2222", "info@lida.example", "Spain", "Active", "Fertilizer and phosphite supplier"),
            ("SUP-0003", "Koppert", "Ana Reyes", "0918-222-3333", "contact@koppert.example", "Netherlands", "Active", "Biological products supplier"),
            ("SUP-0004", "Seracsa", "Pedro Cruz", "0919-333-4444", "sales@seracsa.example", "Spain", "Active", "Plant nutrition supplier"),
        ]
        for row in suppliers:
            execute(
                """
                INSERT INTO suppliers
                (supplier_code, supplier_name, contact_person, contact_number, email, country, status, remarks)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?)
                """,
                row,
            )

    if table_count("products") == 0:
        sid = query_df("SELECT supplier_id, supplier_name FROM suppliers")
        s = dict(zip(sid["supplier_name"], sid["supplier_id"]))
        today = date.today()
        products = [
            ("PRD-0001", "Delfan Plus", "L-Alpha Free Amino Acids", "Biostimulant", s.get("Tradecorp"), "Bottle", 1, "L", 25, 25, 600, "DP-2026-001", str(today + timedelta(days=420)), "PD Office", "For anti-stress sample inventory", None),
            ("PRD-0002", "Fytofert Phosco", "Copper & Phosphite", "Fertilizer", s.get("LIDA Plant Research"), "Bottle", 1, "L", 18, 18, 700, "FP-2026-001", str(today + timedelta(days=300)), "PD Office", "For Sigatoka and pineapple evaluation", None),
            ("PRD-0003", "Polyversum OD", "Pythium oligandrum", "Biological", s.get("Koppert"), "Bottle", 250, "ml", 12, 12, 950, "POD-2026-001", str(today + timedelta(days=80)), "Cold Storage", "Biological for Fusarium testing", None),
            ("PRD-0004", "Turbo Root WG", "Humic/Fulvic + Amino Acids", "Root Enhancer", s.get("Tradecorp"), "Sachet", 500, "g", 30, 30, 450, "TRW-2026-001", str(today + timedelta(days=550)), "PD Office", "Root development sample", None),
            ("PRD-0005", "Verno Cu30 + Zn30", "Copper + Zinc", "Fungicide", s.get("Seracsa"), "Bag", 1, "kg", 0, 0, 500, "VCZ-2025-001", str(today - timedelta(days=20)), "Warehouse", "Expired sample batch", None),
        ]
        for row in products:
            execute(
                """
                INSERT INTO products
                (product_code, product_name, active_ingredient, category, supplier_id, pack_type,
                 pack_size_value, pack_size_unit, beginning_inventory, quantity, unit_cost, batch_number, expiry_date,
                 storage_location, remarks, image_path)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                """,
                row,
            )

# -----------------------------
# BUSINESS LOGIC
# -----------------------------

def products_base_df():
    df = query_df(
        """
        SELECT
            p.product_id,
            p.product_code,
            p.product_name,
            p.active_ingredient,
            p.category,
            COALESCE(s.supplier_name, 'No Supplier') AS supplier,
            p.supplier_id,
            p.pack_type,
            p.pack_size_value,
            p.pack_size_unit,
            COALESCE(p.beginning_inventory, 0) AS beginning_inventory,
            p.quantity,
            p.unit_cost,
            p.batch_number,
            p.expiry_date,
            p.storage_location,
            p.remarks,
            p.image_path,
            p.created_at,
            p.updated_at
        FROM products p
        LEFT JOIN suppliers s ON p.supplier_id = s.supplier_id
        ORDER BY p.product_name, p.expiry_date
        """
    )
    return enrich_products(df)


def enrich_products(df):
    if df.empty:
        return df

    today = pd.Timestamp(date.today())
    df = df.copy()
    # Attach transaction totals for the Products module.
    in_totals = query_df("SELECT product_id, COALESCE(SUM(quantity),0) AS sample_in FROM stock_in GROUP BY product_id")
    out_totals = query_df("SELECT product_id, COALESCE(SUM(quantity),0) AS sample_out FROM stock_out WHERE status='Released' GROUP BY product_id")
    if not in_totals.empty:
        df = df.merge(in_totals, on="product_id", how="left")
    else:
        df["sample_in"] = 0
    if not out_totals.empty:
        df = df.merge(out_totals, on="product_id", how="left")
    else:
        df["sample_out"] = 0
    df["sample_in"] = pd.to_numeric(df.get("sample_in", 0), errors="coerce").fillna(0)
    df["sample_out"] = pd.to_numeric(df.get("sample_out", 0), errors="coerce").fillna(0)
    df["beginning_inventory"] = pd.to_numeric(df.get("beginning_inventory", 0), errors="coerce").fillna(0)
    df["current_inventory"] = pd.to_numeric(df.get("quantity", 0), errors="coerce").fillna(0)
    df["inventory_balance_formula"] = df["beginning_inventory"] + df["sample_in"] - df["sample_out"]

    df["expiry_date_dt"] = pd.to_datetime(df["expiry_date"], errors="coerce")
    df["days_to_expiry"] = (df["expiry_date_dt"] - today).dt.days
    df["expiry_status"] = df["days_to_expiry"].apply(
        lambda x: "Expired" if pd.notna(x) and x < 0 else "Not Expired"
    )
    df["stock_status"] = df["current_inventory"].apply(lambda x: "Available" if float(x or 0) > 0 else "Out of Stock")

    def final_status(row):
        if row["expiry_status"] == "Expired":
            return "Expired"
        if row["stock_status"] == "Out of Stock":
            return "Out of Stock"
        urgent_days = get_int_setting("urgent_expiry_days", 30)
        expiring_days = get_int_setting("expiry_alert_days", 90)
        if pd.notna(row["days_to_expiry"]) and row["days_to_expiry"] <= urgent_days:
            return "Urgent Expiry"
        if pd.notna(row["days_to_expiry"]) and row["days_to_expiry"] <= expiring_days:
            return "Expiring Soon"
        return "Available"

    df["final_status"] = df.apply(final_status, axis=1)
    df["pack_size"] = df.apply(
        lambda r: f"{r['pack_size_value']:g} {r['pack_size_unit']}" if pd.notna(r["pack_size_value"]) else "",
        axis=1,
    )
    df["inventory_value"] = df["current_inventory"].astype(float).fillna(0) * df["unit_cost"].astype(float).fillna(0)
    df["total_content"] = df.apply(format_total_content, axis=1)
    return df


def format_total_content(row):
    qty = float(row.get("current_inventory", row.get("quantity", 0)) or 0)
    size = float(row.get("pack_size_value") or 0)
    unit = str(row.get("pack_size_unit") or "")
    total = qty * size
    if unit.lower() == "ml" and total >= 1000:
        return f"{total/1000:g} L"
    if unit.lower() == "g" and total >= 1000:
        return f"{total/1000:g} kg"
    return f"{total:g} {unit}"


def status_badge(status):
    klass = {
        "Available": "pill-available",
        "Expiring Soon": "pill-expiring",
        "Urgent Expiry": "pill-expiring",
        "Expired": "pill-expired",
        "Out of Stock": "pill-out",
        "Not Expired": "pill-info",
        "Released": "pill-available",
        "Approved": "pill-info",
        "Pending": "pill-expiring",
        "Cancelled": "pill-out",
        "Active": "pill-available",
        "Inactive": "pill-out",
    }.get(status, "pill-info")
    return f"<span class='status-pill {klass}'>{status}</span>"


def clean_number(value):
    """Show integers without decimals while preserving real decimal quantities."""
    try:
        if pd.isna(value):
            return ""
    except Exception:
        pass
    if isinstance(value, (int, float, np.integer, np.floating)):
        num = float(value)
        if abs(num - round(num)) < 1e-9:
            return f"{int(round(num)):,}"
        return f"{num:,.4f}".rstrip("0").rstrip(".")
    return value


def clean_number_for_chart(value):
    try:
        if pd.isna(value):
            return ""
        num = float(value)
        if abs(num - round(num)) < 1e-9:
            return f"{int(round(num)):,}"
        return f"{num:,.4f}".rstrip("0").rstrip(".")
    except Exception:
        return str(value)


def format_display_numbers(df):
    """Return a display/export copy with unnecessary .0 decimals removed."""
    formatted = df.copy()
    for col in formatted.columns:
        if pd.api.types.is_numeric_dtype(formatted[col]):
            formatted[col] = formatted[col].apply(clean_number)
    return formatted


STATUS_COLUMN_NAMES = {
    "status", "final_status", "stock_status", "expiry_status",
    "Status", "Final Status", "Stock Status", "Expiry Status"
}

STATUS_DISPLAY_MAP = {
    "available": "🟢 Available",
    "expiring soon": "🟠 Expiring Soon",
    "urgent expiry": "🟠 Urgent Expiry",
    "expired": "🔴 Expired",
    "out of stock": "⚪ Out of Stock",
    "released": "🔵 Released",
    "pending": "🟠 Pending",
    "active": "🟢 Active",
    "inactive": "⚪ Inactive",
    "not expired": "🟢 Not Expired",
    "low stock": "🟠 Low Stock",
    "disposal": "🟣 Disposal",
}

STATUS_STYLE_MAP = {
    "available": {"bg": "#ecfdf3", "fg": "#067647", "border": "#abefc6"},
    "expiring soon": {"bg": "#fffaeb", "fg": "#b54708", "border": "#fedf89"},
    "urgent expiry": {"bg": "#fff4e5", "fg": "#c2410c", "border": "#fdba74"},
    "expired": {"bg": "#fef3f2", "fg": "#b42318", "border": "#fda29b"},
    "out of stock": {"bg": "#f2f4f7", "fg": "#475467", "border": "#d0d5dd"},
    "released": {"bg": "#eff8ff", "fg": "#175cd3", "border": "#b2ddff"},
    "pending": {"bg": "#fffaeb", "fg": "#b54708", "border": "#fedf89"},
    "active": {"bg": "#ecfdf3", "fg": "#067647", "border": "#abefc6"},
    "inactive": {"bg": "#f2f4f7", "fg": "#475467", "border": "#d0d5dd"},
    "not expired": {"bg": "#ecfdf3", "fg": "#067647", "border": "#abefc6"},
    "low stock": {"bg": "#fffaeb", "fg": "#b54708", "border": "#fedf89"},
    "disposal": {"bg": "#f4f3ff", "fg": "#6d28d9", "border": "#c4b5fd"},
}


def _normalize_status(value):
    text = str(value).strip()
    for prefix in ["🟢", "🟠", "🔴", "⚪", "🔵", "🟣"]:
        text = text.replace(prefix, "")
    return text.strip().lower()


def _format_status_display(value):
    key = _normalize_status(value)
    return STATUS_DISPLAY_MAP.get(key, value)


def _status_cell_style(value):
    key = _normalize_status(value)
    conf = STATUS_STYLE_MAP.get(key)
    if not conf:
        return ""
    return (
        f"background-color: {conf['bg']}; color: {conf['fg']}; border: 1px solid {conf['border']}; "
        "font-weight: 700; border-radius: 999px; text-align: center; padding: 0.2rem 0.55rem;"
    )


def _zebra_row(row):
    bg = "#fcfcfd" if getattr(row, "name", 0) % 2 else "#ffffff"
    return [f"background-color: {bg};" for _ in row]


def _styled_dataframe(df, density="Comfortable"):
    styled_df = format_display_numbers(df)
    status_cols = [c for c in styled_df.columns if c in STATUS_COLUMN_NAMES]
    for col in status_cols:
        styled_df[col] = styled_df[col].apply(_format_status_display)

    non_status_cols = [c for c in styled_df.columns if c not in status_cols]
    font_size = "0.84rem" if str(density).lower() == "compact" else "0.9rem"
    cell_padding = "6px 10px" if str(density).lower() == "compact" else "9px 12px"
    header_padding = "8px 10px" if str(density).lower() == "compact" else "10px 12px"

    styler = styled_df.style
    if non_status_cols:
        styler = styler.apply(_zebra_row, axis=1, subset=non_status_cols)
    if status_cols:
        styler = styler.map(_status_cell_style, subset=status_cols)
    styler = styler.set_properties(**{
        "font-size": font_size,
        "color": "#344054",
        "border-bottom": "1px solid #eef2f6",
    })
    styler = styler.set_table_styles([
        {"selector": "th", "props": [
            ("background-color", "#f8fafc"),
            ("color", "#475467"),
            ("font-weight", "700"),
            ("border-bottom", "1px solid #e4e7ec"),
            ("padding", header_padding),
            ("position", "sticky"),
            ("top", "0"),
            ("z-index", "1"),
        ]},
        {"selector": "td", "props": [
            ("padding", cell_padding),
            ("vertical-align", "middle"),
        ]},
    ])
    return styler


def _table_key_from_context(display, height, table_key=None):
    if table_key:
        return f"table_{table_key}"
    try:
        caller = inspect.currentframe().f_back.f_back
        location = f"{caller.f_code.co_name}_{caller.f_lineno}"
    except Exception:
        location = "table"
    col_hash = hashlib.md5("|".join([str(c) for c in display.columns]).encode("utf-8")).hexdigest()[:8]
    return f"table_{location}_{height}_{col_hash}"


def display_df(df, columns=None, height=430, density=None, customizable=True, table_key=None):
    if columns:
        safe_cols = [c for c in columns if c in df.columns]
        display = df[safe_cols].copy()
    else:
        display = df.copy()

    if display.empty:
        st.dataframe(display, use_container_width=True, height=height, hide_index=True)
        return display

    default_cols = list(display.columns)
    keybase = _table_key_from_context(display, height, table_key)
    cols_key = f"{keybase}_visible_columns"
    density_key = f"{keybase}_density"

    selected_cols = default_cols
    selected_density = density or get_setting("table_density", "Comfortable")

    if customizable and len(default_cols) > 1:
        if cols_key not in st.session_state:
            st.session_state[cols_key] = default_cols
        if density_key not in st.session_state:
            st.session_state[density_key] = selected_density

        # Clean saved selections if the table schema changes.
        st.session_state[cols_key] = [c for c in st.session_state[cols_key] if c in default_cols]
        if not st.session_state[cols_key]:
            st.session_state[cols_key] = default_cols

        with st.expander("Customize Table View", expanded=False):
            st.markdown(
                "<div class='custom-table-help'>Choose which columns to show. This only changes the table view; your database records are not deleted.</div>",
                unsafe_allow_html=True,
            )
            c1, c2, c3 = st.columns([3, 1, 1])

            # IMPORTANT: Reset widget-backed session_state BEFORE the widgets
            # are instantiated. Streamlit raises an exception if a widget key
            # is modified after the widget has already been created in the
            # same run. This fixes the Reset Columns button error.
            if c3.button("Reset Columns", key=f"{keybase}_reset", use_container_width=True):
                st.session_state[cols_key] = default_cols
                st.session_state[density_key] = get_setting("table_density", "Comfortable")
                rerun_app()

            selected_cols = c1.multiselect(
                "Visible Columns",
                options=default_cols,
                default=st.session_state[cols_key],
                key=cols_key,
            )
            selected_density = c2.radio(
                "Table Density",
                ["Comfortable", "Compact"],
                index=0 if st.session_state[density_key] == "Comfortable" else 1,
                key=density_key,
            )

    if not selected_cols:
        st.warning("Please select at least one column to display.")
        selected_cols = default_cols

    display = display[selected_cols].copy()
    display_export = format_display_numbers(display)
    st.dataframe(_styled_dataframe(display, density=selected_density), use_container_width=True, height=height, hide_index=True)
    return display_export


def row_matches_search(row, search_text):
    """Robust search helper that safely handles numbers, blanks, dates, and NaN values."""
    if not search_text:
        return True
    values = []
    for value in row.values:
        try:
            if pd.isna(value):
                values.append("")
            else:
                values.append(str(value))
        except Exception:
            values.append(str(value))
    return str(search_text).lower() in " ".join(values).lower()


def to_excel_bytes(sheets: dict):
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        for name, df in sheets.items():
            sheet_name = name[:31]
            df.to_excel(writer, sheet_name=sheet_name, index=False)
    return output.getvalue()


def to_csv_bytes(df):
    return df.to_csv(index=False).encode("utf-8")


def to_pdf_bytes(df, title="Inventory Report"):
    if not REPORTLAB_AVAILABLE:
        return None
    output = io.BytesIO()
    doc = SimpleDocTemplate(output, pagesize=landscape(letter), rightMargin=20, leftMargin=20, topMargin=20, bottomMargin=20)
    styles = getSampleStyleSheet()
    elements = []
    try:
        if APP_LOGO_PATH.exists():
            logo = Image(str(APP_LOGO_PATH), width=52, height=58)
            logo.hAlign = "LEFT"
            elements.extend([logo, Spacer(1, 6)])
    except Exception:
        pass
    elements.extend([Paragraph(title, styles["Title"]), Spacer(1, 10)])
    if df.empty:
        elements.append(Paragraph("No records found.", styles["Normal"]))
    else:
        pdf_df = df.copy().head(40)
        pdf_df = pdf_df.astype(str)
        data = [list(pdf_df.columns)] + pdf_df.values.tolist()
        table = Table(data, repeatRows=1)
        table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1f7a4d")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("ALIGN", (0, 0), (-1, -1), "LEFT"),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 7),
            ("BOTTOMPADDING", (0, 0), (-1, 0), 8),
            ("BACKGROUND", (0, 1), (-1, -1), colors.whitesmoke),
            ("GRID", (0, 0), (-1, -1), 0.25, colors.lightgrey),
        ]))
        elements.append(table)
    doc.build(elements)
    return output.getvalue()


def _month_range(report_month: date):
    month_start = report_month.replace(day=1)
    if month_start.month == 12:
        next_month = month_start.replace(year=month_start.year + 1, month=1)
    else:
        next_month = month_start.replace(month=month_start.month + 1)
    return month_start, next_month


def build_monthly_report_data(report_month: date):
    month_start, next_month = _month_range(report_month)
    products = products_base_df()
    if products.empty:
        products = pd.DataFrame()
    usable_count = len(products[(products["quantity"] > 0) & (products["expiry_status"] != "Expired")]) if not products.empty else 0
    expiring_count = len(products[products["final_status"].isin(["Expiring Soon", "Urgent Expiry"])]) if not products.empty else 0
    expired_count = len(products[products["expiry_status"] == "Expired"]) if not products.empty else 0

    received_month = query_df("""
        SELECT si.reference_no, si.date_received, p.product_name, COALESCE(s.supplier_name, 'Unknown Supplier') AS supplier,
               si.quantity, p.pack_type, si.remarks
        FROM stock_in si
        LEFT JOIN products p ON si.product_id = p.product_id
        LEFT JOIN suppliers s ON si.supplier_id = s.supplier_id
        WHERE date(si.date_received) >= date(?) AND date(si.date_received) < date(?)
        ORDER BY date(si.date_received) DESC, si.stock_in_id DESC
        """, (str(month_start), str(next_month)))
    released_month = query_df("""
        SELECT so.reference_no, so.date_released, p.product_name, COALESCE(s.supplier_name, 'Unknown Supplier') AS supplier,
               so.quantity, p.pack_type, so.released_by, so.remarks
        FROM stock_out so
        LEFT JOIN products p ON so.product_id = p.product_id
        LEFT JOIN suppliers s ON p.supplier_id = s.supplier_id
        WHERE so.status='Released' AND date(so.date_released) >= date(?) AND date(so.date_released) < date(?)
        ORDER BY date(so.date_released) DESC, so.stock_out_id DESC
        """, (str(month_start), str(next_month)))
    recent_received = query_df("""
        SELECT si.date_received, p.product_name AS product, COALESCE(s.supplier_name, 'Unknown Supplier') AS supplier,
               si.quantity, p.pack_type, si.remarks
        FROM stock_in si
        LEFT JOIN products p ON si.product_id = p.product_id
        LEFT JOIN suppliers s ON si.supplier_id = s.supplier_id
        WHERE date(si.date_received) >= date(?) AND date(si.date_received) < date(?)
        ORDER BY date(si.date_received) DESC, si.stock_in_id DESC
        LIMIT 8
        """, (str(month_start), str(next_month)))
    recent_released = query_df("""
        SELECT so.date_released, p.product_name AS product, COALESCE(s.supplier_name, 'Unknown Supplier') AS supplier,
               so.quantity, p.pack_type, so.remarks
        FROM stock_out so
        LEFT JOIN products p ON so.product_id = p.product_id
        LEFT JOIN suppliers s ON p.supplier_id = s.supplier_id
        WHERE so.status='Released' AND date(so.date_released) >= date(?) AND date(so.date_released) < date(?)
        ORDER BY date(so.date_released) DESC, so.stock_out_id DESC
        LIMIT 8
        """, (str(month_start), str(next_month)))

    cursor = month_start
    for _ in range(11):
        cursor = (cursor.replace(day=1) - timedelta(days=1)).replace(day=1)
    start_12 = cursor
    labels, cur = [], start_12
    for _ in range(12):
        labels.append(cur.strftime("%b %Y"))
        cur = (cur.replace(day=28) + timedelta(days=4)).replace(day=1)
    trend_raw = query_df("""
        SELECT substr(date_received, 1, 7) AS month_key, COUNT(*) AS entries
        FROM stock_in
        WHERE date(date_received) >= date(?) AND date(date_received) < date(?)
        GROUP BY substr(date_received, 1, 7)
        """, (str(start_12), str(next_month)))
    trend_map = {}
    if not trend_raw.empty:
        for _, r in trend_raw.iterrows():
            try:
                trend_map[pd.to_datetime(str(r["month_key"]) + "-01").strftime("%b %Y")] = int(r["entries"])
            except Exception:
                pass
    trend_df = pd.DataFrame({"Month": labels, "Entries": [trend_map.get(x, 0) for x in labels]})

    if not products.empty:
        status_df = pd.DataFrame({"Status": ["Usable Stock", "Expiring Soon", "Expired Records"], "Count": [usable_count, expiring_count, expired_count]})
        available_products = products[(products["quantity"] > 0) & (products["expiry_status"] != "Expired")].copy()
        category_df = available_products.groupby("category", dropna=False).size().reset_index(name="Count").sort_values("Count", ascending=False)
        category_df = category_df.rename(columns={"category": "Category"})
        category_df["Category"] = category_df["Category"].fillna("Unspecified")
        supplier_category_df = available_products.copy()
        supplier_category_df["supplier"] = supplier_category_df["supplier"].fillna("Unknown Supplier").replace("", "Unknown Supplier")
        supplier_category_df["category"] = supplier_category_df["category"].fillna("Unspecified").replace("", "Unspecified")
        supplier_category_df = supplier_category_df.groupby(["supplier", "category"], dropna=False).size().reset_index(name="Count")
        supplier_summary_df = supplier_category_df.groupby("supplier", as_index=False)["Count"].sum().sort_values("Count", ascending=False).head(12)
        supplier_summary_df = supplier_summary_df.rename(columns={"supplier": "Supplier", "Count": "Available Product/Batch Count"})
    else:
        status_df = pd.DataFrame(columns=["Status", "Count"])
        category_df = pd.DataFrame(columns=["Category", "Count"])
        supplier_summary_df = pd.DataFrame(columns=["Supplier", "Available Product/Batch Count"])

    observations = [
        f"Usable stock is at {clean_number(usable_count)} product/batch records available for current operations.",
        f"{clean_number(expiring_count)} records are expiring soon. Review utilization plans to avoid expired samples.",
        f"{clean_number(expired_count)} records are expired. Review disposal or documentation requirements.",
        f"{clean_number(len(received_month))} receiving transactions and {clean_number(len(released_month))} released transactions were recorded for {month_start.strftime('%B %Y')}.",
    ]
    return {"report_month": month_start, "usable_count": usable_count, "expiring_count": expiring_count,
            "expired_count": expired_count, "received_count": len(received_month), "released_count": len(released_month),
            "received_month": received_month, "released_month": released_month, "recent_received": recent_received,
            "recent_released": recent_released, "trend_df": trend_df, "status_df": status_df,
            "category_df": category_df, "supplier_summary_df": supplier_summary_df, "observations": observations}


def _safe_table_df(df, max_rows=8):
    if df is None or df.empty:
        return pd.DataFrame()
    return format_display_numbers(df.copy()).head(max_rows).astype(str)


def premium_monthly_pdf_bytes(report_data):
    if not REPORTLAB_AVAILABLE:
        return None

    month_label = report_data["report_month"].strftime("%B %Y")
    output = io.BytesIO()
    doc = SimpleDocTemplate(
        output,
        pagesize=landscape(letter),
        rightMargin=22,
        leftMargin=22,
        topMargin=24,
        bottomMargin=24,
        title=f"Monthly Inventory Summary Report - {month_label}",
        author="Sample Inventory Management System",
    )
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "PdfTitle", parent=styles["Title"], fontName="Helvetica-Bold", fontSize=22,
        leading=26, textColor=colors.HexColor("#0F3B2E"), alignment=TA_LEFT, spaceAfter=4
    )
    subtitle_style = ParagraphStyle(
        "PdfSubTitle", parent=styles["Normal"], fontName="Helvetica", fontSize=11,
        leading=14, textColor=colors.HexColor("#667085"), alignment=TA_LEFT, spaceAfter=10
    )
    section_style = ParagraphStyle(
        "PdfSection", parent=styles["Heading2"], fontName="Helvetica-Bold", fontSize=14,
        leading=18, textColor=colors.HexColor("#17312A"), spaceBefore=6, spaceAfter=8,
        borderPadding=0
    )
    small_style = ParagraphStyle(
        "PdfSmall", parent=styles["Normal"], fontName="Helvetica", fontSize=9,
        leading=12, textColor=colors.HexColor("#667085")
    )
    body_style = ParagraphStyle(
        "PdfBody", parent=styles["Normal"], fontName="Helvetica", fontSize=10,
        leading=13, textColor=colors.HexColor("#344054")
    )
    bullet_style = ParagraphStyle(
        "PdfBullet", parent=body_style, leftIndent=10, bulletIndent=0, spaceBefore=2, spaceAfter=2
    )

    def _header_footer(canvas, doc):
        canvas.saveState()
        page_w, page_h = landscape(letter)
        canvas.setFillColor(colors.HexColor("#0B4F35"))
        canvas.rect(0, page_h - 12, page_w, 12, fill=1, stroke=0)
        canvas.setFillColor(colors.HexColor("#0B4F35"))
        canvas.setFont("Helvetica-Bold", 8)
        canvas.drawString(doc.leftMargin, 10, "Sample Inventory Management System")
        canvas.setFillColor(colors.HexColor("#667085"))
        canvas.setFont("Helvetica", 8)
        canvas.drawRightString(page_w - doc.rightMargin, 10, f"Page {canvas.getPageNumber()}")
        canvas.restoreState()

    def section_label(text):
        return Table(
            [[Paragraph(f"<font color='#0B4F35'>▌</font> <b>{text}</b>", section_style)]],
            colWidths=[10.9*inch],
            style=TableStyle([("LEFTPADDING", (0,0), (-1,-1), 0), ("RIGHTPADDING", (0,0), (-1,-1), 0), ("TOPPADDING", (0,0), (-1,-1), 0), ("BOTTOMPADDING", (0,0), (-1,-1), 0)])
        )

    def make_kpi_card(title, value, note, fill="#F8FBF9", accent="#16A34A"):
        title_p = Paragraph(f"<font color='#667085'><b>{title}</b></font>", small_style)
        value_p = Paragraph(f"<font color='{accent}'><b>{clean_number(value)}</b></font>", ParagraphStyle("kpi_v", parent=styles["Normal"], fontName="Helvetica-Bold", fontSize=20, leading=24, textColor=colors.HexColor(accent)))
        note_p = Paragraph(f"<font color='#98A2B3'>{note}</font>", ParagraphStyle("kpi_n", parent=small_style, fontSize=8, leading=10))
        tbl = Table([[title_p], [value_p], [note_p]], colWidths=[2.02*inch])
        tbl.setStyle(TableStyle([
            ("BACKGROUND", (0,0), (-1,-1), colors.HexColor(fill)),
            ("BOX", (0,0), (-1,-1), 0.7, colors.HexColor("#D8E6DE")),
            ("LINEABOVE", (0,0), (-1,0), 2.2, colors.HexColor(accent)),
            ("ROUNDEDCORNERS", [10,10,10,10]),
            ("LEFTPADDING", (0,0), (-1,-1), 10), ("RIGHTPADDING", (0,0), (-1,-1), 10),
            ("TOPPADDING", (0,0), (-1,-1), 8), ("BOTTOMPADDING", (0,0), (-1,-1), 8),
        ]))
        return tbl

    def add_table(elements, title, df, max_rows=8, note=None):
        elements.append(section_label(title))
        if note:
            elements.append(Paragraph(note, small_style))
            elements.append(Spacer(1, 4))
        tdf = _safe_table_df(df, max_rows)
        if tdf.empty:
            empty = Table([[Paragraph("No records available.", small_style)]], colWidths=[10.8*inch])
            empty.setStyle(TableStyle([
                ("BACKGROUND", (0,0), (-1,-1), colors.HexColor("#F8FAFC")),
                ("BOX", (0,0), (-1,-1), 0.5, colors.HexColor("#E5E7EB")),
                ("LEFTPADDING", (0,0), (-1,-1), 10), ("RIGHTPADDING", (0,0), (-1,-1), 10),
                ("TOPPADDING", (0,0), (-1,-1), 8), ("BOTTOMPADDING", (0,0), (-1,-1), 8),
            ]))
            elements.append(empty)
            elements.append(Spacer(1, 10))
            return
        data = [list(tdf.columns)] + tdf.values.tolist()
        cols = len(tdf.columns)
        base = [1.0] * cols
        for i, col in enumerate(tdf.columns):
            cname = str(col).lower()
            if "remark" in cname:
                base[i] = 2.6
            elif "product" in cname:
                base[i] = 2.1
            elif "supplier" in cname:
                base[i] = 1.9
            elif "date" in cname:
                base[i] = 1.15
            elif "status" in cname:
                base[i] = 1.4
            elif "category" in cname:
                base[i] = 1.6
            elif "count" in cname or "quantity" in cname:
                base[i] = 1.0
        total = sum(base)
        available_w = 10.8
        col_widths = [available_w * x / total * inch for x in base]
        tbl = Table(data, colWidths=col_widths, repeatRows=1)
        style_cmds = [
            ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#0B4F35")),
            ("TEXTCOLOR", (0,0), (-1,0), colors.white),
            ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
            ("FONTSIZE", (0,0), (-1,0), 9),
            ("ALIGN", (0,0), (-1,0), "LEFT"),
            ("VALIGN", (0,0), (-1,-1), "MIDDLE"),
            ("GRID", (0,0), (-1,-1), 0.35, colors.HexColor("#E3EBE6")),
            ("FONTSIZE", (0,1), (-1,-1), 8),
            ("TEXTCOLOR", (0,1), (-1,-1), colors.HexColor("#344054")),
            ("ROWBACKGROUNDS", (0,1), (-1,-1), [colors.HexColor("#FFFFFF"), colors.HexColor("#F7FBF8")]),
            ("LEFTPADDING", (0,0), (-1,-1), 7),
            ("RIGHTPADDING", (0,0), (-1,-1), 7),
            ("TOPPADDING", (0,0), (-1,-1), 7),
            ("BOTTOMPADDING", (0,0), (-1,-1), 7),
        ]
        tbl.setStyle(TableStyle(style_cmds))
        elements.append(tbl)
        elements.append(Spacer(1, 10))

    elements = []
    generated_text = datetime.now().strftime("%B %d, %Y %I:%M %p")
    logo_path = APP_LOGO_PATH if Path(APP_LOGO_PATH).exists() else None
    if logo_path:
        try:
            img = Image(str(logo_path), width=0.68*inch, height=0.68*inch)
            title_block = [
                [img, Paragraph("Monthly Inventory Summary Report", title_style)],
                ["", Paragraph(f"Product Sample Monitoring • {month_label}<br/><font size='9' color='#667085'>Generated {generated_text}</font>", subtitle_style)],
            ]
            head = Table(title_block, colWidths=[0.85*inch, 10.2*inch])
            head.setStyle(TableStyle([("SPAN", (1,0), (1,0)), ("LEFTPADDING", (0,0), (-1,-1), 0), ("RIGHTPADDING", (0,0), (-1,-1), 0), ("VALIGN", (0,0), (-1,-1), "MIDDLE")]))
            elements.append(head)
        except Exception:
            elements.append(Paragraph("Monthly Inventory Summary Report", title_style))
            elements.append(Paragraph(f"Product Sample Monitoring • {month_label}<br/><font size='9' color='#667085'>Generated {generated_text}</font>", subtitle_style))
    else:
        elements.append(Paragraph("Monthly Inventory Summary Report", title_style))
        elements.append(Paragraph(f"Product Sample Monitoring • {month_label}<br/><font size='9' color='#667085'>Generated {generated_text}</font>", subtitle_style))

    hero = Table([[Paragraph("<b>Management Summary</b><br/><font size='10'>This report consolidates current usable stock, expiry exposure, monthly receiving activity, monthly release activity, and available product/batch composition for management review.</font>", body_style)]], colWidths=[10.8*inch])
    hero.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,-1), colors.HexColor("#EEF7F1")),
        ("BOX", (0,0), (-1,-1), 0.7, colors.HexColor("#D8E6DE")),
        ("LEFTPADDING", (0,0), (-1,-1), 12), ("RIGHTPADDING", (0,0), (-1,-1), 12),
        ("TOPPADDING", (0,0), (-1,-1), 10), ("BOTTOMPADDING", (0,0), (-1,-1), 10),
    ]))
    elements += [hero, Spacer(1, 10)]

    kpi_row = Table([[
        make_kpi_card("Usable Stock", report_data.get("usable_count", 0), "Available product/batch records", fill="#EEF9F2", accent="#16A34A"),
        make_kpi_card("Expiring Soon", report_data.get("expiring_count", 0), "Records needing monitoring", fill="#FFF7E8", accent="#F59E0B"),
        make_kpi_card("Expired Records", report_data.get("expired_count", 0), "Review for action", fill="#FDEDED", accent="#EF4444"),
        make_kpi_card("Received This Month", report_data.get("received_count", 0), "Receiving transactions", fill="#EEF4FF", accent="#2563EB"),
        make_kpi_card("Released This Month", report_data.get("released_count", 0), "Released transactions", fill="#F3EEFF", accent="#7C3AED"),
    ]], colWidths=[2.06*inch]*5)
    kpi_row.setStyle(TableStyle([("LEFTPADDING", (0,0), (-1,-1), 0), ("RIGHTPADDING", (0,0), (-1,-1), 0), ("TOPPADDING", (0,0), (-1,-1), 0), ("BOTTOMPADDING", (0,0), (-1,-1), 0)]))
    elements += [kpi_row, Spacer(1, 12)]

    elements.append(section_label("Management Highlights"))
    for obs in report_data.get("observations", []):
        elements.append(Paragraph(f"• {obs}", bullet_style))
    elements.append(Spacer(1, 8))

    summary_left = _safe_table_df(report_data.get("status_df"), 10)
    summary_right = _safe_table_df(report_data.get("category_df"), 10)
    left_tbl = Table([[Paragraph("<b>Inventory Status Snapshot</b>", body_style)], [Table([list(summary_left.columns)] + summary_left.values.tolist(), colWidths=[1.9*inch, 1.1*inch], repeatRows=1) if not summary_left.empty else Paragraph("No data", small_style)]], colWidths=[5.2*inch])
    right_tbl = Table([[Paragraph("<b>Product Sample Category Composition</b>", body_style)], [Table([list(summary_right.columns)] + summary_right.values.tolist(), colWidths=[2.8*inch, 1.0*inch], repeatRows=1) if not summary_right.empty else Paragraph("No data", small_style)]], colWidths=[5.2*inch])
    for outer in (left_tbl, right_tbl):
        outer.setStyle(TableStyle([("BACKGROUND", (0,0), (-1,-1), colors.white), ("BOX", (0,0), (-1,-1), 0.6, colors.HexColor("#E2EBE5")), ("LEFTPADDING", (0,0), (-1,-1), 8), ("RIGHTPADDING", (0,0), (-1,-1), 8), ("TOPPADDING", (0,0), (-1,-1), 8), ("BOTTOMPADDING", (0,0), (-1,-1), 8)]))
    # inner table styling
    for inner_df, outer in ((summary_left, left_tbl), (summary_right, right_tbl)):
        if not inner_df.empty:
            inner = outer._cellvalues[1][0]
            inner.setStyle(TableStyle([
                ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#EAF3EE")),
                ("TEXTCOLOR", (0,0), (-1,0), colors.HexColor("#17312A")),
                ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
                ("GRID", (0,0), (-1,-1), 0.3, colors.HexColor("#E2EBE5")),
                ("ROWBACKGROUNDS", (0,1), (-1,-1), [colors.white, colors.HexColor("#F8FBF9")]),
                ("FONTSIZE", (0,0), (-1,-1), 9),
                ("LEFTPADDING", (0,0), (-1,-1), 6), ("RIGHTPADDING", (0,0), (-1,-1), 6),
                ("TOPPADDING", (0,0), (-1,-1), 5), ("BOTTOMPADDING", (0,0), (-1,-1), 5),
            ]))
    two_col = Table([[left_tbl, right_tbl]], colWidths=[5.35*inch, 5.35*inch])
    two_col.setStyle(TableStyle([("VALIGN", (0,0), (-1,-1), "TOP"), ("LEFTPADDING", (0,0), (-1,-1), 0), ("RIGHTPADDING", (0,0), (-1,-1), 8), ("TOPPADDING", (0,0), (-1,-1), 0), ("BOTTOMPADDING", (0,0), (-1,-1), 0)]))
    elements += [two_col, Spacer(1, 12)]

    add_table(elements, "Sample-In Entries — Last 12 Months", report_data.get("trend_df"), 12, note="Counts receiving transactions per month for the current month and previous 11 months.")
    add_table(elements, "Available Product/Batch Samples by Supplier", report_data.get("supplier_summary_df"), 12, note="Each available product/batch record is counted as 1. Expired and out-of-stock records are excluded.")
    add_table(elements, "Recently Received Samples", report_data.get("recent_received"), 8, note="Latest receiving transactions recorded in the selected reporting month.")
    add_table(elements, "Recently Released Samples", report_data.get("recent_released"), 8, note="Latest release transactions recorded in the selected reporting month.")

    recommendations = [
        "Review expired product/batch records for disposal, retention, or documentation requirements.",
        "Prioritize utilization of expiring-soon samples to avoid avoidable losses.",
        "Validate key receiving and release records before final monthly submission.",
        "Continue regular database backup after each reporting cycle.",
    ]
    elements.append(section_label("Recommended Actions"))
    rec_data = [[Paragraph(f"<b>{i+1}.</b> {txt}", body_style)] for i, txt in enumerate(recommendations)]
    rec_tbl = Table(rec_data, colWidths=[10.8*inch])
    rec_tbl.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,-1), colors.HexColor("#FCFDFD")),
        ("BOX", (0,0), (-1,-1), 0.6, colors.HexColor("#E2EBE5")),
        ("ROWBACKGROUNDS", (0,0), (-1,-1), [colors.HexColor("#FFFFFF"), colors.HexColor("#F8FBF9")]),
        ("LEFTPADDING", (0,0), (-1,-1), 10), ("RIGHTPADDING", (0,0), (-1,-1), 10),
        ("TOPPADDING", (0,0), (-1,-1), 8), ("BOTTOMPADDING", (0,0), (-1,-1), 8),
    ]))
    elements.append(rec_tbl)
    elements.append(Spacer(1, 6))
    elements.append(Paragraph(f"Report generated on {generated_text}", ParagraphStyle("pdf_footer_note", parent=small_style, alignment=TA_RIGHT)))

    doc.build(elements, onFirstPage=_header_footer, onLaterPages=_header_footer)
    return output.getvalue()



def _ppt_color(hex_color):
    return RGBColor.from_string(str(hex_color).replace("#", ""))


PPT_FONT_HEAD = "Aptos Display"
PPT_FONT_BODY = "Aptos"


def _ppt_add_picture(slide, image_path, x, y, w, h):
    try:
        if image_path and Path(image_path).exists():
            return slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), Inches(w), Inches(h))
    except Exception:
        return None
    return None


def _ppt_add_textbox(slide, text, x, y, w, h, size=18, bold=False, color="12352B", align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = PPT_FONT_HEAD if bold or size >= 16 else PPT_FONT_BODY
    run.font.color.rgb = _ppt_color(color) if isinstance(color, str) else color
    if align:
        p.alignment = align
    return box


def _ppt_add_shape(slide, x, y, w, h, fill="FFFFFF", line="E5E7EB", radius=False, transparency=0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _ppt_color(fill)
    shp.fill.transparency = transparency
    shp.line.color.rgb = _ppt_color(line)
    shp.line.width = Pt(0.65)
    try:
        shp.shadow.inherit = False
    except Exception:
        pass
    return shp


def _ppt_add_footer(slide, page_no=None, note="For Internal Monthly Reporting"):
    _ppt_add_shape(slide, 0, 7.12, 13.333, 0.38, fill="0B4F35", line="0B4F35", radius=False)
    _ppt_add_textbox(slide, "Sample Inventory Management System", 0.42, 7.18, 3.2, 0.22, size=7.5, bold=True, color="FFFFFF")
    _ppt_add_textbox(slide, f"{note}  •  Generated {datetime.now().strftime('%b %d, %Y %I:%M %p')}", 4.2, 7.18, 4.6, 0.22, size=7.5, color="EAF7F0", align=PP_ALIGN.CENTER)
    if page_no:
        _ppt_add_textbox(slide, f"Page {page_no}", 11.9, 7.18, 1.0, 0.22, size=7.5, color="FFFFFF", align=PP_ALIGN.RIGHT)


def _ppt_add_header(slide, title, subtitle="", page_no=None):
    _ppt_add_shape(slide, 0, 0, 13.333, 0.18, fill="0B4F35", line="0B4F35", radius=False)
    _ppt_add_textbox(slide, title, 0.55, 0.38, 6.2, 0.42, size=22, bold=True, color="0F3B2E")
    if subtitle:
        _ppt_add_textbox(slide, subtitle, 0.57, 0.83, 6.8, 0.28, size=10.5, color="667085")
    _ppt_add_footer(slide, page_no=page_no)


def _ppt_add_icon_circle(slide, icon, x, y, accent="16A34A"):
    """Add a clean square icon badge instead of a circle to prevent the overly rounded look."""
    soft = {"16A34A":"EAF7F0", "F59E0B":"FFF3D6", "EF4444":"FDECEC", "2563EB":"EAF2FF", "7C3AED":"F3ECFF"}.get(accent, "EAF7F0")
    badge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.56), Inches(0.56))
    badge.fill.solid()
    badge.fill.fore_color.rgb = _ppt_color(soft)
    badge.line.color.rgb = _ppt_color("DDE8E1")
    badge.line.width = Pt(0.6)
    _ppt_add_textbox(slide, icon, x+0.05, y+0.12, 0.46, 0.28, size=15, bold=True, color=accent, align=PP_ALIGN.CENTER)


def _ppt_add_kpi_card(slide, title, value, note, icon, x, y, accent="16A34A", w=2.25, h=1.38):
    _ppt_add_shape(slide, x, y, w, h, fill="FFFFFF", line="DCE8E1", radius=True)
    # top accent
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x+0.03), Inches(y), Inches(w-0.06), Inches(0.075))
    bar.fill.solid(); bar.fill.fore_color.rgb = _ppt_color(accent); bar.line.color.rgb = _ppt_color(accent)
    _ppt_add_icon_circle(slide, icon, x+0.16, y+0.23, accent)
    _ppt_add_textbox(slide, title, x+0.86, y+0.22, w-1.0, 0.3, size=10.8, bold=True, color="17312A")
    _ppt_add_textbox(slide, clean_number(value), x+0.86, y+0.56, w-1.0, 0.48, size=28, bold=True, color=accent)
    _ppt_add_textbox(slide, note, x+0.86, y+1.08, w-1.0, 0.22, size=7.2, bold=True, color="98A2B3")


def _ppt_set_cell_text(cell, text, font_size=9, bold=False, color="344054", align=PP_ALIGN.LEFT):
    tf = cell.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.075)
    tf.margin_right = Inches(0.055)
    tf.margin_top = Inches(0.035)
    tf.margin_bottom = Inches(0.035)
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = PPT_FONT_HEAD if bold else PPT_FONT_BODY
    run.font.color.rgb = _ppt_color(color) if isinstance(color, str) else color


def _ppt_add_table(slide, df, x, y, w, h, max_rows=8, font_size=9.0):
    tdf = _safe_table_df(df, max_rows)
    if tdf.empty or len(tdf.columns) == 0:
        _ppt_add_shape(slide, x, y, w, 0.55, fill="F8FAFC", line="E5E7EB")
        _ppt_add_textbox(slide, "No records found.", x+0.15, y+0.16, w-0.3, 0.22, size=10.5, color="667085")
        return
    rows, cols = len(tdf) + 1, len(tdf.columns)
    table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table

    # Wider, presentation-friendly column distribution.
    widths = []
    for col in tdf.columns:
        cname = str(col).lower()
        if "remark" in cname:
            widths.append(2.95)
        elif "product" in cname:
            widths.append(2.35)
        elif "supplier" in cname:
            widths.append(2.15)
        elif "date" in cname:
            widths.append(1.15)
        elif "quantity" in cname or "qty" in cname:
            widths.append(0.95)
        elif "pack" in cname:
            widths.append(0.95)
        else:
            widths.append(1.0)
    scale = w / sum(widths)
    for i, cw in enumerate(widths):
        table.columns[i].width = Inches(cw * scale)

    try:
        table.rows[0].height = Inches(0.48)
        body_h = max(0.38, min(0.56, (h - 0.48) / max(1, rows - 1)))
        for r in range(1, rows):
            table.rows[r].height = Inches(body_h)
    except Exception:
        pass

    header_size = 10.0
    body_size = 9.0

    for c, col in enumerate(tdf.columns):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _ppt_color("0B4F35")
        # python-pptx table cells do not expose a .line property; borders are handled by fill/spacing.
        _ppt_set_cell_text(
            cell,
            str(col).replace("_", " ").title(),
            font_size=header_size,
            bold=True,
            color=RGBColor(255, 255, 255),
            align=PP_ALIGN.LEFT,
        )

    for r, row in enumerate(tdf.values.tolist(), start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _ppt_color("FFFFFF" if r % 2 else "F6FBF8")
            # Avoid cell.line here; _Cell has no line attribute in python-pptx.
            sval = str(val)
            limit = 42 if "remark" in str(tdf.columns[c]).lower() else 32
            _ppt_set_cell_text(cell, sval[:limit], font_size=body_size, color="344054", align=PP_ALIGN.LEFT)


def _ppt_add_bar_chart(slide, df, label_col, value_col, x, y, w, h, title="", color="16A34A"):
    if df is None or df.empty:
        _ppt_add_textbox(slide, "No chart data available.", x, y, w, 0.35, size=10, color="667085")
        return
    chart_data = CategoryChartData()
    chart_data.categories = [str(v) for v in df[label_col].tolist()]
    chart_data.add_series(title or value_col, [float(v or 0) for v in df[value_col].tolist()])
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), chart_data).chart
    chart.has_legend = False
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.number_format = '0'
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.font.size = Pt(9)
    chart.plots[0].data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = _ppt_color(color)


def _ppt_add_line_chart(slide, df, x, y, w, h, title="Entries"):
    if df is None or df.empty:
        _ppt_add_textbox(slide, "No trend data available.", x, y, w, 0.35, size=10, color="667085")
        return
    chart_data = CategoryChartData()
    chart_data.categories = [str(v) for v in df["Month"].tolist()]
    chart_data.add_series(title, [float(v or 0) for v in df["Entries"].tolist()])
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(x), Inches(y), Inches(w), Inches(h), chart_data).chart
    chart.has_legend = False
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.number_format = '0'
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.series[0].format.line.color.rgb = _ppt_color("16A34A")
    chart.series[0].format.line.width = Pt(2.6)
    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.font.size = Pt(9)


def _ppt_add_donut_chart(slide, status_df, x, y, w, h):
    if status_df is None or status_df.empty:
        _ppt_add_textbox(slide, "No status data available.", x, y, w, 0.35, size=10, color="667085")
        return
    chart_data = CategoryChartData()
    chart_data.categories = [str(v) for v in status_df["Status"].tolist()]
    chart_data.add_series("Status", [float(v or 0) for v in status_df["Count"].tolist()])
    chart = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(x), Inches(y), Inches(w), Inches(h), chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.font.size = Pt(9)
    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.show_percentage = True
    chart.plots[0].data_labels.font.size = Pt(9)


def _ppt_panel(slide, title, subtitle, x, y, w, h):
    _ppt_add_shape(slide, x, y, w, h, fill="FFFFFF", line="DDE8E1", radius=False)
    _ppt_add_textbox(slide, title, x+0.22, y+0.18, w-0.44, 0.28, size=16, bold=True, color="0F3B2E")
    if subtitle:
        _ppt_add_textbox(slide, subtitle, x+0.22, y+0.48, w-0.44, 0.25, size=12, color="667085")


# PPT_REPORT_REFERENCE_DESIGN = "monthly_inventory_report_2026_05 (2).pptx"
# The PPT generator intentionally mirrors the uploaded reference deck details:
# Executive Summary, Monthly KPI Overview, Sample-In Transaction Trend,
# Inventory Composition, Recent Sample Movement, and Actions and Follow-up.
def premium_monthly_pptx_bytes(report_data):
    """Generate a polished management-style monthly inventory PPT report.

    Design goals:
    - less text crowding than the previous version
    - clearer slide story: summary -> movement -> stock status -> details -> actions
    - stronger chart/table hierarchy for management review
    """
    if not PPTX_AVAILABLE:
        return None

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    month_label = report_data["report_month"].strftime("%B %Y")
    generated = datetime.now().strftime("%b %d, %Y %I:%M %p")
    brand = "Sample Inventory Management System"
    ppt_logo_path = APP_DIR / "assets" / "ppt_farmfix_reference_logo.png"

    green = "0B4F35"
    green2 = "1F7A4D"
    pale = "EEF8F2"
    text = "17312A"
    muted = "667085"
    amber = "F59E0B"
    red = "EF4444"
    blue = "2563EB"
    purple = "7C3AED"

    usable = int(report_data.get("usable_count", 0) or 0)
    expiring = int(report_data.get("expiring_count", 0) or 0)
    expired = int(report_data.get("expired_count", 0) or 0)
    received = int(report_data.get("received_count", 0) or 0)
    released = int(report_data.get("released_count", 0) or 0)
    total_records = usable + expiring + expired

    def add_bg(slide, fill="FFFFFF"):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = _ppt_color(fill)
        # Clean executive-report background: no circles/ovals, only slim accent bars.
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.10))
        band.fill.solid(); band.fill.fore_color.rgb = _ppt_color(green); band.line.color.rgb = _ppt_color(green)
        top_soft = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.10), Inches(13.333), Inches(0.06))
        top_soft.fill.solid(); top_soft.fill.fore_color.rgb = _ppt_color("DFF2E7"); top_soft.line.color.rgb = _ppt_color("DFF2E7")
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(6.74), Inches(12.2), Inches(0.015))
        line.fill.solid(); line.fill.fore_color.rgb = _ppt_color("D7E9DE"); line.line.color.rgb = _ppt_color("D7E9DE")

    def add_footer(slide, page):
        _ppt_add_shape(slide, 0, 7.04, 13.333, 0.46, fill=green, line=green, radius=False)
        logo_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.11), Inches(7.09), Inches(0.33), Inches(0.33))
        logo_bg.fill.solid(); logo_bg.fill.fore_color.rgb = _ppt_color("FFFFFF"); logo_bg.fill.transparency = 8; logo_bg.line.color.rgb = _ppt_color("BFE6CB")
        _ppt_add_picture(slide, ppt_logo_path, 0.14, 7.11, 0.27, 0.28)
        _ppt_add_textbox(slide, brand, 0.55, 7.17, 3.5, 0.20, size=7.5, bold=True, color="FFFFFF")
        _ppt_add_textbox(slide, f"For Internal Monthly Reporting  •  Generated {generated}", 4.05, 7.17, 5.7, 0.20, size=7.5, color="EAF7F0", align=PP_ALIGN.CENTER)
        _ppt_add_textbox(slide, f"Page {page}", 11.83, 7.17, 1.0, 0.20, size=7.5, bold=True, color="FFFFFF", align=PP_ALIGN.RIGHT)

    def add_title(slide, title, subtitle, page):
        add_bg(slide)
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.38), Inches(0.06), Inches(0.72))
        accent.fill.solid(); accent.fill.fore_color.rgb = _ppt_color(green2); accent.line.color.rgb = _ppt_color(green2)
        _ppt_add_textbox(slide, title, 0.75, 0.32, 7.0, 0.46, size=24, bold=True, color=text)
        _ppt_add_textbox(slide, subtitle, 0.77, 0.83, 7.8, 0.28, size=12, color=muted)
        add_footer(slide, page)

    def metric_card(slide, x, y, w, h, title, value, note, icon, color):
        _ppt_add_shape(slide, x+0.045, y+0.055, w, h, fill="EAF3ED", line="EAF3ED", radius=False, transparency=16)
        _ppt_add_shape(slide, x, y, w, h, fill="FFFFFF", line="DDE8E1", radius=False)
        # color bar + icon
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.075))
        bar.fill.solid(); bar.fill.fore_color.rgb = _ppt_color(color); bar.line.color.rgb = _ppt_color(color)
        _ppt_add_icon_circle(slide, icon, x+0.18, y+0.25, color)
        _ppt_add_textbox(slide, title, x+0.82, y+0.20, w-0.96, 0.36, size=11, bold=True, color=text)
        _ppt_add_textbox(slide, clean_number(value), x+0.82, y+0.60, w-0.96, 0.46, size=30, bold=True, color=color)
        _ppt_add_textbox(slide, note, x+0.82, y+1.12, w-0.96, 0.24, size=8, color="98A2B3")
        underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x+0.82), Inches(y+h-0.15), Inches(0.55), Inches(0.025))
        underline.fill.solid(); underline.fill.fore_color.rgb = _ppt_color(color); underline.line.color.rgb = _ppt_color(color)

    def insight_line(slide, y, icon, heading, body, color=green2):
        _ppt_add_icon_circle(slide, icon, 0.82, y, color)
        _ppt_add_textbox(slide, heading, 1.55, y+0.00, 3.4, 0.22, size=10.5, bold=True, color=text)
        _ppt_add_textbox(slide, body, 1.55, y+0.25, 4.25, 0.30, size=7.6, color="344054")

    def panel(slide, title, subtitle, x, y, w, h):
        _ppt_add_shape(slide, x+0.04, y+0.05, w, h, fill="EAF3ED", line="EAF3ED", radius=False, transparency=18)
        _ppt_add_shape(slide, x, y, w, h, fill="FFFFFF", line="DDE8E1", radius=False)
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x+0.22), Inches(y+0.20), Inches(0.045), Inches(0.36))
        accent.fill.solid(); accent.fill.fore_color.rgb = _ppt_color(green2); accent.line.color.rgb = _ppt_color(green2)
        _ppt_add_textbox(slide, title, x+0.34, y+0.15, w-0.6, 0.30, size=13.2, bold=True, color=text)
        if subtitle:
            _ppt_add_textbox(slide, subtitle, x+0.34, y+0.60, w-0.6, 0.24, size=7.9, color=muted)

    kpis = [
        ("Usable Stock", usable, "Stock + not expired", "✅", "16A34A"),
        ("Expiring Soon", expiring, "Needs review", "⏳", amber),
        ("Expired Records", expired, "Past expiry", "⚠️", red),
        ("Received This Month", received, "Transactions", "⬇️", blue),
        ("Released This Month", released, "Transactions", "⬆️", purple),
    ]

    # Slide 1: Executive dashboard
    slide = prs.slides.add_slide(blank)
    add_bg(slide, "F7FBF8")
    # FarmFix logo and title layout copied from the user's edited reference deck.
    _ppt_add_picture(slide, ppt_logo_path, 0.78, 0.48, 0.92, 1.05)
    _ppt_add_textbox(slide, "Inventory\nPerformance Report", 1.95, 0.38, 5.95, 1.35, size=36, bold=True, color=text)
    _ppt_add_textbox(slide, f"Executive Summary  •  {month_label}", 1.95, 1.62, 5.8, 0.35, size=18, color=muted)
    _ppt_add_shape(slide, 7.41, 0.42, 5.26, 1.28, fill="FFFFFF", line="DDE8E1")
    _ppt_add_textbox(slide, "Management Snapshot", 7.54, 0.58, 2.75, 0.34, size=13.0, bold=True, color=text)
    _ppt_add_textbox(slide, f"{total_records} current product/batch records reviewed for this monthly inventory report.", 7.62, 0.96, 2.70, 0.58, size=8.4, color="344054")
    pill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.44), Inches(0.68), Inches(1.30), Inches(0.34))
    pill.fill.solid(); pill.fill.fore_color.rgb = _ppt_color(green2); pill.line.color.rgb = _ppt_color(green2)
    _ppt_add_textbox(slide, "Internal reporting", 10.49, 0.755, 1.20, 0.18, size=7.4, bold=True, color="FFFFFF", align=PP_ALIGN.CENTER)
    _ppt_add_shape(slide, 11.92, 0.63, 0.60, 0.60, fill="EAF7F0", line="CFE6D8")
    _ppt_add_textbox(slide, "📦", 12.02, 0.77, 0.40, 0.18, size=15, align=PP_ALIGN.CENTER)
    for i, (t, v, n, icon, a) in enumerate(kpis):
        metric_card(slide, 0.58 + i*2.48, 2.18, 2.18, 1.42, t, v, n, icon, a)
    panel(slide, "Key Message", "What management should focus on", 0.58, 4.0, 6.1, 2.18)
    key_msg = f"Inventory is usable for current operations, but {expired} expired records and {expiring} expiring soon records require follow-up before the next reporting cycle."
    _ppt_add_textbox(slide, key_msg, 0.95, 4.72, 5.25, 0.72, size=12.4, bold=True, color=text)
    _ppt_add_textbox(slide, f"Monthly movement: {received} receiving transaction(s) and {released} released transaction(s) for {month_label}.", 0.95, 5.55, 5.25, 0.38, size=9.5, color="344054")
    panel(slide, "Management Highlights", "Auto-generated from inventory data", 6.9, 4.0, 5.85, 2.35)
    for i, obs in enumerate(report_data.get("observations", [])[:3]):
        _ppt_add_textbox(slide, f"• {obs}", 7.18, 4.78+i*0.48, 5.25, 0.38, size=7.6, color="344054")
    add_footer(slide, 1)

    # Slide 2: KPI deep dive
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Monthly KPI Overview", f"Inventory movement and status for {month_label}", 2)
    for i, (t, v, n, icon, a) in enumerate(kpis):
        metric_card(slide, 0.55 + i*2.55, 1.32, 2.25, 1.55, t, v, n, icon, a)
    panel(slide, "Interpretation", "How to read this month’s metrics", 0.55, 3.18, 5.95, 3.00)
    lines = [
        ("✅", "Usable Stock", "Records with inventory balance and not past expiry."),
        ("⬇️", "Monthly Movement", "Received and released figures count transaction records, not total packs."),
        ("⚠️", "Expired Records", "Expired records require disposition review, documentation, or disposal plan."),
    ]
    for i, (ic, head, body) in enumerate(lines):
        insight_line(slide, 4.05+i*0.66, ic, head, body, [green2, blue, red][i])
    panel(slide, "Attention Needed", "Priority review items", 6.85, 3.18, 5.95, 3.00)
    _ppt_add_textbox(slide, f"Expired Records: {expired}", 7.25, 3.92, 3.2, 0.32, size=15.5, bold=True, color=red)
    _ppt_add_textbox(slide, f"Expiring Soon: {expiring}", 7.25, 4.42, 3.2, 0.32, size=15.5, bold=True, color=amber)
    _ppt_add_textbox(slide, "Recommended action: validate utilization, disposal, and redistribution plan before next reporting period.", 7.25, 5.02, 4.85, 0.70, size=8.8, color="344054")

    # Slide 3: Trends
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Sample-In Transaction Trend", "Receiving activity across the current month and previous 11 months", 3)
    panel(slide, "Sample-In Entries — Last 12 Months", "Counts receiving transaction records per month", 0.55, 1.25, 8.1, 5.15)
    _ppt_add_line_chart(slide, report_data.get("trend_df", pd.DataFrame()), 0.9, 1.95, 7.35, 3.75)
    total_entries = report_data.get("trend_df", pd.DataFrame()).get("Entries", pd.Series(dtype=float)).sum() if report_data.get("trend_df", pd.DataFrame()) is not None else 0
    _ppt_add_textbox(slide, f"Total entries shown: {clean_number(total_entries)}", 1.0, 5.88, 3.4, 0.25, size=10.5, bold=True, color=text)
    panel(slide, "Trend Insight", "Quick interpretation", 8.95, 1.25, 3.8, 5.15)
    _ppt_add_textbox(slide, "Use this chart to identify months with unusually high receiving activity and validate if all supporting documents are complete.", 9.25, 1.92, 3.2, 0.98, size=9.2, color="344054")
    _ppt_add_textbox(slide, f"For {month_label}, there were {received} receiving transaction(s).", 9.25, 3.10, 3.2, 0.58, size=12.5, bold=True, color=blue)
    _ppt_add_textbox(slide, "Action: review bulk receiving months for complete batch number, supplier, expiry, and remarks information.", 9.25, 4.12, 3.2, 0.88, size=9.0, color="344054")

    # Slide 4: Category and status
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Inventory Composition", "Product categories and status distribution", 4)
    panel(slide, "Product Sample Category", "Available product/batch records by category", 0.55, 1.25, 6.1, 5.2)
    cat = report_data.get("category_df", pd.DataFrame())
    if cat is not None and not cat.empty:
        _ppt_add_bar_chart(slide, cat.head(8), "Category", "Count", 0.92, 1.95, 5.45, 3.8, color=green2)
    panel(slide, "Product Status Summary", "Distribution of current product/batch records", 6.9, 1.25, 5.85, 5.2)
    _ppt_add_donut_chart(slide, report_data.get("status_df", pd.DataFrame()), 7.25, 1.85, 2.8, 2.8)
    _ppt_add_textbox(slide, f"Total Records\n{total_records}", 10.35, 2.06, 1.75, 0.86, size=15.5, bold=True, color=text, align=PP_ALIGN.CENTER)
    _ppt_add_textbox(slide, f"Usable: {usable}\nExpiring: {expiring}\nExpired: {expired}", 10.16, 3.18, 2.05, 0.88, size=10.5, color="344054", align=PP_ALIGN.CENTER)
    _ppt_add_textbox(slide, "Focus: expired and expiring samples should be reviewed before the next monthly reporting cycle.", 7.35, 5.12, 4.7, 0.62, size=9.2, color="344054")

    # Slide 5: Recent receiving details
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Recent Sample Movement", f"Receiving records for {month_label}", 5)
    panel(slide, "Recently Received Samples", "Latest receiving transactions in the selected reporting window", 0.55, 1.18, 12.25, 5.7)
    recent = report_data.get("recent_received", pd.DataFrame())
    recent_cols = [c for c in ["date_received", "product", "supplier", "quantity", "pack_type", "remarks"] if c in getattr(recent, "columns", [])]
    recent_view = recent[recent_cols].copy() if recent_cols else recent
    _ppt_add_table(slide, recent_view, 0.82, 1.92, 11.7, 3.85, max_rows=5, font_size=10.0)
    _ppt_add_textbox(slide, "Tip: Detailed line-item records remain available in the app report preview and exported Excel files.", 0.9, 6.18, 8.6, 0.25, size=9, color=muted)
    _ppt_add_textbox(slide, "Report table intentionally shows top records only to keep the presentation readable.", 8.55, 6.18, 3.7, 0.25, size=8.5, color="98A2B3", align=PP_ALIGN.RIGHT)

    # Slide 6: Releases and recommendations
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Actions and Follow-up", "Released samples, risks, and next-step recommendations", 6)
    panel(slide, "Released Samples This Month", "Sample Out transaction summary", 0.55, 1.22, 5.9, 2.35)
    rel = report_data.get("released_month", pd.DataFrame())
    if rel is not None and not rel.empty:
        rel_cols = [c for c in ["date_released", "product_name", "quantity", "released_by", "remarks"] if c in rel.columns]
        _ppt_add_table(slide, rel[rel_cols].copy() if rel_cols else rel, 0.82, 1.9, 5.35, 1.45, max_rows=3, font_size=8.7)
    else:
        _ppt_add_textbox(slide, "No released transactions were recorded for the selected month.", 0.9, 2.15, 4.8, 0.28, size=11.5, color=muted)
    panel(slide, "Recommended Actions", "Suggested follow-up for next reporting cycle", 6.75, 1.22, 6.05, 4.85)
    actions = [
        ("Review Expired Records", f"Validate {expired} expired product/batch records for disposal, retention, or documentation.", red, "⚠️"),
        ("Monitor Expiring Soon", f"Plan utilization or redistribution for {expiring} expiring soon records.", amber, "⏳"),
        ("Validate Monthly Movement", f"Check {received} received and {released} released transaction records for completeness.", blue, "📋"),
        ("Maintain Backup", "Download database backup after monthly reporting and store it in the official shared folder.", green2, "🛡️"),
    ]
    for i, (head, body, color, icon) in enumerate(actions):
        y = 2.12 + i*0.82
        _ppt_add_shape(slide, 6.98, y-0.02, 5.48, 0.74, fill="F8FBF9", line="E5EFE8", radius=False)
        _ppt_add_icon_circle(slide, icon, 7.10, y+0.06, color)
        _ppt_add_textbox(slide, head, 7.86, y+0.03, 4.3, 0.22, size=10.5, bold=True, color=text)
        _ppt_add_textbox(slide, body, 7.86, y+0.28, 4.18, 0.36, size=7.4, color="344054")
    _ppt_add_shape(slide, 0.55, 4.03, 5.9, 2.04, fill="FFFFFF", line="DDE8E1")
    _ppt_add_textbox(slide, "Monthly Release Summary", 0.9, 4.34, 2.8, 0.25, size=15, bold=True, color=text)
    _ppt_add_textbox(slide, f"Released transactions: {released}", 0.9, 4.82, 3.8, 0.30, size=18, bold=True, color=purple)
    _ppt_add_textbox(slide, "Review sample-out records for complete requester, approver, and purpose details.", 0.9, 5.28, 4.85, 0.52, size=9.0, color="344054")

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def rerun_app():
    try:
        st.rerun()
    except Exception:
        st.experimental_rerun()


def page_header(title, subtitle):
    st.markdown(
        f"""
        <div class="app-header">
            <div>
                <div class="app-title">{title}</div>
                <div class="app-subtitle">{subtitle}</div>
            </div>
            <div class="small-note">📦 Inventory monitoring • SQLite database • Streamlit app</div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def module_hero(title: str, subtitle: str, icon: str = "📦", tag: str = ""):
    tag_html = f"<div class='module-hero-tag'>{tag}</div>" if tag else ""
    st.markdown(
        f"""
        <div class="module-hero-card">
            <div class="module-hero-main">
                <div class="module-hero-icon">{icon}</div>
                <div>
                    <div class="module-hero-title">{title}</div>
                    <div class="module-hero-subtitle">{subtitle}</div>
                </div>
            </div>
            {tag_html}
        </div>
        """,
        unsafe_allow_html=True,
    )


def section_shell(title: str, subtitle: str = ""):
    subtitle_html = f"<div class='section-shell-subtitle'>{subtitle}</div>" if subtitle else ""
    st.markdown(f"<div class='section-shell'><div class='section-shell-title'>{title}</div>{subtitle_html}", unsafe_allow_html=True)


def end_section_shell():
    st.markdown("</div>", unsafe_allow_html=True)


def report_summary_card(label, value, note=""):
    st.markdown(
        f"""
        <div class="report-summary-card">
            <div class="report-summary-label">{label}</div>
            <div class="report-summary-value">{value}</div>
            <div class="report-summary-note">{note}</div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def chart_note(text):
    if text:
        st.markdown(f"<div class='chart-note'>{text}</div>", unsafe_allow_html=True)


def open_details(title, df):
    def content():
        if df.empty:
            st.info("No records found.")
        else:
            display_df(df, height=380)
            st.download_button(
                "Download CSV",
                to_csv_bytes(df),
                file_name=f"{title.lower().replace(' ', '_')}.csv",
                mime="text/csv",
                use_container_width=True,
            )
    if hasattr(st, "dialog"):
        @st.dialog(title, width="large")
        def _dialog():
            content()
        _dialog()
    else:
        st.subheader(title)
        content()


def get_supplier_options():
    df = query_df("SELECT supplier_id, supplier_name FROM suppliers ORDER BY supplier_name")
    if df.empty:
        return {}, []
    mapping = {row["supplier_name"]: int(row["supplier_id"]) for _, row in df.iterrows()}
    return mapping, list(mapping.keys())


def get_product_options():
    df = products_base_df()
    if df.empty:
        return {}, []
    df["label"] = df.apply(
        lambda r: (
            f"{r['product_name']} | Batch: {r['batch_number'] or '-'} | "
            f"Pack Size: {r.get('pack_size', '') or '-'} | "
            f"Loc: {r.get('storage_location', '') or '-'} | "
            f"Qty: {clean_number(r['current_inventory'])} {str(r['pack_type']).lower()}(s) | "
            f"ID: {int(r['product_id'])}"
        ),
        axis=1,
    )
    mapping = {row["label"]: int(row["product_id"]) for _, row in df.iterrows()}
    return mapping, list(mapping.keys())


def get_product(product_id):
    df = products_base_df()
    df = df[df["product_id"] == product_id]
    if df.empty:
        return None
    return df.iloc[0]


def recommended_fefo(product_name):
    df = products_base_df()
    candidates = df[
        (df["product_name"] == product_name)
        & (df["quantity"] > 0)
        & (df["expiry_status"] != "Expired")
    ].sort_values("expiry_date_dt")
    if candidates.empty:
        return None
    return candidates.iloc[0]



def dashboard_visual_header():
    left, right = st.columns([1.0, 1.35])
    with left:
        st.markdown("<div class='dashboard-heading'>Dashboard</div>", unsafe_allow_html=True)
    with right:
        search_query = st.text_input(
            "Dashboard Search",
            placeholder="🔎 Search products, suppliers, batch numbers, sample-in/out records...",
            label_visibility="collapsed",
            key="dashboard_global_search",
        )
    st.markdown(
        """
        <div class="premium-hero">
            <div class="hero-icon">🌿</div>
            <div class="hero-content">
                <div class="hero-title">Inventory Overview</div>
                <p class="hero-subtitle">Monitor product samples, stock levels, expiry status, and sample movements in real time.</p>
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )
    return search_query


def chart_card_heading(title, subtitle=""):
    subtitle_html = f"<div class='visual-card-subtitle'>{subtitle}</div>" if subtitle else ""
    st.markdown(
        f"<div class='visual-card-title'>{title}</div>{subtitle_html}",
        unsafe_allow_html=True,
    )


def build_dashboard_search_results(search_query):
    q = str(search_query or "").strip()
    if not q:
        return None
    like = f"%{q}%"

    products = query_df(
        """
        SELECT
            'Product / Batch' AS "Result Type",
            p.product_name AS "Product",
            p.batch_number AS "Batch No.",
            COALESCE(s.supplier_name, '') AS "Supplier",
            p.category AS "Category",
            CASE
                WHEN p.pack_size_value IS NOT NULL AND p.pack_size_unit IS NOT NULL
                THEN printf('%g %s', p.pack_size_value, p.pack_size_unit)
                ELSE ''
            END AS "Pack Size",
            p.pack_type AS "Pack Type",
            p.quantity AS "Current Qty",
            p.expiry_date AS "Date",
            p.storage_location AS "Storage / Purpose",
            p.remarks AS "Remarks"
        FROM products p
        LEFT JOIN suppliers s ON p.supplier_id = s.supplier_id
        WHERE p.product_name LIKE ? OR p.batch_number LIKE ? OR p.product_code LIKE ?
           OR p.active_ingredient LIKE ? OR p.category LIKE ? OR COALESCE(s.supplier_name, '') LIKE ?
        ORDER BY p.updated_at DESC
        LIMIT 50
        """,
        (like, like, like, like, like, like),
    )

    suppliers = query_df(
        """
        SELECT
            'Supplier' AS "Result Type",
            '' AS "Product",
            '' AS "Batch No.",
            supplier_name AS "Supplier",
            country AS "Category",
            '' AS "Pack Size",
            '' AS "Pack Type",
            '' AS "Current Qty",
            created_at AS "Date",
            contact_person AS "Storage / Purpose",
            remarks AS "Remarks"
        FROM suppliers
        WHERE supplier_name LIKE ? OR supplier_code LIKE ? OR contact_person LIKE ?
           OR contact_number LIKE ? OR email LIKE ? OR country LIKE ? OR remarks LIKE ?
        ORDER BY supplier_name
        LIMIT 50
        """,
        (like, like, like, like, like, like, like),
    )

    sample_in = query_df(
        """
        SELECT
            'Sample In' AS "Result Type",
            COALESCE(p.product_name, '') AS "Product",
            COALESCE(p.batch_number, si.batch_number, '') AS "Batch No.",
            COALESCE(s.supplier_name, '') AS "Supplier",
            COALESCE(p.category, '') AS "Category",
            CASE
                WHEN p.pack_size_value IS NOT NULL AND p.pack_size_unit IS NOT NULL
                THEN printf('%g %s', p.pack_size_value, p.pack_size_unit)
                ELSE ''
            END AS "Pack Size",
            COALESCE(p.pack_type, '') AS "Pack Type",
            si.quantity AS "Current Qty",
            si.date_received AS "Date",
            si.received_by AS "Storage / Purpose",
            si.remarks AS "Remarks"
        FROM stock_in si
        LEFT JOIN products p ON si.product_id = p.product_id
        LEFT JOIN suppliers s ON si.supplier_id = s.supplier_id
        WHERE si.reference_no LIKE ? OR COALESCE(p.product_name, '') LIKE ?
           OR COALESCE(p.batch_number, si.batch_number, '') LIKE ? OR COALESCE(s.supplier_name, '') LIKE ?
           OR si.received_by LIKE ? OR si.document_no LIKE ? OR si.remarks LIKE ?
        ORDER BY si.created_at DESC
        LIMIT 50
        """,
        (like, like, like, like, like, like, like),
    )

    sample_out = query_df(
        """
        SELECT
            'Sample Out' AS "Result Type",
            COALESCE(p.product_name, '') AS "Product",
            COALESCE(p.batch_number, '') AS "Batch No.",
            COALESCE(s.supplier_name, '') AS "Supplier",
            COALESCE(p.category, '') AS "Category",
            CASE
                WHEN p.pack_size_value IS NOT NULL AND p.pack_size_unit IS NOT NULL
                THEN printf('%g %s', p.pack_size_value, p.pack_size_unit)
                ELSE ''
            END AS "Pack Size",
            COALESCE(p.pack_type, '') AS "Pack Type",
            so.quantity AS "Current Qty",
            so.date_released AS "Date",
            so.purpose AS "Storage / Purpose",
            so.remarks AS "Remarks"
        FROM stock_out so
        LEFT JOIN products p ON so.product_id = p.product_id
        LEFT JOIN suppliers s ON p.supplier_id = s.supplier_id
        WHERE so.reference_no LIKE ? OR COALESCE(p.product_name, '') LIKE ?
           OR COALESCE(p.batch_number, '') LIKE ? OR COALESCE(s.supplier_name, '') LIKE ?
           OR so.purpose LIKE ? OR so.requested_by LIKE ? OR so.released_by LIKE ? OR so.status LIKE ? OR so.remarks LIKE ?
        ORDER BY so.created_at DESC
        LIMIT 50
        """,
        (like, like, like, like, like, like, like, like, like),
    )

    frames = [df for df in [products, suppliers, sample_in, sample_out] if not df.empty]
    if not frames:
        return pd.DataFrame()
    return pd.concat(frames, ignore_index=True)


def render_clickable_kpi_card(title, value, note, icon, accent_class, query_key):
    """Render original premium KPI card with isolated hover wrapper.

    v1.0 keeps the original HTML KPI card design and safe Streamlit button
    click handling. Each KPI is placed in its own keyed container so hover only
    affects the matching KPI card.
    """
    try:
        kpi_container = st.container(key=f"kpi_wrap_{query_key}")
    except TypeError:
        kpi_container = st.container()

    with kpi_container:
        st.markdown(
            f"""
            <div class="kpi-card {accent_class}" role="button" aria-label="Open {title} details">
                <div class="kpi-icon">{icon}</div>
                <div class="kpi-title">{title}</div>
                <div class="kpi-value">{clean_number(value)}</div>
                <div class="kpi-note">{note}</div>
                <div class="kpi-chevron">›</div>
            </div>
            """,
            unsafe_allow_html=True,
        )

        clicked = st.button(
            f"Open {title} details",
            key=f"dashboard_kpi_btn_{query_key}",
            use_container_width=True,
        )

    if clicked:
        st.session_state["selected_dashboard_kpi"] = query_key
    return clicked


# -----------------------------
# TRANSACTION EDIT HELPERS
# -----------------------------

def _safe_date_value(value, fallback=None, allow_blank=False):
    """Return a safe date object for Streamlit date_input.

    Some imported rows carried placeholder expiry dates such as 2001-12-30.
    Treat those as blank/missing instead of showing them as real expiry dates.
    """
    if fallback is None and not allow_blank:
        fallback = date.today()
    if value in [None, "", "None", "nan", "NaT"]:
        return None if allow_blank else fallback
    dt = pd.to_datetime(value, errors="coerce")
    if pd.isna(dt):
        return None if allow_blank else fallback
    d = dt.date()
    if d <= date(2001, 12, 31):
        return None if allow_blank else fallback
    return d


def _date_to_db(value):
    """Save optional date values cleanly to SQLite."""
    return str(value) if value else ""


def stock_in_edit_section():
    st.markdown("<div class='section-title'>Edit Sample-In Record</div>", unsafe_allow_html=True)
    st.caption("Use this when you need to correct typo errors such as date received, reference number, quantity, batch number, received by, or remarks.")
    records = query_df(
        """
        SELECT si.stock_in_id, si.reference_no, si.date_received, si.product_id, p.product_name,
               si.batch_number, si.quantity, si.unit_cost, si.expiry_date, si.received_by, si.remarks
        FROM stock_in si
        LEFT JOIN products p ON si.product_id = p.product_id
        ORDER BY date(si.date_received) DESC, si.created_at DESC
        LIMIT 300
        """
    )
    if records.empty:
        st.info("No Sample-In records available to edit.")
        return
    records["label"] = records.apply(lambda r: f"{r['reference_no']} | {r['date_received']} | {r['product_name']} | Qty: {float(r['quantity'] or 0):g}", axis=1)
    label_to_id = dict(zip(records["label"], records["stock_in_id"]))
    selected_label = st.selectbox("Select Sample-In record to edit", records["label"].tolist(), key="edit_stock_in_select")
    stock_in_id = int(label_to_id[selected_label])
    old = query_df("SELECT * FROM stock_in WHERE stock_in_id=?", (stock_in_id,)).iloc[0]
    # Build product selector using product_id values, not labels.
    # This prevents stale Streamlit selectbox state and duplicate/old labels from
    # showing the wrong product/batch when a Sample-In record is selected.
    products_for_select = products_base_df().copy()
    if products_for_select.empty:
        product_ids = []
        product_labels = {}
    else:
        products_for_select["select_label"] = products_for_select.apply(
            lambda r: (
                f"{r['product_name']} | Batch: {r['batch_number'] or '-'} | "
                f"Pack Size: {r.get('pack_size', '') or '-'} | "
                f"Loc: {r.get('storage_location', '') or '-'} | "
                f"Qty: {clean_number(r['current_inventory'])} {str(r['pack_type'] or '').lower()}(s) | "
                f"ID: {int(r['product_id'])}"
            ),
            axis=1,
        )
        product_ids = [int(x) for x in products_for_select["product_id"].tolist()]
        product_labels = dict(zip(product_ids, products_for_select["select_label"].tolist()))

    old_product_id = int(old["product_id"]) if pd.notna(old["product_id"]) else (product_ids[0] if product_ids else None)
    product_index = product_ids.index(old_product_id) if old_product_id in product_ids else 0

    # If the stock-in expiry is blank, use the selected product/batch expiry as the
    # starting value. If both are blank, show the field as blank.
    old_product_row = products_for_select[products_for_select["product_id"].astype(int) == int(old_product_id)] if product_ids and old_product_id is not None else pd.DataFrame()
    product_expiry_value = old_product_row.iloc[0].get("expiry_date", "") if not old_product_row.empty else ""
    starting_expiry_value = old["expiry_date"] if str(old.get("expiry_date", "") or "").strip() else product_expiry_value

    with st.form(f"edit_stock_in_form_{stock_in_id}"):
        c1, c2 = st.columns(2)
        reference_no = c1.text_input("Reference No.", value=str(old["reference_no"] or ""), key=f"edit_stock_in_ref_{stock_in_id}")
        date_received = c2.date_input("Date Received", value=_safe_date_value(old["date_received"]), key=f"edit_stock_in_date_{stock_in_id}")
        selected_product_id = st.selectbox(
            "Product / Batch",
            product_ids,
            index=product_index,
            format_func=lambda pid: product_labels.get(int(pid), str(pid)),
            key=f"edit_stock_in_product_id_{stock_in_id}_{old_product_id}",
        ) if product_ids else None
        c3, c4, c5 = st.columns(3)
        quantity = c3.number_input("Quantity / No. of Packs", min_value=0.0, value=float(old["quantity"] or 0), step=0.25, key=f"edit_stock_in_qty_{stock_in_id}")
        unit_cost = c4.number_input("Unit Cost per Pack", min_value=0.0, value=float(old["unit_cost"] or 0), step=10.0, key=f"edit_stock_in_cost_{stock_in_id}")
        expiry_date = c5.date_input("Expiry Date", value=_safe_date_value(starting_expiry_value, allow_blank=True), key=f"edit_stock_in_expiry_{stock_in_id}_{old_product_id}")
        if expiry_date is None:
            c5.caption("No valid expiry date encoded.")
        c6, c7 = st.columns(2)
        batch_number = c6.text_input("Batch Number", value=str(old["batch_number"] or ""), key=f"edit_stock_in_batch_{stock_in_id}")
        received_by = c7.text_input("Received By", value=str(old["received_by"] or ""), key=f"edit_stock_in_received_by_{stock_in_id}")
        remarks = st.text_area("Remarks", value=str(old["remarks"] or ""), key=f"edit_stock_in_remarks_{stock_in_id}")
        save_edit = st.form_submit_button("Update Sample-In Record", use_container_width=True)

    if save_edit:
        if quantity <= 0:
            st.error("Quantity must be greater than zero.")
            return
        if not reference_no.strip():
            st.error("Reference No. is required.")
            return
        new_product_id = int(selected_product_id) if selected_product_id is not None else int(old["product_id"])
        conflict = query_df("SELECT COUNT(*) AS n FROM stock_in WHERE reference_no=? AND stock_in_id<>?", (reference_no.strip(), stock_in_id)).loc[0, "n"]
        if int(conflict) > 0:
            st.error("Reference No. already exists in another Sample-In record.")
            return
        old_qty = float(old["quantity"] or 0)
        old_product_id = int(old["product_id"])
        supplier_row = query_df("SELECT supplier_id FROM products WHERE product_id=?", (int(new_product_id),))
        new_supplier_id = int(supplier_row.loc[0, "supplier_id"]) if not supplier_row.empty and pd.notna(supplier_row.loc[0, "supplier_id"]) else None
        try:
            execute_transaction([
                ("UPDATE products SET quantity = quantity - ?, updated_at=CURRENT_TIMESTAMP WHERE product_id=?", (old_qty, old_product_id)),
                ("UPDATE products SET quantity = quantity + ?, updated_at=CURRENT_TIMESTAMP WHERE product_id=?", (float(quantity), int(new_product_id))),
                (
                    """
                    UPDATE stock_in SET reference_no=?, date_received=?, product_id=?, supplier_id=?, quantity=?, unit_cost=?,
                        batch_number=?, expiry_date=?, received_by=?, remarks=?
                    WHERE stock_in_id=?
                    """,
                    (reference_no.strip(), str(date_received), int(new_product_id), new_supplier_id, float(quantity), float(unit_cost),
                     batch_number.strip(), _date_to_db(expiry_date), received_by.strip(), remarks, stock_in_id),
                ),
            ])
        except Exception as exc:
            st.error(f"Update failed. No inventory changes were saved. Error: {exc}")
            return
        log_action("Sample In", "Edit", f"Edited {reference_no.strip()} from old qty {old_qty:g} to {float(quantity):g}", get_setting("current_user", "Admin User"))
        st.success("Sample-In record updated successfully and product balance recalculated.")
        rerun_app()


def stock_out_edit_section():
    st.markdown("<div class='section-title'>Edit Sample-Out Record</div>", unsafe_allow_html=True)
    st.caption("Use this when you need to correct typo errors such as release date, quantity, pullout by, status, or remarks.")
    records = query_df(
        """
        SELECT so.stock_out_id, so.reference_no, so.date_released, so.product_id, p.product_name,
               p.batch_number, so.quantity, so.status, so.released_by, so.remarks
        FROM stock_out so
        LEFT JOIN products p ON so.product_id = p.product_id
        ORDER BY date(so.date_released) DESC, so.created_at DESC
        LIMIT 300
        """
    )
    if records.empty:
        st.info("No Sample-Out records available to edit.")
        return
    records["label"] = records.apply(lambda r: f"{r['reference_no']} | {r['date_released']} | {r['product_name']} | Qty: {float(r['quantity'] or 0):g} | {r['status']}", axis=1)
    label_to_id = dict(zip(records["label"], records["stock_out_id"]))
    selected_label = st.selectbox("Select Sample-Out record to edit", records["label"].tolist(), key="edit_stock_out_select")
    stock_out_id = int(label_to_id[selected_label])
    old = query_df("SELECT * FROM stock_out WHERE stock_out_id=?", (stock_out_id,)).iloc[0]
    products_for_select = products_base_df().copy()
    if products_for_select.empty:
        product_ids = []
        product_labels = {}
    else:
        products_for_select["select_label"] = products_for_select.apply(
            lambda r: (
                f"{r['product_name']} | Batch: {r['batch_number'] or '-'} | "
                f"Pack Size: {r.get('pack_size', '') or '-'} | "
                f"Loc: {r.get('storage_location', '') or '-'} | "
                f"Qty: {clean_number(r['current_inventory'])} {str(r['pack_type'] or '').lower()}(s) | "
                f"ID: {int(r['product_id'])}"
            ),
            axis=1,
        )
        product_ids = [int(x) for x in products_for_select["product_id"].tolist()]
        product_labels = dict(zip(product_ids, products_for_select["select_label"].tolist()))
    old_product_id = int(old["product_id"]) if pd.notna(old["product_id"]) else (product_ids[0] if product_ids else None)
    product_index = product_ids.index(old_product_id) if old_product_id in product_ids else 0
    status_options = get_dropdown_options("sample_out_status")
    if str(old["status"] or "") not in status_options:
        status_options = [str(old["status"] or "Released")] + status_options

    selected_product_id = None
    if product_ids:
        selected_product_id = st.selectbox(
            "Product / Batch",
            product_ids,
            index=product_index,
            format_func=lambda pid: product_labels.get(int(pid), str(pid)),
            key=f"edit_stock_out_product_id_picker_{stock_out_id}_{old_product_id}",
            help="This selector is outside the form so the selected product details update immediately."
        )
        selected_row = get_product(int(selected_product_id))
        if selected_row is not None:
            st.info(
                f"Selected: {selected_row['product_name']} | Batch: {selected_row['batch_number']} | "
                f"Pack Size: {selected_row.get('pack_size', '') or '-'} | "
                f"Available: {clean_number(selected_row['quantity'])} {str(selected_row['pack_type']).lower()}(s) | "
                f"Status: {selected_row['final_status']}"
            )

    with st.form(f"edit_stock_out_form_{stock_out_id}"):
        c1, c2 = st.columns(2)
        reference_no = c1.text_input("Reference No.", value=str(old["reference_no"] or ""), key=f"edit_stock_out_ref_{stock_out_id}")
        date_released = c2.date_input("Date Released", value=_safe_date_value(old["date_released"]), key=f"edit_stock_out_date_{stock_out_id}")
        c3, c4, c5 = st.columns(3)
        quantity = c3.number_input("Quantity Released / No. of Packs", min_value=0.0, value=float(old["quantity"] or 0), step=0.25, key=f"edit_stock_out_qty_{stock_out_id}")
        purpose_options = get_dropdown_options("sample_out_purpose")
        old_purpose = str(old["purpose"] or "") if "purpose" in old.index else ""
        if old_purpose and old_purpose not in purpose_options:
            purpose_options = [old_purpose] + purpose_options
        purpose_index = purpose_options.index(old_purpose) if old_purpose in purpose_options else 0
        purpose = c4.selectbox("Purpose", purpose_options, index=purpose_index, key=f"edit_stock_out_purpose_{stock_out_id}")
        status = c5.selectbox("Status", status_options, index=(status_options.index(str(old["status"] or "Released")) if str(old["status"] or "Released") in status_options else 0), key=f"edit_stock_out_status_{stock_out_id}")
        c6, c7, c8 = st.columns(3)
        requested_by = c6.text_input("Requested By", value=str(old["requested_by"] or ""), key=f"edit_stock_out_requested_by_{stock_out_id}")
        released_by = c7.text_input("Released By", value=str(old["released_by"] or ""), key=f"edit_stock_out_released_by_{stock_out_id}")
        approved_by = c8.text_input("Approved By", value=str(old["approved_by"] or ""), key=f"edit_stock_out_approved_by_{stock_out_id}")
        remarks = st.text_area("Remarks", value=str(old["remarks"] or ""), key=f"edit_stock_out_remarks_{stock_out_id}")
        save_edit = st.form_submit_button("Update Sample-Out Record", use_container_width=True)

    if save_edit:
        if quantity <= 0:
            st.error("Quantity must be greater than zero.")
            return
        if not reference_no.strip():
            st.error("Reference No. is required.")
            return
        new_product_id = int(selected_product_id) if selected_product_id is not None else int(old["product_id"])
        conflict = query_df("SELECT COUNT(*) AS n FROM stock_out WHERE reference_no=? AND stock_out_id<>?", (reference_no.strip(), stock_out_id)).loc[0, "n"]
        if int(conflict) > 0:
            st.error("Reference No. already exists in another Sample-Out record.")
            return
        old_qty = float(old["quantity"] or 0)
        old_product_id = int(old["product_id"])
        old_status = str(old["status"] or "")
        # Validate the new released quantity against the balance after reversing the old effect.
        available_now = query_df("SELECT COALESCE(quantity,0) AS qty FROM products WHERE product_id=?", (int(new_product_id),)).loc[0, "qty"]
        available_after_reverse = float(available_now) + (old_qty if old_status == "Released" and old_product_id == int(new_product_id) else 0)
        if status == "Released" and float(quantity) > float(available_after_reverse):
            st.error(f"Cannot save. Available balance after correction is only {float(available_after_reverse):g} pack(s).")
            return

        operations = []
        if old_status == "Released":
            operations.append(("UPDATE products SET quantity = quantity + ?, updated_at=CURRENT_TIMESTAMP WHERE product_id=?", (old_qty, old_product_id)))
        if status == "Released":
            operations.append(("UPDATE products SET quantity = quantity - ?, updated_at=CURRENT_TIMESTAMP WHERE product_id=?", (float(quantity), int(new_product_id))))
        operations.append((
            """
            UPDATE stock_out SET reference_no=?, date_released=?, product_id=?, quantity=?, purpose=?,
                requested_by=?, released_by=?, approved_by=?, status=?, remarks=?
            WHERE stock_out_id=?
            """,
            (reference_no.strip(), str(date_released), int(new_product_id), float(quantity), purpose,
             requested_by.strip(), released_by.strip(), approved_by.strip(), status, remarks, stock_out_id),
        ))
        try:
            execute_transaction(operations)
        except Exception as exc:
            st.error(f"Update failed. No inventory changes were saved. Error: {exc}")
            return
        log_action("Sample Out", "Edit", f"Edited {reference_no.strip()} from old qty {old_qty:g} to {float(quantity):g}", get_setting("current_user", "Admin User"))
        st.success("Sample-Out record updated successfully and product balance recalculated.")
        rerun_app()

# -----------------------------
# PAGES
# -----------------------------

def dashboard_page():
    dashboard_search = dashboard_visual_header()
    df = products_base_df()
    today = date.today()
    month_start = today.replace(day=1)
    next_month_start = (pd.Timestamp(month_start) + pd.DateOffset(months=1)).date()

    stock_in_month_qty = query_df("SELECT COALESCE(SUM(quantity), 0) AS qty FROM stock_in WHERE date(date_received) >= date(?) AND date(date_received) < date(?)", (str(month_start), str(next_month_start))).loc[0, "qty"]
    stock_out_month_qty = query_df("SELECT COALESCE(SUM(quantity), 0) AS qty FROM stock_out WHERE date(date_released) >= date(?) AND date(date_released) < date(?) AND status = 'Released'", (str(month_start), str(next_month_start))).loc[0, "qty"]
    stock_in_month = query_df("SELECT COUNT(*) AS n FROM stock_in WHERE date(date_received) >= date(?) AND date(date_received) < date(?)", (str(month_start), str(next_month_start))).loc[0, "n"]
    stock_out_month = query_df("SELECT COUNT(*) AS n FROM stock_out WHERE date(date_released) >= date(?) AND date(date_released) < date(?) AND status = 'Released'", (str(month_start), str(next_month_start))).loc[0, "n"]

    usable_stock_df = df[(df["quantity"] > 0) & (df["expiry_status"] != "Expired")]
    expiry_alert_days = get_int_setting("expiry_alert_days", 90)
    expiring_df = df[(df["expiry_status"] != "Expired") & (df["days_to_expiry"].notna()) & (df["days_to_expiry"] <= expiry_alert_days)]
    expired_records_df = df[df["expiry_status"] == "Expired"]
    stock_in_df = query_df(
        """
        SELECT si.reference_no, si.date_received, p.product_name, p.batch_number, si.quantity,
               p.pack_type, s.supplier_name AS supplier, si.received_by, si.remarks
        FROM stock_in si
        LEFT JOIN products p ON si.product_id = p.product_id
        LEFT JOIN suppliers s ON si.supplier_id = s.supplier_id
        WHERE date(si.date_received) >= date(?)
          AND date(si.date_received) < date(?)
        ORDER BY si.date_received DESC
        """,
        (str(month_start), str(next_month_start))
    )
    stock_out_df = query_df(
        """
        SELECT so.reference_no, so.date_released, p.product_name, p.batch_number, so.quantity,
               p.pack_type, so.purpose, so.requested_by, so.released_by, so.status
        FROM stock_out so
        LEFT JOIN products p ON so.product_id = p.product_id
        WHERE date(so.date_released) >= date(?)
          AND date(so.date_released) < date(?)
          AND so.status = 'Released'
        ORDER BY so.date_released DESC
        """,
        (str(month_start), str(next_month_start))
    )

    if dashboard_search:
        search_results = build_dashboard_search_results(dashboard_search)
        with st.container(border=True):
            chart_card_heading("Search Results", f"Results for: {dashboard_search}")

            if search_results is None or search_results.empty:
                st.info("No matching product/batch records found.")
            else:
                # v1.0 dashboard search fix:
                # Apply the dashboard filter BEFORE KPI cards and BEFORE the table.
                # Default rule: Product / Batch only + Current Qty > 0.
                include_out_of_stock = st.checkbox(
                    "Include out-of-stock records",
                    value=False,
                    key="dashboard_search_include_out_of_stock_v1.0",
                    help="Default OFF. Dashboard search shows only Product / Batch records with Current Qty greater than 0.",
                )

                filtered_search_results = search_results.copy()

                # 1. Keep Product / Batch rows only.
                # This removes Supplier, Sample In, and Sample Out rows from the dashboard search results.
                if "Result Type" in filtered_search_results.columns:
                    filtered_search_results = filtered_search_results[
                        filtered_search_results["Result Type"]
                        .astype(str)
                        .str.strip()
                        .str.lower()
                        .eq("product / batch")
                    ].copy()

                # 2. Convert Current Qty to numeric and hide zero-stock rows by default.
                if "Current Qty" in filtered_search_results.columns:
                    filtered_search_results["_current_qty_numeric"] = pd.to_numeric(
                        filtered_search_results["Current Qty"],
                        errors="coerce",
                    ).fillna(0)

                    if not include_out_of_stock:
                        filtered_search_results = filtered_search_results[
                            filtered_search_results["_current_qty_numeric"] > 0
                        ].copy()

                # 3. Add dashboard status.
                # Expired products are NOT automatically hidden if they still have stock.
                def _dashboard_search_status(row):
                    qty = float(row.get("_current_qty_numeric", 0) or 0)
                    if qty <= 0:
                        return "Out of Stock"

                    try:
                        expiry_dt = pd.to_datetime(row.get("Date", ""), errors="coerce")
                        if pd.notna(expiry_dt):
                            days_left = (expiry_dt.date() - date.today()).days
                            if days_left < 0:
                                return "Expired"
                            if days_left <= expiry_alert_days:
                                return "Expiring Soon"
                    except Exception:
                        pass

                    return "Available"

                filtered_search_results["Status"] = filtered_search_results.apply(_dashboard_search_status, axis=1)
                filtered_search_results = filtered_search_results.drop(columns=["_current_qty_numeric"], errors="ignore")

                if include_out_of_stock:
                    st.warning("Available stock only is OFF. Showing Product / Batch records including out-of-stock.")
                else:
                    st.info("Showing Product / Batch records with available stock only. Expired products with remaining stock are still included.")

                # 4. KPI cards must use the filtered result, not the old combined search_results.
                r1, r2, r3, r4, r5 = st.columns(5)
                r1.metric("Total Results", len(filtered_search_results))
                r2.metric("Products", len(filtered_search_results))
                r3.metric(
                    "Suppliers",
                    int(filtered_search_results["Supplier"].nunique()) if "Supplier" in filtered_search_results.columns and not filtered_search_results.empty else 0,
                )
                r4.metric("Sample In", 0)
                r5.metric("Sample Out", 0)

                # 5. Table must display filtered_search_results, not search_results.
                display_columns = [
                    "Result Type",
                    "Product",
                    "Batch No.",
                    "Supplier",
                    "Category",
                    "Current Qty",
                    "Pack Size",
                    "Pack Type",
                    "Date",
                    "Status",
                    "Storage / Purpose",
                    "Remarks",
                ]
                display_columns = [col for col in display_columns if col in filtered_search_results.columns]
                filtered_search_results = filtered_search_results[display_columns].reset_index(drop=True)

                if filtered_search_results.empty:
                    st.info("No available Product / Batch records found. Check 'Include out-of-stock records' if you also want to see zero-stock batches.")
                else:
                    display_df(filtered_search_results, height=360, table_key="dashboard_search_available_only_v1.0")

    cols = st.columns(5)
    card_items = [
        ("Usable Stock", len(usable_stock_df), "With stock + not expired", "✅", "kpi-available", "available_products", usable_stock_df),
        ("Expiring Soon", len(expiring_df), f"Within {expiry_alert_days} days", "⏳", "kpi-expiring", "expiring_soon", expiring_df),
        ("Expired Records", len(expired_records_df), "All expired batches", "⚠️", "kpi-expired", "expired_products", expired_records_df),
        ("Received This Month", int(stock_in_month), "Receiving transactions", "⬇️", "kpi-received", "stock_in_month", stock_in_df),
        ("Released This Month", int(stock_out_month), "Released transactions", "⬆️", "kpi-out", "stock_out_month", stock_out_df),
    ]
    for col, (title, value, note, icon, accent_class, key, details_df) in zip(cols, card_items):
        with col:
            render_clickable_kpi_card(title, value, note, icon, accent_class, key)

    selected_kpi = st.session_state.get("selected_dashboard_kpi", "")
    kpi_lookup = {key: (title, details_df) for title, value, note, icon, accent_class, key, details_df in card_items}
    if selected_kpi in kpi_lookup:
        selected_title, selected_df = kpi_lookup[selected_kpi]
        st.session_state["selected_dashboard_kpi"] = ""
        open_details(selected_title, selected_df)

    # Visual spacing between the KPI card row and the chart section.
    st.markdown("<div style='height: 28px;'></div>", unsafe_allow_html=True)

    c1, c2 = st.columns(2)

    with c1:
        with st.container(border=True):
            chart_card_heading("Product Sample Category", "Count of available product samples by category")
            if df.empty:
                st.info("No product records yet.")
            else:
                available_category_df = df[(df["quantity"].astype(float) > 0) & (df["expiry_status"] == "Not Expired")]
                cat_df = available_category_df.groupby("category", dropna=False)["product_id"].count().reset_index(name="available_products")
                chart_type, palette_name, color_sequence = get_chart_pref("category")
                fig = make_custom_chart(cat_df, "category", "available_products", "", chart_type, color_sequence)
                if fig:
                    st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})
                    chart_note("Shows the distribution of product sample status based on current records.")
                    chart_note("Counts available product records only.")

    with c2:
        with st.container(border=True):
            chart_card_heading("Product Status Summary", "Distribution of product samples by status")
            if df.empty:
                st.info("No product status data yet.")
            else:
                status_df = df.groupby("final_status")["product_id"].count().reset_index(name="count")
                chart_type, palette_name, color_sequence = get_chart_pref("status")
                fig = make_custom_chart(status_df, "final_status", "count", "", chart_type, color_sequence)
                if fig:
                    st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})

    with st.container(border=True):
        chart_card_heading(
            "Available Product/Batch Samples by Supplier Category",
            "Counts available product/batch records by supplier and product sample category"
        )
        if df.empty:
            st.info("No product records yet.")
        else:
            supplier_category_df = df[
                (pd.to_numeric(df["quantity"], errors="coerce").fillna(0) > 0)
                & (df["expiry_status"] != "Expired")
            ].copy()
            if supplier_category_df.empty:
                st.info("No available product/batch samples to summarize by supplier category.")
            else:
                supplier_category_df["supplier"] = supplier_category_df["supplier"].fillna("No Supplier").replace("", "No Supplier")
                supplier_category_df["category"] = supplier_category_df["category"].fillna("Uncategorized").replace("", "Uncategorized")
                supplier_category_summary = (
                    supplier_category_df
                    .groupby(["supplier", "category"], dropna=False)["product_id"]
                    .count()
                    .reset_index(name="product_batch_count")
                )
                supplier_totals = (
                    supplier_category_summary
                    .groupby("supplier", as_index=False)["product_batch_count"]
                    .sum()
                    .sort_values("product_batch_count", ascending=False)
                )
                top_suppliers = supplier_totals.head(12)["supplier"].tolist()
                supplier_category_summary = supplier_category_summary[supplier_category_summary["supplier"].isin(top_suppliers)]
                supplier_category_summary["supplier"] = pd.Categorical(
                    supplier_category_summary["supplier"], categories=list(reversed(top_suppliers)), ordered=True
                )
                supplier_category_summary["label_value"] = supplier_category_summary["product_batch_count"].apply(clean_number)
                _, _, color_sequence = get_chart_pref("category")
                fig = px.bar(
                    supplier_category_summary,
                    x="product_batch_count",
                    y="supplier",
                    color="category",
                    orientation="h",
                    text="label_value",
                    color_discrete_sequence=color_sequence,
                    labels={
                        "product_batch_count": "Product/Batch Count",
                        "supplier": "Supplier",
                        "category": "Product Sample Category",
                    },
                )
                fig.update_traces(
                    textposition="inside",
                    textfont=dict(color="white", size=11),
                    textangle=0,
                    insidetextanchor="middle",
                    marker_line_color="rgba(255,255,255,.85)",
                    marker_line_width=1,
                    hovertemplate="Supplier: %{y}<br>Count: %{x:.0f}<br>Category: %{legendgroup}<extra></extra>",
                )
                fig.update_layout(
                    barmode="stack",
                    yaxis={"categoryorder": "array", "categoryarray": list(reversed(top_suppliers))},
                    legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1),
                    height=max(360, 42 * max(4, len(top_suppliers))),
                )
                fig = apply_chart_layout(fig)
                fig.update_xaxes(title_text="Available Product/Batch Count", tickformat=",.0f")
                fig.update_yaxes(title_text="Supplier")
                st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})
                chart_note("Each available product/batch record is counted as 1, regardless of quantity or pack size.")
                chart_note("Expired and out-of-stock records are excluded. Only the top 12 suppliers are shown for readability.")

    with st.container(border=True):
        chart_card_heading("Recently Received Samples — 70 Days", "Latest received product samples with supplier, quantity, and remarks")
        seventy_days_ago = today - timedelta(days=70)
        recent_received_70 = query_df(
            """
            SELECT si.date_received AS "Date Received",
                   p.product_name AS "Product",
                   s.supplier_name AS "Supplier",
                   si.quantity AS qty,
                   p.pack_type AS pack_type,
                   si.remarks AS "Remarks"
            FROM stock_in si
            LEFT JOIN products p ON si.product_id = p.product_id
            LEFT JOIN suppliers s ON si.supplier_id = s.supplier_id
            WHERE date(si.date_received) >= date(?)
            ORDER BY date(si.date_received) DESC, si.created_at DESC
            """,
            (str(seventy_days_ago),),
        )
        if recent_received_70.empty:
            st.info("No samples received in the last 70 days.")
        else:
            recent_received_70["Quantity"] = recent_received_70.apply(
                lambda r: f"{clean_number(r['qty'])} {str(r['pack_type'] or '').strip()}".strip(), axis=1
            )
            display_df(recent_received_70[["Date Received", "Product", "Supplier", "Quantity", "Remarks"]], height=320)

    with st.container(border=True):
        chart_card_heading("Samples Received — Last 12 Months", "Smooth trend of receiving transactions per month")
        month_start_12 = (pd.Timestamp(today.replace(day=1)) - pd.DateOffset(months=11)).date()
        month_end = (pd.Timestamp(today.replace(day=1)) + pd.DateOffset(months=1)).date()
        monthly_received = query_df(
            """
            SELECT strftime('%Y-%m', date_received) AS month_key,
                   COUNT(*) AS receiving_transactions
            FROM stock_in
            WHERE date(date_received) >= date(?)
              AND date(date_received) < date(?)
            GROUP BY strftime('%Y-%m', date_received)
            ORDER BY month_key
            """,
            (str(month_start_12), str(month_end)),
        )
        month_index = pd.period_range(start=pd.Period(month_start_12, freq="M"), periods=12, freq="M")
        month_df = pd.DataFrame({"month_key": month_index.astype(str)})
        if monthly_received.empty:
            month_df["receiving_transactions"] = 0
        else:
            monthly_received["receiving_transactions"] = pd.to_numeric(monthly_received["receiving_transactions"], errors="coerce").fillna(0)
            month_df = month_df.merge(monthly_received, on="month_key", how="left").fillna({"receiving_transactions": 0})
        month_df["Month"] = pd.to_datetime(month_df["month_key"] + "-01").dt.strftime("%b %Y")
        chart_type, palette_name, color_sequence = get_chart_pref("expiry")
        fig = make_custom_chart(month_df, "Month", "receiving_transactions", "", chart_type, color_sequence)
        if fig:
            fig.update_yaxes(title_text="Receiving transactions")
            st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})
            chart_note("Smooth area trend counts Sample-In transaction records for the current month and previous 11 months.")
        else:
            st.info("No sample received data available.")

    st.markdown("<div class='section-title'>Recent Activity</div>", unsafe_allow_html=True)
    st.caption("Recent Activity is based on the actual transaction dates: Date Received for Sample In and Date Released for Sample Out.")
    recent_in = query_df(
        """
        SELECT
            si.date_received AS transaction_date,
            si.reference_no,
            'Sample In' AS type,
            p.product_name,
            si.quantity,
            p.pack_type,
            si.remarks,
            si.created_at AS encoded_at
        FROM stock_in si
        LEFT JOIN products p ON si.product_id = p.product_id
        ORDER BY date(si.date_received) DESC, si.created_at DESC
        LIMIT 8
        """
    )
    recent_out = query_df(
        """
        SELECT
            so.date_released AS transaction_date,
            so.reference_no,
            'Sample Out' AS type,
            p.product_name,
            so.quantity,
            p.pack_type,
            so.remarks,
            so.created_at AS encoded_at
        FROM stock_out so
        LEFT JOIN products p ON so.product_id = p.product_id
        ORDER BY date(so.date_released) DESC, so.created_at DESC
        LIMIT 8
        """
    )
    recent = pd.concat([recent_in, recent_out], ignore_index=True)
    if not recent.empty:
        recent["_sort_date"] = pd.to_datetime(recent["transaction_date"], errors="coerce")
        recent["_sort_encoded"] = pd.to_datetime(recent["encoded_at"], errors="coerce")
        recent = recent.sort_values(["_sort_date", "_sort_encoded"], ascending=[False, False]).head(8)
        recent = recent.drop(columns=["_sort_date", "_sort_encoded", "encoded_at"], errors="ignore")
    display_df(recent, height=260)


def products_page():
    page_header("Products", "Register and manage product/batch records with beginning inventory, sample in, sample out, and current inventory.")
    st.markdown(
        """
        <div class="module-hero-card">
            <div class="module-hero-title">Products Module</div>
            <div class="module-hero-subtitle">Manage all product and batch records in one place. Review inventory position, maintain batch details, and keep product information clean and organized without changing the core inventory logic.</div>
        </div>
        """,
        unsafe_allow_html=True,
    )

    tabs = st.tabs(["Product List", "Add Product", "Edit / Delete", "Movement History"])

    with tabs[0]:
        df = products_base_df()
        total_products = len(df)
        available_products = len(df[df["final_status"] == "Available"]) if not df.empty else 0
        expiring_products = len(df[df["final_status"].isin(["Expiring Soon", "Urgent Expiry"])]) if not df.empty else 0
        out_products = len(df[df["final_status"] == "Out of Stock"]) if not df.empty else 0

        s1, s2, s3, s4 = st.columns(4)
        stats = [
            (s1, "Product / Batch Records", total_products, "All registered product batches"),
            (s2, "Available", available_products, "Ready for use"),
            (s3, "Expiring Soon", expiring_products, "Needs monitoring"),
            (s4, "Out of Stock", out_products, "For replenishment"),
        ]
        for col, label, value, note in stats:
            with col:
                st.markdown(
                    f"""
                    <div class="mini-stat-card">
                        <div class="mini-stat-label">{label}</div>
                        <div class="mini-stat-value">{value}</div>
                        <div class="mini-stat-note">{note}</div>
                    </div>
                    """,
                    unsafe_allow_html=True,
                )

        st.markdown("<div style='height: 0.45rem;'></div>", unsafe_allow_html=True)

        with st.container(border=True):
            st.markdown("<div class='filter-shell-title'>Filter Products</div><div class='filter-shell-subtitle'>Search and narrow down the product list by category, supplier, and inventory status.</div>", unsafe_allow_html=True)
            f1, f2, f3, f4 = st.columns(4)
            search = f1.text_input("Search product", placeholder="Delfan, batch no., active ingredient...")
            category = f2.selectbox("Category", ["All"] + sorted([x for x in df["category"].dropna().unique()]) if not df.empty else ["All"])
            supplier = f3.selectbox("Supplier", ["All"] + sorted([x for x in df["supplier"].dropna().unique()]) if not df.empty else ["All"])
            status = f4.selectbox("Status", ["All", "Available", "Expiring Soon", "Urgent Expiry", "Expired", "Out of Stock"])

        view = df.copy()
        if search and not view.empty:
            mask = view.apply(lambda row: row_matches_search(row, search), axis=1)
            view = view[mask]
        if category != "All" and not view.empty:
            view = view[view["category"] == category]
        if supplier != "All" and not view.empty:
            view = view[view["supplier"] == supplier]
        if status != "All" and not view.empty:
            view = view[view["final_status"] == status]

        display_cols = [
            "product_code", "product_name", "active_ingredient", "category", "supplier", "pack_type",
            "pack_size", "beginning_inventory", "sample_in", "sample_out", "current_inventory", "total_content", "batch_number", "expiry_date", "final_status", "storage_location"
        ]
        product_display = view[display_cols].copy() if not view.empty else view.copy()
        product_display = product_display.rename(columns={
            "product_code": "Product Code",
            "product_name": "Product Name",
            "active_ingredient": "Active Ingredient",
            "category": "Category",
            "supplier": "Supplier",
            "pack_type": "Pack Type",
            "pack_size": "Pack Size",
            "beginning_inventory": "Beginning Inventory",
            "sample_in": "Sample In",
            "sample_out": "Sample Out",
            "current_inventory": "Current Inventory",
            "total_content": "Total Content",
            "batch_number": "Batch Number",
            "expiry_date": "Expiry Date",
            "final_status": "Status",
            "storage_location": "Storage Location",
        })

        with st.container(border=True):
            st.markdown("<div class='form-shell-title'>Product Inventory List</div><div class='form-shell-subtitle'>Viewing product information, inventory summary, and batch-level details.</div>", unsafe_allow_html=True)
            visible_product_display = display_df(product_display, height=520, table_key="product_list")
            st.download_button("Download Product List CSV", to_csv_bytes(visible_product_display), "products.csv", "text/csv")

    with tabs[1]:
        supplier_map, suppliers = get_supplier_options()
        if not suppliers:
            st.warning("Please add a supplier first in the Suppliers module.")
        st.markdown("<div class='form-shell-title'>Add Product</div><div class='form-shell-subtitle'>Create a new product and batch record with packaging, inventory, and storage details.</div>", unsafe_allow_html=True)
        with st.form("add_product_form", clear_on_submit=True):
            st.markdown("<div class='form-section-heading'>A. Product Information</div>", unsafe_allow_html=True)
            c1, c2 = st.columns(2)
            product_code = c1.text_input("Product Code", value=generate_code("PRD", "products", "product_id"))
            product_name = c2.text_input("Product Name *", placeholder="Delfan Plus")
            c3, c4 = st.columns(2)
            active_ingredient = c3.text_input("Active Ingredient", placeholder="L-Alpha Free Amino Acids")
            category_options = get_dropdown_options("category")
            pack_options = get_dropdown_options("pack_type")
            unit_options = get_dropdown_options("pack_size_unit")
            category = c4.selectbox("Category", category_options)

            st.markdown("<div class='form-section-heading'>B. Supplier and Packaging</div>", unsafe_allow_html=True)
            c5, c6, c7 = st.columns(3)
            supplier_name = c5.selectbox("Supplier", suppliers if suppliers else ["No Supplier"])
            pack_type = c6.selectbox("Pack Type", pack_options)
            pack_size_value = c7.number_input("Pack Size Value", min_value=0.0, value=1.0, step=0.25)
            c8, c9, c10 = st.columns(3)
            pack_size_unit = c8.selectbox("Pack Size Unit", unit_options)
            beginning_inventory = c9.number_input("Beginning Inventory / No. of Packs", min_value=0.0, value=0.0, step=1.0)
            unit_cost = c10.number_input("Unit Cost per Pack", min_value=0.0, value=0.0, step=10.0)

            st.markdown("<div class='form-section-heading'>C. Batch and Storage Details</div>", unsafe_allow_html=True)
            c11, c12 = st.columns(2)
            batch_number = c11.text_input("Batch Number", placeholder="DP-2026-001")
            expiry_date = c12.date_input("Expiry Date", value=date.today() + timedelta(days=365))
            storage_options = get_dropdown_options("storage_location")
            storage_location = st.selectbox("Storage Location", storage_options)
            remarks = st.text_area("Remarks")
            image = st.file_uploader("Product Image / Label", type=["png", "jpg", "jpeg", "webp"])
            submitted = st.form_submit_button("Save Product", use_container_width=True)

        if submitted:
            if not product_name.strip():
                st.error("Product Name is required.")
            else:
                image_path = save_upload(image, "product_images")
                supplier_id = supplier_map.get(supplier_name)
                try:
                    execute(
                        """
                        INSERT INTO products
                        (product_code, product_name, active_ingredient, category, supplier_id, pack_type,
                         pack_size_value, pack_size_unit, beginning_inventory, quantity, unit_cost, batch_number, expiry_date,
                         storage_location, remarks, image_path)
                        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                        """,
                        (product_code, product_name, active_ingredient, category, supplier_id, pack_type,
                         pack_size_value, pack_size_unit, beginning_inventory, beginning_inventory, unit_cost, batch_number, str(expiry_date),
                         storage_location, remarks, image_path),
                    )
                    log_action("Products", "Create", f"Added {product_name} / batch {batch_number}")
                    st.success("Product saved successfully.")
                except sqlite3.IntegrityError:
                    st.error("Product code already exists. Please use another code.")

    with tabs[2]:
        product_map, options = get_product_options()
        if not options:
            st.info("No products available to edit.")
        else:
            selected = st.selectbox("Select product/batch to edit", options)
            product_id = product_map[selected]
            row = get_product(product_id)
            supplier_map, suppliers = get_supplier_options()
            current_supplier = row["supplier"] if row["supplier"] in suppliers else (suppliers[0] if suppliers else "No Supplier")
            st.markdown(
                f"""
                <div class="product-summary-card">
                    <div class="product-summary-title">{row['product_name']} — Batch {row['batch_number']}</div>
                    <div class="product-summary-subtitle">Review and update the product information below while keeping inventory calculations intact.</div>
                    <div class="product-summary-grid">
                        <div class="product-summary-chip"><div class="product-summary-chip-label">Category</div><div class="product-summary-chip-value">{row['category']}</div></div>
                        <div class="product-summary-chip"><div class="product-summary-chip-label">Supplier</div><div class="product-summary-chip-value">{row['supplier']}</div></div>
                        <div class="product-summary-chip"><div class="product-summary-chip-label">Current Inventory</div><div class="product-summary-chip-value">{float(row['current_inventory'] or 0):g} {str(row['pack_type']).lower()}(s)</div></div>
                        <div class="product-summary-chip"><div class="product-summary-chip-label">Status</div><div class="product-summary-chip-value">{row['final_status']}</div></div>
                    </div>
                </div>
                """,
                unsafe_allow_html=True,
            )
            with st.form("edit_product_form"):
                st.markdown("<div class='form-section-heading'>A. Product Information</div>", unsafe_allow_html=True)
                c1, c2 = st.columns(2)
                product_code = c1.text_input("Product Code", value=row["product_code"])
                product_name = c2.text_input("Product Name *", value=row["product_name"])
                c3, c4 = st.columns(2)
                active_ingredient = c3.text_input("Active Ingredient", value=row["active_ingredient"] or "")
                category_options = get_dropdown_options("category")
                category = c4.selectbox("Category", category_options, index=category_options.index(row["category"]) if row["category"] in category_options else 0)

                st.markdown("<div class='form-section-heading'>B. Supplier and Packaging</div>", unsafe_allow_html=True)
                c5, c6, c7 = st.columns(3)
                supplier_name = c5.selectbox("Supplier", suppliers if suppliers else ["No Supplier"], index=(suppliers.index(current_supplier) if current_supplier in suppliers else 0))
                pack_options = get_dropdown_options("pack_type")
                pack_type = c6.selectbox("Pack Type", pack_options, index=pack_options.index(row["pack_type"]) if row["pack_type"] in pack_options else 0)
                pack_size_value = c7.number_input("Pack Size Value", min_value=0.0, value=float(row["pack_size_value"] or 0), step=0.25)
                c8, c9, c10 = st.columns(3)
                unit_options = get_dropdown_options("pack_size_unit")
                pack_size_unit = c8.selectbox("Pack Size Unit", unit_options, index=unit_options.index(row["pack_size_unit"]) if row["pack_size_unit"] in unit_options else 0)
                beginning_inventory = c9.number_input("Beginning Inventory / No. of Packs", min_value=0.0, value=float(row["beginning_inventory"] or 0), step=1.0)
                c9.caption(f"Current Inventory: {float(row['current_inventory'] or 0):g} {str(row['pack_type']).lower()}(s)")
                unit_cost = c10.number_input("Unit Cost per Pack", min_value=0.0, value=float(row["unit_cost"] or 0), step=10.0)

                st.markdown("<div class='form-section-heading'>C. Batch and Storage Details</div>", unsafe_allow_html=True)
                c11, c12 = st.columns(2)
                batch_number = c11.text_input("Batch Number", value=row["batch_number"] or "")
                expiry_val = pd.to_datetime(row["expiry_date"], errors="coerce")
                expiry_date = c12.date_input("Expiry Date", value=expiry_val.date() if pd.notna(expiry_val) else date.today())
                storage_options = get_dropdown_options("storage_location")
                current_storage = row["storage_location"] if row["storage_location"] in storage_options else (storage_options[0] if storage_options else "")
                storage_location = st.selectbox("Storage Location", storage_options if storage_options else [current_storage], index=(storage_options.index(current_storage) if current_storage in storage_options else 0))
                remarks = st.text_area("Remarks", value=row["remarks"] or "")
                csave, cdelete = st.columns(2)
                save = csave.form_submit_button("Update Product", use_container_width=True)
                delete = cdelete.form_submit_button("Delete Product", use_container_width=True)

            if save:
                supplier_id = supplier_map.get(supplier_name)
                totals = query_df(
                    """
                    SELECT
                        COALESCE((SELECT SUM(quantity) FROM stock_in WHERE product_id=?), 0) AS sample_in,
                        COALESCE((SELECT SUM(quantity) FROM stock_out WHERE product_id=? AND status='Released'), 0) AS sample_out
                    """,
                    (product_id, product_id),
                )
                new_current_inventory = float(beginning_inventory or 0) + float(totals.loc[0, "sample_in"] or 0) - float(totals.loc[0, "sample_out"] or 0)
                if new_current_inventory < 0:
                    st.warning("The calculated current inventory is below zero. It will be saved as 0. Please review Sample Out records.")
                    new_current_inventory = 0
                execute(
                    """
                    UPDATE products SET
                        product_code=?, product_name=?, active_ingredient=?, category=?, supplier_id=?, pack_type=?,
                        pack_size_value=?, pack_size_unit=?, beginning_inventory=?, quantity=?, unit_cost=?, batch_number=?, expiry_date=?,
                        storage_location=?, remarks=?, updated_at=CURRENT_TIMESTAMP
                    WHERE product_id=?
                    """,
                    (product_code, product_name, active_ingredient, category, supplier_id, pack_type,
                     pack_size_value, pack_size_unit, beginning_inventory, new_current_inventory, unit_cost, batch_number, str(expiry_date),
                     storage_location, remarks, product_id),
                )
                log_action("Products", "Update", f"Updated {product_name} / batch {batch_number}")
                st.success("Product updated successfully.")
                rerun_app()

            if delete:
                execute("DELETE FROM products WHERE product_id=?", (product_id,))
                log_action("Products", "Delete", f"Deleted product ID {product_id}")
                st.success("Product deleted successfully.")
                rerun_app()

    with tabs[3]:
        product_map, options = get_product_options()
        if not options:
            st.info("No products available.")
        else:
            selected = st.selectbox("Select product/batch", options, key="movement_select")
            product_id = product_map[selected]
            product = get_product(product_id)
            st.markdown(
                f"""
                <div class="product-summary-card">
                    <div class="product-summary-title">{product['product_name']} — Batch {product['batch_number']}</div>
                    <div class="product-summary-subtitle">Quick overview of product movement and inventory position.</div>
                    <div class="product-summary-grid">
                        <div class="product-summary-chip"><div class="product-summary-chip-label">Pack</div><div class="product-summary-chip-value">{product['pack_type']} | {product['pack_size']}</div></div>
                        <div class="product-summary-chip"><div class="product-summary-chip-label">Supplier</div><div class="product-summary-chip-value">{product['supplier']}</div></div>
                        <div class="product-summary-chip"><div class="product-summary-chip-label">Expiry Date</div><div class="product-summary-chip-value">{product['expiry_date']}</div></div>
                        <div class="product-summary-chip"><div class="product-summary-chip-label">Storage</div><div class="product-summary-chip-value">{product['storage_location']}</div></div>
                    </div>
                </div>
                """,
                unsafe_allow_html=True,
            )
            c1, c2, c3, c4, c5 = st.columns(5)
            c1.metric("Beginning Inventory", clean_number(product['beginning_inventory']))
            c2.metric("Sample In", clean_number(product['sample_in']))
            c3.metric("Sample Out", clean_number(product['sample_out']))
            c4.metric("Current Inventory", f"{clean_number(product['current_inventory'])} {str(product['pack_type']).lower()}(s)")
            c5.metric("Total Content", product["total_content"])
            ins = query_df("SELECT date_received AS date, reference_no, 'Sample In' AS movement, quantity, remarks FROM stock_in WHERE product_id=?", (product_id,))
            outs = query_df("SELECT date_released AS date, reference_no, 'Sample Out' AS movement, quantity, remarks FROM stock_out WHERE product_id=?", (product_id,))
            movement = pd.concat([ins, outs], ignore_index=True).sort_values("date", ascending=False) if not ins.empty or not outs.empty else pd.DataFrame()
            with st.container(border=True):
                st.markdown("<div class='form-shell-title'>Movement History</div><div class='form-shell-subtitle'>Chronological record of Sample In and Sample Out transactions for this batch.</div>", unsafe_allow_html=True)
                display_df(movement, height=350)


def stock_in_page():
    page_header(
        "Sample In",
        "Create new product/batch records directly from Sample In, or add quantity to an existing batch."
    )
    module_hero("Sample In Workspace", "Receive new product samples, encode new batches, and keep receiving records complete with QC, cost, and document details.", "⬇️", "Receiving")

    today = date.today()
    month_start = today.replace(day=1)
    next_month_start = (pd.Timestamp(month_start) + pd.DateOffset(months=1)).date()
    received_today_qty = query_df("SELECT COALESCE(SUM(quantity),0) AS qty FROM stock_in WHERE date(date_received)=date(?)", (str(today),)).loc[0, "qty"]
    received_month_qty = query_df("SELECT COALESCE(SUM(quantity),0) AS qty FROM stock_in WHERE date(date_received)>=date(?) AND date(date_received)<date(?)", (str(month_start), str(next_month_start))).loc[0, "qty"]
    entries_today = query_df("SELECT COUNT(*) AS n FROM stock_in WHERE date(date_received)=date(?)", (str(today),)).loc[0, "n"]
    entries_month = query_df("SELECT COUNT(*) AS n FROM stock_in WHERE date(date_received)>=date(?) AND date(date_received)<date(?)", (str(month_start), str(next_month_start))).loc[0, "n"]

    c1, c2, c3 = st.columns(3)
    c1.metric("Received Today", int(entries_today))
    c2.metric("Sample Received This Month", int(entries_month))
    c3.metric("Packs Received This Month", clean_number(received_month_qty))

    supplier_map, suppliers = get_supplier_options()
    product_map, options = get_product_options()

    left, right = st.columns([1.02, 0.98])

    with left:
        st.markdown("<div class='section-title'>Sample-In Entry</div>", unsafe_allow_html=True)

        with st.form("stock_in_create_direct_form", clear_on_submit=False):
            category_options = get_dropdown_options("category")
            pack_options = get_dropdown_options("pack_type")
            unit_options = get_dropdown_options("pack_size_unit")
            storage_options = get_dropdown_options("storage_location")
            currency_options = get_dropdown_options("currency")
            cost_type_options = get_dropdown_options("cost_type")
            document_type_options = get_dropdown_options("document_type")
            receiving_condition_options = get_dropdown_options("receiving_condition")
            qc_status_options = get_dropdown_options("qc_status")
            default_currency = get_setting("default_currency", "PHP")

            st.markdown("##### A. Receiving Information")
            c0, c1, c2 = st.columns([1, 1, 1])
            reference_no = c0.text_input("Reference No.", value=generate_code("SIN", "stock_in", "stock_in_id"))
            date_received = c1.date_input("Date Received", value=date.today())
            received_by = c2.text_input("Received By", value="Admin User")

            entry_mode = st.radio(
                "Entry Type",
                ["Create new product/batch directly from Sample In", "Add sample to existing product/batch"],
                index=0,
                horizontal=False,
                help="Choose 'Create new' when the product, supplier, batch, pack size, and initial quantity should be encoded from this Sample In page."
            )

            st.markdown("##### B. Product / Supplier Information")

            # Default variables used during submit.
            selected = None
            product_id = None
            row = None
            supplier_id = None
            supplier_name_for_log = ""
            image = None
            storage_location = ""
            active_ingredient = ""
            category = ""
            product_code = ""
            product_name = ""
            pack_type = ""
            pack_size_value = 0.0
            pack_size_unit = "L"
            batch_number = ""
            expiry_date = date.today() + timedelta(days=365)

            if entry_mode == "Add sample to existing product/batch":
                if not options:
                    st.warning("No existing product/batch yet. Use the 'Create new product/batch directly from Sample-In' option.")
                else:
                    selected = st.selectbox("Product / Batch", options)
                    product_id = product_map[selected]
                    row = get_product(product_id)
                    supplier_id = int(row["supplier_id"]) if pd.notna(row["supplier_id"]) else None
                    supplier_name_for_log = row["supplier"]
                    product_name = row["product_name"]
                    active_ingredient = row["active_ingredient"]
                    category = row["category"]
                    pack_type = row["pack_type"]
                    pack_size_value = float(row["pack_size_value"] or 0)
                    pack_size_unit = row["pack_size_unit"]
                    batch_number = row["batch_number"] or ""
                    expiry_dt = pd.to_datetime(row["expiry_date"], errors="coerce")
                    expiry_date = expiry_dt.date() if pd.notna(expiry_dt) else date.today() + timedelta(days=365)

                    st.info(
                        f"Selected: {row['product_name']} | Supplier: {row['supplier']} | "
                        f"Batch: {row['batch_number']} | Pack: {row['pack_type']} | "
                        f"{row['pack_size']} per pack | Current Qty: {row['quantity']:g}"
                    )

                    c3, c4 = st.columns(2)
                    batch_number = c3.text_input("Batch Number", value=batch_number)
                    expiry_date = c4.date_input("Expiry Date", value=expiry_date)

            else:
                c3, c4 = st.columns(2)
                product_code = c3.text_input("Product Code", value=generate_code("PRD", "products", "product_id"))
                product_name = c4.text_input("Product Name *", placeholder="Example: Delfan Plus")

                c5, c6 = st.columns(2)
                active_ingredient = c5.text_input("Active Ingredient", placeholder="Example: L-Alpha Free Amino Acids")
                category = c6.selectbox("Category", category_options)

                c7, c8 = st.columns(2)
                supplier_choice = c7.selectbox("Supplier", (["Add New Supplier"] + suppliers) if suppliers else ["Add New Supplier"])
                new_supplier_name = c8.text_input("New Supplier Name", placeholder="Type supplier name if adding new")

                c9, c10 = st.columns(2)
                batch_number = c9.text_input("Batch Number *", placeholder="Example: DP-2026-001")
                expiry_date = c10.date_input("Expiry Date", value=date.today() + timedelta(days=365))

                storage_options = get_dropdown_options("storage_location")
                storage_location = st.selectbox("Storage Location", storage_options)
                image = st.file_uploader("Product Image / Label", type=["png", "jpg", "jpeg", "webp"])

            st.markdown("##### C. Packaging and Quantity")
            c11, c12, c13, c14 = st.columns([1, 1, 1, 1])
            if entry_mode == "Add sample to existing product/batch" and row is not None:
                pack_type = c11.text_input("Pack Type", value=str(pack_type or ""))
                pack_size_value = c12.number_input("Pack Size Value", min_value=0.0, value=float(pack_size_value or 0), step=0.25)
                pack_size_unit = c13.selectbox("Pack Size Unit", unit_options, index=(unit_options.index(pack_size_unit) if pack_size_unit in unit_options else 0))
            else:
                pack_type = c11.selectbox("Pack Type", pack_options)
                pack_size_value = c12.number_input("Pack Size Value", min_value=0.0, value=1.0, step=0.25)
                pack_size_unit = c13.selectbox("Pack Size Unit", unit_options)

            quantity = c14.number_input("Quantity / No. of Packs", min_value=0.0, value=1.0, step=1.0)

            total_content_value = float(quantity or 0) * float(pack_size_value or 0)
            total_content_display = f"{total_content_value:g} {pack_size_unit}"
            if str(pack_size_unit).lower() == "ml" and total_content_value >= 1000:
                total_content_display = f"{total_content_value/1000:g} L"
            elif str(pack_size_unit).lower() == "g" and total_content_value >= 1000:
                total_content_display = f"{total_content_value/1000:g} kg"

            st.markdown(
                f"""
                <div class='soft-card'>
                    <b>Auto Summary</b><br>
                    Pack: {pack_type} | {float(pack_size_value or 0):g} {pack_size_unit} per pack<br>
                    Quantity to receive: {float(quantity or 0):g} pack(s)<br>
                    Total content added: <b>{total_content_display}</b>
                </div>
                """,
                unsafe_allow_html=True,
            )

            st.markdown("##### D. Cost and Document Information")
            c15, c16, c17 = st.columns(3)
            unit_cost = c15.number_input("Unit Cost per Pack", min_value=0.0, value=0.0, step=10.0)
            currency = c16.selectbox("Currency", currency_options, index=(currency_options.index(default_currency) if default_currency in currency_options else 0))
            cost_type = c17.selectbox("Cost Type", cost_type_options)

            total_cost = float(quantity or 0) * float(unit_cost or 0)
            st.caption(f"Auto Total Cost: {currency} {total_cost:,.2f}")

            c18, c19 = st.columns(2)
            document_type = c18.selectbox("Document Type", document_type_options)
            document_no = c19.text_input("Document No.", placeholder="DR / Invoice / Sample Reference No.")
            attachment = st.file_uploader("Attachment / Delivery Receipt / COA", type=["pdf", "png", "jpg", "jpeg", "xlsx", "docx"])

            st.markdown("##### E. Quality Check")
            c20, c21 = st.columns(2)
            receiving_condition = c20.selectbox("Receiving Condition", receiving_condition_options)
            qc_status = c21.selectbox("QC Status", qc_status_options)
            remarks = st.text_area("Remarks", placeholder="Notes about receiving, storage, document, or QC observation")

            b1, b2 = st.columns(2)
            submitted = b1.form_submit_button("Save Sample-In", use_container_width=True)
            submitted_add_another = b2.form_submit_button("Save and Add Another", use_container_width=True)

        if submitted or submitted_add_another:
            if quantity <= 0:
                st.error("Quantity must be greater than zero.")
            elif not batch_number.strip():
                st.error("Batch Number is required.")
            else:
                attachment_path = save_upload(attachment, "stock_in_attachments")

                if entry_mode == "Add sample to existing product/batch":
                    if product_id is None or row is None:
                        st.error("Please select an existing product/batch or create a new record.")
                    else:
                        execute(
                            """
                            INSERT INTO stock_in
                            (reference_no, date_received, product_id, supplier_id, quantity, unit_cost,
                             batch_number, expiry_date, received_by, remarks, attachment_path,
                             document_type, document_no, currency, cost_type, receiving_condition, qc_status)
                            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                            """,
                            (reference_no, str(date_received), product_id, supplier_id, quantity, unit_cost,
                             batch_number, str(expiry_date), received_by, remarks, attachment_path,
                             document_type, document_no, currency, cost_type, receiving_condition, qc_status),
                        )
                        execute(
                            """
                            UPDATE products
                            SET quantity = quantity + ?, unit_cost=?, batch_number=?, expiry_date=?,
                                pack_type=?, pack_size_value=?, pack_size_unit=?, updated_at=CURRENT_TIMESTAMP
                            WHERE product_id=?
                            """,
                            (quantity, unit_cost, batch_number, str(expiry_date), pack_type, pack_size_value, pack_size_unit, product_id),
                        )
                        log_action("Sample In", "Create", f"{reference_no}: Received {quantity:g} pack(s) of {product_name}", received_by)
                        st.success("Sample-in saved and existing product quantity updated.")
                        rerun_app()

                else:
                    if not product_name.strip():
                        st.error("Product Name is required.")
                    else:
                        # Duplicate batch warning: the same product and batch should normally be updated through Existing Product mode.
                        dup = query_df(
                            "SELECT COUNT(*) AS n FROM products WHERE lower(product_name)=lower(?) AND lower(COALESCE(batch_number,''))=lower(?)",
                            (product_name.strip(), batch_number.strip()),
                        ).loc[0, "n"]
                        if int(dup) > 0:
                            st.error("This Product Name and Batch Number already exist. Use 'Add sample to existing product/batch' instead.")
                        elif 'supplier_choice' in locals() and supplier_choice == "Add New Supplier" and not new_supplier_name.strip():
                            st.error("Supplier name is required when adding a new supplier.")
                        else:
                            if 'supplier_choice' in locals() and supplier_choice == "Add New Supplier":
                                supplier_code = generate_code("SUP", "suppliers", "supplier_id")
                                supplier_id = execute(
                                    """
                                    INSERT INTO suppliers (supplier_code, supplier_name, status, remarks)
                                    VALUES (?, ?, 'Active', ?)
                                    """,
                                    (supplier_code, new_supplier_name.strip(), "Created from Sample-In entry"),
                                )
                                supplier_name_for_log = new_supplier_name.strip()
                            else:
                                supplier_id = supplier_map.get(supplier_choice)
                                supplier_name_for_log = supplier_choice

                            image_path = save_upload(image, "product_images")
                            product_id = execute(
                                """
                                INSERT INTO products
                                (product_code, product_name, active_ingredient, category, supplier_id, pack_type,
                                 pack_size_value, pack_size_unit, quantity, unit_cost, batch_number, expiry_date,
                                 storage_location, remarks, image_path)
                                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                                """,
                                (product_code, product_name.strip(), active_ingredient, category, supplier_id, pack_type,
                                 pack_size_value, pack_size_unit, 0, unit_cost, batch_number.strip(), str(expiry_date),
                                 storage_location, remarks, image_path),
                            )

                            execute(
                                """
                                INSERT INTO stock_in
                                (reference_no, date_received, product_id, supplier_id, quantity, unit_cost,
                                 batch_number, expiry_date, received_by, remarks, attachment_path,
                                 document_type, document_no, currency, cost_type, receiving_condition, qc_status)
                                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                                """,
                                (reference_no, str(date_received), product_id, supplier_id, quantity, unit_cost,
                                 batch_number.strip(), str(expiry_date), received_by, remarks, attachment_path,
                                 document_type, document_no, currency, cost_type, receiving_condition, qc_status),
                            )
                            execute("UPDATE products SET quantity = quantity + ?, updated_at=CURRENT_TIMESTAMP WHERE product_id=?", (quantity, product_id))
                            log_action("Sample In", "Create", f"{reference_no}: Created {product_name.strip()} / {batch_number.strip()} from {supplier_name_for_log} and received {quantity:g} pack(s)", received_by)
                            st.success("New product/batch record created directly from Sample In.")

                            if submitted:
                                rerun_app()
                            else:
                                st.info("Saved. You may now encode the next item.")

    with right:
        st.markdown("<div class='section-title'>Recent Sample-In Transactions</div>", unsafe_allow_html=True)
        recent = query_df(
            """
            SELECT si.reference_no, si.date_received, p.product_name, si.batch_number,
                   si.quantity, p.pack_type, s.supplier_name AS supplier, si.received_by,
                   si.document_type, si.document_no, si.receiving_condition, si.qc_status, si.remarks
            FROM stock_in si
            LEFT JOIN products p ON si.product_id = p.product_id
            LEFT JOIN suppliers s ON si.supplier_id = s.supplier_id
            ORDER BY si.created_at DESC
            LIMIT 15
            """
        )
        display_df(recent, height=380)

        with st.expander("✏️ Edit Sample-In Record", expanded=False):
            stock_in_edit_section()

        st.markdown("<div class='section-title'>Direct Sample-In Workflow</div>", unsafe_allow_html=True)
        st.markdown(
            """
            <div class='soft-card'>
            <b>Use this page to create everything in one transaction:</b><br><br>
            1. Select <b>Create new product/batch directly from Sample In</b><br>
            2. Encode Product Name, Supplier, Batch, Pack Size, and Quantity<br>
            3. Add document/QC details<br>
            4. Click <b>Save Sample-In</b><br><br>
            The app will automatically create the supplier if needed, create the product/batch, create the sample-in record, and update the quantity.
            </div>
            """,
            unsafe_allow_html=True,
        )


def stock_out_page():
    page_header("Sample Out", "Release samples with stock validation and FEFO batch recommendation.")
    module_hero("Sample Out Workspace", "Release samples with better visibility of stock status, FEFO guidance, and recent sample-out transactions.", "⬆️", "Releasing")
    today = date.today()
    month_start = today.replace(day=1)
    next_month_start = (pd.Timestamp(month_start) + pd.DateOffset(months=1)).date()
    released_today = query_df("SELECT COALESCE(SUM(quantity),0) AS qty FROM stock_out WHERE date(date_released)=date(?) AND status='Released'", (str(today),)).loc[0,"qty"]
    released_month = query_df("SELECT COALESCE(SUM(quantity),0) AS qty FROM stock_out WHERE date(date_released)>=date(?) AND date(date_released)<date(?) AND status='Released'", (str(month_start), str(next_month_start))).loc[0,"qty"]
    pending = query_df("SELECT COUNT(*) AS n FROM stock_out WHERE status='Pending'").loc[0,"n"]
    completed = query_df("SELECT COUNT(*) AS n FROM stock_out WHERE status='Released'").loc[0,"n"]
    c1, c2, c3, c4 = st.columns(4)
    c1.metric("Sample Out Today", clean_number(released_today))
    c2.metric("Sample Out This Month", clean_number(released_month))
    c3.metric("Pending Approval", int(pending))
    c4.metric("Completed Sample Out", int(completed))

    product_map, options = get_product_options()
    if not options:
        st.warning("Please add products first before recording sample-out.")
        return

    left, right = st.columns([0.95, 1.05])
    with left:
        section_shell("Release Sample", "Record a sample release request with product validation, approvals, and FEFO suggestions.")
        product_ids = list(product_map.values())
        product_labels = {pid: label for label, pid in product_map.items()}
        selected_product_id = st.selectbox(
            "Product / Batch",
            product_ids,
            format_func=lambda pid: product_labels.get(int(pid), str(pid)),
            key="stock_out_product_id_picker",
            help="This selector is outside the form so product details update immediately when you change the product/batch."
        )
        product_id = int(selected_product_id)
        row = get_product(product_id)
        st.info(f"Selected: {row['product_name']} | Batch: {row['batch_number']} | Pack Size: {row.get('pack_size', '') or '-'} | Available: {clean_number(row['quantity'])} {str(row['pack_type']).lower()}(s) | Status: {row['final_status']}")
        fefo = recommended_fefo(row["product_name"])
        if fefo is not None and int(fefo["product_id"]) != int(product_id):
            st.warning(f"FEFO Suggestion: Use batch {fefo['batch_number']} first. It expires on {fefo['expiry_date']} and has {fefo['quantity']:g} pack(s).")

        with st.form("stock_out_form", clear_on_submit=True):
            reference_no = st.text_input("Reference No.", value=generate_code("SOUT", "stock_out", "stock_out_id"))
            date_released = st.date_input("Date Released", value=date.today())
            c1, c2 = st.columns(2)
            quantity = c1.number_input("Quantity Released / No. of Packs", min_value=0.0, value=1.0, step=1.0)
            purpose = c2.selectbox("Purpose", get_dropdown_options("sample_out_purpose"))
            c3, c4, c5 = st.columns(3)
            requested_by = c3.text_input("Requested By")
            released_by = c4.text_input("Released By", value="Admin User")
            approved_by = c5.text_input("Approved By")
            status_options = get_dropdown_options("sample_out_status")
            status = st.selectbox("Status", status_options, index=(status_options.index("Released") if "Released" in status_options else 0))
            remarks = st.text_area("Remarks")
            submitted = st.form_submit_button("Release Sample", use_container_width=True)
        end_section_shell()

        if submitted:
            available = float(row["quantity"] or 0)
            if status == "Released" and quantity > available:
                st.error(f"Cannot release {clean_number(quantity)} pack(s). Available balance is only {clean_number(available)} pack(s).")
            elif status == "Released" and row["expiry_status"] == "Expired" and purpose != "Disposal":
                st.error("This batch is already expired. Change the purpose to Disposal if you need to remove expired stock from inventory.")
            else:
                execute(
                    """
                    INSERT INTO stock_out
                    (reference_no, date_released, product_id, quantity, purpose, requested_by, released_by, approved_by, status, remarks)
                    VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                    """,
                    (reference_no, str(date_released), product_id, quantity, purpose, requested_by, released_by, approved_by, status, remarks),
                )
                if status == "Released":
                    execute("UPDATE products SET quantity = quantity - ?, updated_at=CURRENT_TIMESTAMP WHERE product_id=?", (quantity, product_id))
                log_action("Sample Out", "Create", f"{reference_no}: {status} {clean_number(quantity)} pack(s) of {row['product_name']}", released_by)
                st.success("Sample-out transaction saved." + (" Sample quantity updated." if status == "Released" else " Quantity not deducted until status is Released."))
                rerun_app()

    with right:
        section_shell("Recent Sample-Out Transactions", "Review the latest sample releases and transaction details.")
        recent = query_df(
            """
            SELECT so.reference_no, so.date_released, p.product_name, p.batch_number,
                   so.quantity, p.pack_type, so.purpose, so.requested_by, so.released_by, so.status, so.remarks
            FROM stock_out so
            LEFT JOIN products p ON so.product_id = p.product_id
            ORDER BY so.created_at DESC
            LIMIT 50
            """
        )
        display_df(recent, height=520)
        with st.expander("✏️ Edit Sample-Out Record", expanded=False):
            stock_out_edit_section()
        end_section_shell()


def stock_balance_page():
    page_header("Stock Balance", "Real-time balance by product/batch, expiry status, and stock status.")
    module_hero("Stock Balance Overview", "Monitor inventory position by product and batch with clearer filters, balance view, and export actions.", "📋", "Balance")
    df = products_base_df()
    if df.empty:
        st.info("No products available.")
        return

    usable_stock_df = df[(df["quantity"] > 0) & (df["expiry_status"] != "Expired")]
    expired_records_df = df[df["expiry_status"] == "Expired"]
    c1, c2, c3, c4 = st.columns(4)
    c1.metric("Usable Stock", len(usable_stock_df))
    c2.metric("Out of Stock", len(df[df["final_status"] == "Out of Stock"]))
    c3.metric("Expired Records", len(expired_records_df))
    c4.metric("Not Expired Records", len(df[df["expiry_status"] == "Not Expired"]))

    section_shell("Filter and Review Stock Balance", "Use the filters below to narrow the inventory view and export the current result.")
    f1, f2, f3, f4 = st.columns(4)
    product_search = f1.text_input("Product Search", placeholder="Search product or batch")
    supplier = f2.selectbox("Supplier", ["All"] + sorted(df["supplier"].dropna().unique().tolist()))
    category = f3.selectbox("Category", ["All"] + sorted(df["category"].dropna().unique().tolist()))
    status = f4.selectbox("Final Status", ["All"] + sorted(df["final_status"].dropna().unique().tolist()))

    view = df.copy()
    if product_search:
        view = view[view.apply(lambda row: row_matches_search(row, product_search), axis=1)]
    if supplier != "All":
        view = view[view["supplier"] == supplier]
    if category != "All":
        view = view[view["category"] == category]
    if status != "All":
        view = view[view["final_status"] == status]

    stock_in_sum = query_df("SELECT product_id, COALESCE(SUM(quantity),0) AS total_stock_in FROM stock_in GROUP BY product_id")
    stock_out_sum = query_df("SELECT product_id, COALESCE(SUM(quantity),0) AS total_stock_out FROM stock_out WHERE status='Released' GROUP BY product_id")
    view = view.merge(stock_in_sum, on="product_id", how="left").merge(stock_out_sum, on="product_id", how="left")
    view["total_stock_in"] = view["total_stock_in"].fillna(0)
    view["total_stock_out"] = view["total_stock_out"].fillna(0)

    display_cols = [
        "product_name", "supplier", "category", "pack_type", "pack_size", "beginning_inventory",
        "total_stock_in", "total_stock_out", "current_inventory", "total_content",
        "batch_number", "expiry_date", "expiry_status", "stock_status", "final_status", "inventory_value"
    ]
    balance_display = view[display_cols].copy()
    balance_display = balance_display.rename(columns={
        "product_name": "Product Name",
        "supplier": "Supplier",
        "category": "Category",
        "pack_type": "Pack Type",
        "pack_size": "Pack Size",
        "beginning_inventory": "Beginning Inventory",
        "total_stock_in": "Sample In",
        "total_stock_out": "Sample Out",
        "current_inventory": "Current Inventory",
        "total_content": "Total Content",
        "batch_number": "Batch Number",
        "expiry_date": "Expiry Date",
        "expiry_status": "Expiry Status",
        "stock_status": "Stock Status",
        "final_status": "Final Status",
        "inventory_value": "Inventory Value",
    })
    visible_balance_display = display_df(balance_display, height=540, table_key="stock_balance")

    cdl1, cdl2 = st.columns(2)
    cdl1.download_button("Download Stock Balance CSV", to_csv_bytes(visible_balance_display), "stock_balance.csv", "text/csv", use_container_width=True)
    cdl2.download_button("Download Stock Balance Excel", to_excel_bytes({"Stock Balance": visible_balance_display}), "stock_balance.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", use_container_width=True)
    end_section_shell()


def suppliers_page():
    page_header("Suppliers", "Manage supplier details and review products supplied.")
    module_hero("Supplier Management", "Maintain supplier profiles, search supplier records, and review products and sample-in history from each partner.", "🏢", "Partners")
    tabs = st.tabs(["Supplier List", "Add Supplier", "Supplier Details"])

    with tabs[0]:
        section_shell("Supplier List", "Search, review, and export supplier records.")
        df = query_df("SELECT * FROM suppliers ORDER BY supplier_name")
        search = st.text_input("Search supplier", placeholder="Supplier name, contact person, country...")
        view = df.copy()
        if search and not view.empty:
            view = view[view.apply(lambda row: row_matches_search(row, search), axis=1)]
        display_df(view, height=480)
        st.download_button("Download Suppliers CSV", to_csv_bytes(view), "suppliers.csv", "text/csv")
        end_section_shell()

    with tabs[1]:
        section_shell("Add Supplier", "Capture supplier profile and contact details for future sample transactions.")
        with st.form("supplier_form", clear_on_submit=True):
            c1, c2 = st.columns(2)
            supplier_code = c1.text_input("Supplier Code", value=generate_code("SUP", "suppliers", "supplier_id"))
            supplier_name = c2.text_input("Supplier Name *")
            c3, c4 = st.columns(2)
            contact_person = c3.text_input("Contact Person")
            contact_number = c4.text_input("Contact Number")
            c5, c6 = st.columns(2)
            email = c5.text_input("Email Address")
            country = c6.text_input("Country / Origin")
            status = st.selectbox("Status", ["Active", "Inactive"])
            remarks = st.text_area("Remarks")
            submitted = st.form_submit_button("Save Supplier", use_container_width=True)
        end_section_shell()
        if submitted:
            if not supplier_name.strip():
                st.error("Supplier Name is required.")
            else:
                try:
                    execute(
                        """
                        INSERT INTO suppliers
                        (supplier_code, supplier_name, contact_person, contact_number, email, country, status, remarks)
                        VALUES (?, ?, ?, ?, ?, ?, ?, ?)
                        """,
                        (supplier_code, supplier_name, contact_person, contact_number, email, country, status, remarks),
                    )
                    log_action("Suppliers", "Create", f"Added supplier {supplier_name}")
                    st.success("Supplier saved successfully.")
                except sqlite3.IntegrityError:
                    st.error("Supplier code already exists. Please use another code.")

    with tabs[2]:
        section_shell("Supplier Details", "View supplied products and sample-in activity for a selected supplier.")
        supplier_map, suppliers = get_supplier_options()
        if not suppliers:
            st.info("No suppliers available.")
        else:
            selected = st.selectbox("Select Supplier", suppliers)
            supplier_id = supplier_map[selected]
            products = products_base_df()
            products = products[products["supplier_id"] == supplier_id]
            stockin = query_df(
                """
                SELECT si.reference_no, si.date_received, p.product_name, p.batch_number, si.quantity, p.pack_type, si.received_by
                FROM stock_in si LEFT JOIN products p ON si.product_id = p.product_id
                WHERE si.supplier_id=?
                ORDER BY si.date_received DESC
                """, (supplier_id,)
            )
            st.markdown(f"### Products Supplied by {selected}")
            display_df(products, ["product_code", "product_name", "category", "pack_type", "pack_size", "quantity", "batch_number", "expiry_date", "final_status"], height=260)
            st.markdown(f"### Sample-In History from {selected}")
            display_df(stockin, height=260)
        end_section_shell()


def render_report_preview_chart(report_type, df):
    if df is None or df.empty:
        return
    fig = None
    if report_type == "Monthly Movement Report" and {"month", "stock_in", "stock_out"}.issubset(df.columns):
        chart_df = df.copy().sort_values("month")
        chart_df = chart_df.rename(columns={"month": "Month", "stock_in": "Sample In", "stock_out": "Sample Out"})
        long_df = chart_df.melt(id_vars="Month", value_vars=["Sample In", "Sample Out"], var_name="Movement", value_name="Quantity")
        fig = px.bar(long_df, x="Month", y="Quantity", color="Movement", barmode="group", color_discrete_sequence=["#16a34a", "#2563eb"])
        fig = apply_chart_layout(fig)
    elif report_type == "Supplier Report" and {"supplier", "products"}.issubset(df.columns):
        chart_df = df.copy().sort_values("products", ascending=False).head(10)
        fig = px.bar(chart_df, x="supplier", y="products", color_discrete_sequence=["#16a34a"])
        fig = apply_chart_layout(fig)
    elif report_type == "Inventory Value Report" and {"product_name", "inventory_value"}.issubset(df.columns):
        chart_df = df.copy().sort_values("inventory_value", ascending=False).head(10)
        fig = px.bar(chart_df, x="product_name", y="inventory_value", color_discrete_sequence=["#16a34a"])
        fig = apply_chart_layout(fig)
    elif report_type == "Inventory Balance Report" and "category" in df.columns:
        # Count actual product/batch records per category.
        # This intentionally counts rows/batches, not the sum of current inventory quantity.
        chart_df = (
            df.copy()
            .assign(category=lambda x: x["category"].fillna("Unspecified").astype(str).replace({"": "Unspecified"}))
            .groupby("category", dropna=False)
            .size()
            .reset_index(name="product_batch_count")
            .sort_values("product_batch_count", ascending=False)
        )
        fig = px.bar(
            chart_df,
            x="category",
            y="product_batch_count",
            text="product_batch_count",
            labels={"category": "Category", "product_batch_count": "Product/Batch Count"},
            color_discrete_sequence=["#16a34a"],
        )
        fig.update_traces(texttemplate="%{text:.0f}", textposition="outside", cliponaxis=False)
        fig = apply_chart_layout(fig)
        fig.update_yaxes(title_text="Product/Batch Count", tickformat=",.0f")
        fig.update_xaxes(title_text="Category")
    if fig is not None:
        st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})


def reports_page():
    page_header("Reports", "Generate exportable inventory, sample-in, sample-out, supplier, and audit reports.")
    module_hero("Reporting Center", "Choose a report type, apply filters, and export clean report outputs for operations and management use.", "📑", "Export")

    section_shell("Build Report", "Select a report type, apply filters, and preview the report before exporting.")
    report_type = st.selectbox(
        "Select Report",
        [
            "Inventory Balance Report",
            "Sample-In Report",
            "Sample-Out Report",
            "Expiring Products Report",
            "Expired Products Report",
            "Supplier Report",
            "Monthly Movement Report",
            "Inventory Value Report",
            "Audit Trail Report",
        ],
    )

    df_products = products_base_df()
    c1, c2, c3, c4 = st.columns(4)
    start_date = c1.date_input("Start Date", value=date.today().replace(day=1))
    end_date = c2.date_input("End Date", value=date.today())
    supplier_filter = c3.selectbox("Supplier", ["All"] + sorted(df_products["supplier"].dropna().unique().tolist()) if not df_products.empty else ["All"])
    category_filter = c4.selectbox("Category", ["All"] + sorted(df_products["category"].dropna().unique().tolist()) if not df_products.empty else ["All"])

    df = build_report(report_type, start_date, end_date, supplier_filter, category_filter)

    s1, s2, s3 = st.columns(3)
    with s1:
        report_summary_card("Rows in Preview", len(df), "Current filtered result")
    with s2:
        report_summary_card("Report Period", f"{start_date.strftime('%b %d')} - {end_date.strftime('%b %d')}", "Selected date range")
    active_filters = []
    if supplier_filter != "All":
        active_filters.append(supplier_filter)
    if category_filter != "All":
        active_filters.append(category_filter)
    with s3:
        report_summary_card("Active Filters", len(active_filters), ", ".join(active_filters) if active_filters else "No extra filters")
    end_section_shell()

    section_shell("Report Preview", "Preview chart insights and table output before exporting.")
    preview_left, preview_right = st.columns([1.3, 1])
    with preview_left:
        st.markdown(f"<div class='section-title'>{report_type}</div>", unsafe_allow_html=True)
        chart_note("Preview chart adapts automatically for supported report types such as monthly movement, supplier, inventory balance, and inventory value.")
    with preview_right:
        preview_density = st.radio("Table View Mode", ["Comfortable", "Compact"], horizontal=True, index=(0 if get_setting("table_density", "Comfortable") == "Comfortable" else 1), key="report_table_density")
        st.markdown("<div class='table-toolbar-note'>Sticky headers stay visible while you scroll inside the table.</div>", unsafe_allow_html=True)

    render_report_preview_chart(report_type, df)
    visible_report_df = display_df(df, height=560, density=preview_density, table_key="report_preview")

    x1, x2, x3 = st.columns(3)
    x1.download_button("Export CSV", to_csv_bytes(visible_report_df), f"{report_type.lower().replace(' ', '_')}.csv", "text/csv", use_container_width=True)
    x2.download_button("Export Excel", to_excel_bytes({report_type: visible_report_df}), f"{report_type.lower().replace(' ', '_')}.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", use_container_width=True)
    pdf_bytes = to_pdf_bytes(visible_report_df, report_type)
    if pdf_bytes:
        x3.download_button("Export PDF", pdf_bytes, f"{report_type.lower().replace(' ', '_')}.pdf", "application/pdf", use_container_width=True)
    else:
        x3.info("Install reportlab to enable PDF export.")
    end_section_shell()

    section_shell("Premium Monthly Report Generator", "Generate a management-style monthly report in PDF or PowerPoint format.")
    report_month = st.date_input("Report Month", value=date.today().replace(day=1), key="premium_report_month")
    premium_data = build_monthly_report_data(report_month)
    p1, p2, p3, p4, p5 = st.columns(5)
    p1.metric("Usable Stock", clean_number(premium_data["usable_count"]))
    p2.metric("Expiring Soon", clean_number(premium_data["expiring_count"]))
    p3.metric("Expired", clean_number(premium_data["expired_count"]))
    p4.metric("Received", clean_number(premium_data["received_count"]))
    p5.metric("Released", clean_number(premium_data["released_count"]))
    d1, d2 = st.columns(2)
    pdf_report = premium_monthly_pdf_bytes(premium_data)
    st.caption("PPT design follows the reference layout with polished typography, cleaner tables, improved KPI cards, chart number formatting, and stronger management-report spacing.")
    ppt_report = premium_monthly_pptx_bytes(premium_data)
    month_file = premium_data["report_month"].strftime("%Y_%m")
    if pdf_report:
        d1.download_button("Download Premium PDF Report", pdf_report, f"monthly_inventory_report_{month_file}.pdf", "application/pdf", use_container_width=True)
    else:
        d1.info("Install reportlab to enable PDF export.")
    if ppt_report:
        d2.download_button("Download PPT Presentation Report", ppt_report, f"monthly_inventory_report_{month_file}.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation", use_container_width=True)
    else:
        d2.info("Install python-pptx to enable PPT export.")
    end_section_shell()

    st.markdown("---")
    with st.expander("Backup Database"):
        if DB_PATH.exists():
            st.download_button("Download SQLite Database Backup", DB_PATH.read_bytes(), "inventory_app_backup.db", "application/octet-stream")
        st.caption("Keep this backup file in a safe folder. You can replace your local inventory_app.db with this file if you need to restore data manually.")


def build_report(report_type, start_date, end_date, supplier_filter, category_filter):
    products = products_base_df()
    if not products.empty:
        if supplier_filter != "All":
            products = products[products["supplier"] == supplier_filter]
        if category_filter != "All":
            products = products[products["category"] == category_filter]

    if report_type == "Inventory Balance Report":
        return products[["product_code", "product_name", "supplier", "category", "pack_type", "pack_size", "beginning_inventory", "sample_in", "sample_out", "current_inventory", "total_content", "batch_number", "expiry_date", "final_status"]] if not products.empty else products

    if report_type == "Expiring Products Report":
        expiry_alert_days = get_int_setting("expiry_alert_days", 90)
        df = products[(products["expiry_status"] != "Expired") & (products["days_to_expiry"].notna()) & (products["days_to_expiry"] <= expiry_alert_days)] if not products.empty else products
        return df[["product_code", "product_name", "supplier", "category", "quantity", "batch_number", "expiry_date", "days_to_expiry", "final_status"]] if not df.empty else df

    if report_type == "Expired Products Report":
        df = products[products["expiry_status"] == "Expired"] if not products.empty else products
        return df[["product_code", "product_name", "supplier", "category", "quantity", "batch_number", "expiry_date", "final_status"]] if not df.empty else df

    if report_type == "Supplier Report":
        if products.empty:
            return products
        return products.groupby("supplier").agg(
            products=("product_id", "count"),
            total_quantity=("quantity", "sum"),
            inventory_value=("inventory_value", "sum"),
            expired_products=("expiry_status", lambda s: (s == "Expired").sum()),
            out_of_stock=("stock_status", lambda s: (s == "Out of Stock").sum()),
        ).reset_index()

    if report_type == "Inventory Value Report":
        if products.empty:
            return products
        return products[["product_code", "product_name", "supplier", "category", "quantity", "unit_cost", "inventory_value", "batch_number", "final_status"]].sort_values("inventory_value", ascending=False)

    if report_type == "Sample-In Report":
        df = query_df(
            """
            SELECT si.reference_no, si.date_received, p.product_name, p.category, s.supplier_name AS supplier,
                   si.quantity, p.pack_type, si.unit_cost, si.batch_number, si.expiry_date, si.received_by, si.remarks
            FROM stock_in si
            LEFT JOIN products p ON si.product_id = p.product_id
            LEFT JOIN suppliers s ON si.supplier_id = s.supplier_id
            WHERE date(si.date_received) BETWEEN date(?) AND date(?)
            ORDER BY si.date_received DESC
            """, (str(start_date), str(end_date))
        )
    elif report_type == "Sample-Out Report":
        df = query_df(
            """
            SELECT so.reference_no, so.date_released, p.product_name, p.category, s.supplier_name AS supplier,
                   so.quantity, p.pack_type, so.purpose, so.requested_by, so.released_by, so.approved_by, so.status, so.remarks
            FROM stock_out so
            LEFT JOIN products p ON so.product_id = p.product_id
            LEFT JOIN suppliers s ON p.supplier_id = s.supplier_id
            WHERE date(so.date_released) BETWEEN date(?) AND date(?)
            ORDER BY so.date_released DESC
            """, (str(start_date), str(end_date))
        )
    elif report_type == "Monthly Movement Report":
        stock_in = query_df("SELECT substr(date_received, 1, 7) AS month, SUM(quantity) AS stock_in FROM stock_in GROUP BY substr(date_received, 1, 7)")
        stock_out = query_df("SELECT substr(date_released, 1, 7) AS month, SUM(quantity) AS stock_out FROM stock_out WHERE status='Released' GROUP BY substr(date_released, 1, 7)")
        df = pd.merge(stock_in, stock_out, on="month", how="outer").fillna(0).sort_values("month", ascending=False)
        return df
    elif report_type == "Audit Trail Report":
        return query_df("SELECT action_time, module, action, details, user_name FROM audit_trail ORDER BY action_time DESC LIMIT 500")
    else:
        df = pd.DataFrame()

    if not df.empty:
        if supplier_filter != "All" and "supplier" in df.columns:
            df = df[df["supplier"] == supplier_filter]
        if category_filter != "All" and "category" in df.columns:
            df = df[df["category"] == category_filter]
    return df

def settings_page():
    page_header("Settings", "Manage app preferences, dropdown lists, users, backup, and restore tools.")
    module_hero("Application Settings", "Control app preferences, chart behavior, dropdown lists, users, and backup tools from one place.", "⚙️", "Configuration")

    tab1, tab2, tab3, tab4, tab5 = st.tabs([
        "App Preferences",
        "Chart Settings",
        "Dropdown Lists",
        "Users & Roles",
        "Backup / Restore",
    ])

    with tab1:
        section_shell("App Preferences", "Update the app name, company name, default currency, alert days, and default user.")
        with st.form("settings_preferences_form"):
            c1, c2 = st.columns(2)
            app_name = c1.text_input("App Name", value=get_setting("app_name", "Sample Inventory App"))
            company_name = c2.text_input("Company Name", value=get_setting("company_name", "Farmfix"))
            c3, c4, c5 = st.columns(3)
            currency_options = get_dropdown_options("currency")
            default_currency_current = get_setting("default_currency", "PHP")
            default_currency = c3.selectbox("Default Currency", currency_options, index=(currency_options.index(default_currency_current) if default_currency_current in currency_options else 0))
            expiry_alert_days = c4.number_input("Expiring Soon Alert Days", min_value=1, max_value=365, value=get_int_setting("expiry_alert_days", 90), step=1)
            urgent_expiry_days = c5.number_input("Urgent Expiry Alert Days", min_value=1, max_value=365, value=get_int_setting("urgent_expiry_days", 30), step=1)
            c6, c7 = st.columns(2)
            current_user = c6.text_input("Default User Name", value=get_setting("current_user", "Admin User"))
            table_density = c7.selectbox("Default Table View", ["Comfortable", "Compact"], index=(0 if get_setting("table_density", "Comfortable") == "Comfortable" else 1))
            save_pref = st.form_submit_button("Save Preferences", use_container_width=True)
        end_section_shell()
        if save_pref:
            set_setting("app_name", app_name)
            set_setting("company_name", company_name)
            set_setting("default_currency", default_currency)
            set_setting("expiry_alert_days", int(expiry_alert_days))
            set_setting("urgent_expiry_days", int(urgent_expiry_days))
            set_setting("current_user", current_user)
            set_setting("table_density", table_density)
            log_action("Settings", "Update", "Updated app preferences", current_user)
            st.success("Settings saved. The sidebar/app labels will refresh now.")
            rerun_app()
        section_shell("Current Alert Logic", "Reference the active alert thresholds used by the app.")
        st.info(
            f"Products are marked as Urgent Expiry when expiry is within {get_int_setting('urgent_expiry_days', 30)} days, "
            f"and Expiring Soon when expiry is within {get_int_setting('expiry_alert_days', 90)} days."
        )
        end_section_shell()
        section_shell("Mobile Layout Optimization", "This version includes improved spacing, card sizing, tabs, and responsive layouts for smaller screens.")
        st.caption("Tip: On phones, collapse the sidebar when you finish navigating to maximize workspace for forms, charts, and tables.")
        end_section_shell()

    with tab2:
        section_shell("Dashboard Chart Settings", "Control chart types and color palettes used on the dashboard.")
        chart_preferences_form("settings")
        st.info("These settings control the chart type and color palette shown in the Dashboard. You can choose preset palettes or select Custom to choose multiple colors.")
        end_section_shell()

    with tab3:
        section_shell("Dropdown List Management", "Maintain reusable dropdown values across the app.")
        list_labels = {
            "category": "Product Categories",
            "pack_type": "Pack Types",
            "pack_size_unit": "Pack Size Units",
            "storage_location": "Storage Locations",
            "currency": "Currencies",
            "cost_type": "Cost Types",
            "document_type": "Document Types",
            "receiving_condition": "Receiving Conditions",
            "qc_status": "QC Status",
            "sample_out_purpose": "Sample Out Purposes",
            "sample_out_status": "Sample Out Statuses",
            "user_role": "User Roles",
        }
        selected_list = st.selectbox("Select dropdown list to manage", list(list_labels.keys()), format_func=lambda x: list_labels[x])
        c1, c2 = st.columns([2, 1])
        new_value = c1.text_input("Add new option", placeholder="Type new dropdown option")
        if c2.button("Add Option", use_container_width=True):
            if new_value.strip():
                add_dropdown_option(selected_list, new_value)
                log_action("Settings", "Add Dropdown Option", f"Added {new_value} to {selected_list}", get_setting("current_user", "Admin User"))
                st.success("Option added.")
                rerun_app()
            else:
                st.warning("Please type an option first.")
        options_df = query_df(
            """
            SELECT option_value, is_active, sort_order
            FROM dropdown_options
            WHERE list_name=?
            ORDER BY is_active DESC, sort_order, option_value
            """,
            (selected_list,),
        )
        if options_df.empty:
            st.info("No options found for this list.")
        else:
            st.caption("Deactivate options you no longer want to show. Existing records will not be deleted.")
            for _, row in options_df.iterrows():
                col_a, col_b, col_c = st.columns([3, 1, 1])
                col_a.write(row["option_value"])
                col_b.write("Active" if int(row["is_active"]) == 1 else "Inactive")
                if int(row["is_active"]) == 1:
                    if col_c.button("Deactivate", key=f"deactivate_{selected_list}_{row['option_value']}"):
                        set_dropdown_active(selected_list, row["option_value"], False)
                        rerun_app()
                else:
                    if col_c.button("Activate", key=f"activate_{selected_list}_{row['option_value']}"):
                        set_dropdown_active(selected_list, row["option_value"], True)
                        rerun_app()
        end_section_shell()

    with tab4:
        section_shell("Users and Roles", "Manage locally stored user names and role assignments.")
        role_options = get_dropdown_options("user_role")
        with st.form("add_user_form", clear_on_submit=True):
            c1, c2, c3 = st.columns(3)
            user_name = c1.text_input("User Name")
            role = c2.selectbox("Role", role_options)
            status = c3.selectbox("Status", ["Active", "Inactive"])
            add_user = st.form_submit_button("Add User", use_container_width=True)
        if add_user:
            if not user_name.strip():
                st.error("User Name is required.")
            else:
                try:
                    execute(
                        """
                        INSERT INTO app_users (user_name, role, status)
                        VALUES (?, ?, ?)
                        """,
                        (user_name.strip(), role, status),
                    )
                    log_action("Settings", "Add User", f"Added user {user_name.strip()} as {role}", get_setting("current_user", "Admin User"))
                    st.success("User added.")
                    rerun_app()
                except sqlite3.IntegrityError:
                    st.error("User already exists.")
        users = query_df("SELECT user_name, role, status, created_at FROM app_users ORDER BY user_name")
        display_df(users, height=320)
        st.caption("This first version stores user roles for tracking. Full login/password protection can be added later.")
        end_section_shell()

    with tab5:
        section_shell("Backup and Restore", "Create backups, export all data, restore a previous database, or clear demo records.")
        c1, c2 = st.columns(2)
        if DB_PATH.exists():
            c1.download_button(
                "Download SQLite Database Backup",
                DB_PATH.read_bytes(),
                "inventory_app_backup.db",
                "application/octet-stream",
                use_container_width=True,
            )
        c2.download_button(
            "Export All Data to Excel",
            export_all_data_bytes(),
            "inventory_app_all_data.xlsx",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            use_container_width=True,
        )

        st.markdown("#### Restore Database")
        uploaded_db = st.file_uploader("Upload previous inventory_app_backup.db", type=["db", "sqlite", "sqlite3"])
        if uploaded_db is not None:
            if st.button("Restore Uploaded Database", use_container_width=True):
                backup_name = f"inventory_app_before_restore_{datetime.now().strftime('%Y%m%d%H%M%S')}.db"
                if DB_PATH.exists():
                    Path(backup_name).write_bytes(DB_PATH.read_bytes())
                DB_PATH.write_bytes(uploaded_db.getbuffer())
                st.success("Database restored. Please stop the app with Ctrl+C and run it again to fully reload the restored database.")

        st.markdown("#### Clear Demo Data")
        st.warning("This removes all products, suppliers, sample-in, sample-out, and audit records. Settings/dropdown options are kept.")
        confirm_clear = st.checkbox("I understand that this will delete inventory records.")
        if st.button("Clear Inventory Records", disabled=not confirm_clear, use_container_width=True):
            for table in ["stock_out", "stock_in", "products", "suppliers", "audit_trail"]:
                execute(f"DELETE FROM {table}")
            log_action("Settings", "Clear Data", "Cleared inventory records", get_setting("current_user", "Admin User"))
            st.success("Inventory records cleared.")
            rerun_app()
        end_section_shell()


# -----------------------------

# -----------------------------
# LOGIN, VERSION, AND AUTO BACKUP
# -----------------------------
USER_ACCOUNTS = {
    "admin": {"password": "admin123", "role": "Admin", "display_name": "Admin User"},
    "viewer": {"password": "viewer123", "role": "Viewer", "display_name": "Viewer User"},
}





# -----------------------------
# SMARTER OFFLINE AI ASSISTANT - VERSION 1.2
# Rule-based decision assistant. No API key, no internet, no external data sharing.
# -----------------------------

def _ai_text(value):
    if value is None:
        return ""
    try:
        if pd.isna(value):
            return ""
    except Exception:
        pass
    return str(value).strip()


def _ai_normalize(value):
    return re.sub(r"\s+", " ", _ai_text(value).lower()).strip()


def _ai_money_or_number(value):
    try:
        return clean_number(value)
    except Exception:
        return str(value)


def _ai_extract_keywords(question):
    q = _ai_normalize(question)
    stop_words = {
        "what", "which", "show", "give", "list", "products", "product", "batch", "batches",
        "sample", "samples", "stock", "available", "expired", "expiring", "soon", "details",
        "of", "the", "are", "is", "in", "from", "with", "qty", "quantity", "current", "how",
        "many", "supplier", "suppliers", "records", "record", "status", "inventory", "please",
        "can", "you", "me", "for", "and", "or", "still", "have", "has", "this", "that",
        "summarize", "summary", "report", "release", "released", "received", "month"
    }
    return [
        t for t in re.findall(r"[a-zA-Z0-9\-]+", q)
        if len(t) >= 3 and t not in stop_words
    ]


def _ai_detail_table(df, max_rows=120):
    """Prepare readable product/batch records for the AI Assistant results table."""
    if df is None or df.empty:
        return pd.DataFrame()

    temp = df.copy()
    preferred = [
        "product_name", "batch_number", "supplier", "category", "active_ingredient",
        "current_inventory", "pack_size", "pack_type", "expiry_date", "days_to_expiry",
        "expiry_status", "final_status", "storage_location", "remarks"
    ]
    cols = [c for c in preferred if c in temp.columns]
    temp = temp[cols].head(max_rows).copy()

    rename_map = {
        "product_name": "Product",
        "batch_number": "Batch No.",
        "supplier": "Supplier",
        "category": "Category",
        "active_ingredient": "Active Ingredient",
        "current_inventory": "Current Qty",
        "pack_size": "Pack Size",
        "pack_type": "Pack Type",
        "expiry_date": "Expiry Date",
        "days_to_expiry": "Days to Expiry",
        "expiry_status": "Expiry Status",
        "final_status": "Status",
        "storage_location": "Storage Location",
        "remarks": "Remarks",
    }
    return temp.rename(columns=rename_map)


def _ai_stock_in_month_df():
    today = date.today()
    month_start = date(today.year, today.month, 1)
    return query_df(
        """
        SELECT
            si.date_received AS "Date Received",
            si.reference_no AS "Reference No.",
            COALESCE(p.product_name, '') AS "Product",
            COALESCE(p.batch_number, si.batch_number, '') AS "Batch No.",
            COALESCE(s.supplier_name, 'Unknown Supplier') AS "Supplier",
            si.quantity AS "Quantity",
            CASE
                WHEN p.pack_size_value IS NOT NULL AND p.pack_size_unit IS NOT NULL
                THEN printf('%g %s', p.pack_size_value, p.pack_size_unit)
                ELSE ''
            END AS "Pack Size",
            COALESCE(p.pack_type, '') AS "Pack Type",
            COALESCE(si.received_by, '') AS "Received By",
            COALESCE(si.remarks, '') AS "Remarks"
        FROM stock_in si
        LEFT JOIN products p ON si.product_id = p.product_id
        LEFT JOIN suppliers s ON si.supplier_id = s.supplier_id
        WHERE date(si.date_received) >= date(?)
        ORDER BY date(si.date_received) DESC, si.created_at DESC
        LIMIT 150
        """,
        (str(month_start),),
    )


def _ai_stock_out_month_df():
    today = date.today()
    month_start = date(today.year, today.month, 1)
    return query_df(
        """
        SELECT
            so.date_released AS "Date Released",
            so.reference_no AS "Reference No.",
            COALESCE(p.product_name, '') AS "Product",
            COALESCE(p.batch_number, '') AS "Batch No.",
            COALESCE(s.supplier_name, 'Unknown Supplier') AS "Supplier",
            so.quantity AS "Quantity",
            COALESCE(p.pack_type, '') AS "Pack Type",
            COALESCE(so.purpose, '') AS "Purpose",
            COALESCE(so.requested_by, '') AS "Requested By",
            COALESCE(so.released_by, '') AS "Released By",
            COALESCE(so.status, '') AS "Status",
            COALESCE(so.remarks, '') AS "Remarks"
        FROM stock_out so
        LEFT JOIN products p ON so.product_id = p.product_id
        LEFT JOIN suppliers s ON p.supplier_id = s.supplier_id
        WHERE date(so.date_released) >= date(?)
        ORDER BY date(so.date_released) DESC, so.created_at DESC
        LIMIT 150
        """,
        (str(month_start),),
    )


def _ai_status_summary(df):
    if df is None or df.empty:
        return {
            "total_records": 0, "usable": 0, "expiring": 0, "urgent": 0,
            "expired": 0, "expired_with_stock": 0, "out": 0,
            "suppliers": 0, "categories": 0, "total_qty": 0,
        }

    current = pd.to_numeric(df.get("current_inventory", 0), errors="coerce").fillna(0)
    return {
        "total_records": len(df),
        "usable": len(df[(current > 0) & (df["expiry_status"] != "Expired")]),
        "expiring": len(df[df["final_status"] == "Expiring Soon"]),
        "urgent": len(df[df["final_status"] == "Urgent Expiry"]),
        "expired": len(df[df["expiry_status"] == "Expired"]),
        "expired_with_stock": len(df[(df["expiry_status"] == "Expired") & (current > 0)]),
        "out": len(df[current <= 0]),
        "suppliers": df["supplier"].fillna("").nunique() if "supplier" in df.columns else 0,
        "categories": df["category"].fillna("").nunique() if "category" in df.columns else 0,
        "total_qty": current.sum(),
    }


def _ai_find_entity(question, values, min_len=3):
    """Find the best matching entity from known database values."""
    q = _ai_normalize(question)
    if not values:
        return None

    candidates = []
    for value in values:
        text = _ai_text(value)
        if len(text) < min_len:
            continue
        norm = _ai_normalize(text)
        if not norm:
            continue
        if norm in q:
            candidates.append((len(norm), text))
        else:
            # Partial word detection for product-like names.
            tokens = [t for t in re.findall(r"[a-zA-Z0-9]+", norm) if len(t) >= min_len]
            if tokens and any(t in q for t in tokens):
                candidates.append((max(len(t) for t in tokens), text))

    if not candidates:
        return None
    candidates.sort(reverse=True)
    return candidates[0][1]


def _ai_detect_entities(question, df):
    """Detect product, supplier, category, batch, and status from the question."""
    if df is None or df.empty:
        return {}

    entities = {}
    q = _ai_normalize(question)

    product_values = sorted(df["product_name"].dropna().astype(str).unique().tolist(), key=len, reverse=True) if "product_name" in df.columns else []
    supplier_values = sorted(df["supplier"].dropna().astype(str).unique().tolist(), key=len, reverse=True) if "supplier" in df.columns else []
    category_values = sorted(df["category"].dropna().astype(str).unique().tolist(), key=len, reverse=True) if "category" in df.columns else []
    batch_values = sorted(df["batch_number"].dropna().astype(str).unique().tolist(), key=len, reverse=True) if "batch_number" in df.columns else []

    entities["product"] = _ai_find_entity(question, product_values, min_len=3)
    entities["supplier"] = _ai_find_entity(question, supplier_values, min_len=3)
    entities["category"] = _ai_find_entity(question, category_values, min_len=3)
    entities["batch"] = _ai_find_entity(question, batch_values, min_len=3)

    if "urgent" in q:
        entities["status"] = "Urgent Expiry"
    elif "expiring" in q or "near expiry" in q or "near expiration" in q:
        entities["status"] = "Expiring Soon"
    elif "expired" in q or "expiration" in q:
        entities["status"] = "Expired"
    elif "out of stock" in q or "zero stock" in q or "no stock" in q:
        entities["status"] = "Out of Stock"
    elif "available" in q or "usable" in q or "with stock" in q:
        entities["status"] = "Available"

    return {k: v for k, v in entities.items() if v}


def _ai_detect_intent(question):
    q = _ai_normalize(question)

    if any(k in q for k in ["management summary", "management report", "summary", "overview", "inventory status", "status report"]):
        return "management_summary"
    if any(k in q for k in ["fefo", "release first", "prioritize", "priority", "use first"]):
        return "fefo"
    if any(k in q for k in ["received this month", "sample in this month", "stock in this month", "received samples"]):
        return "received_this_month"
    if any(k in q for k in ["released this month", "sample out this month", "stock out this month", "released samples"]):
        return "released_this_month"
    if any(k in q for k in ["expired with stock", "expired still have", "expired products still", "expired batches still"]):
        return "expired_with_stock"
    if any(k in q for k in ["expired", "expiration"]):
        return "expired"
    if any(k in q for k in ["expiring", "expire soon", "near expiry", "near expiration", "expiry alert", "urgent expiry"]):
        return "expiring"
    if any(k in q for k in ["out of stock", "zero stock", "no stock"]):
        return "out_of_stock"
    if any(k in q for k in ["available", "usable", "with stock", "still have stock", "current stock"]):
        return "available"
    if any(k in q for k in ["supplier", "suppliers", "by supplier"]):
        return "supplier_summary"
    if any(k in q for k in ["category", "by category"]):
        return "category_summary"
    return "search"


def _ai_apply_entities(df, entities):
    """Apply product/supplier/category/batch filters only."""
    if df is None or df.empty:
        return pd.DataFrame()

    result = df.copy()

    if entities.get("product") and "product_name" in result.columns:
        product = _ai_normalize(entities["product"])
        result = result[result["product_name"].fillna("").astype(str).str.lower().str.contains(re.escape(product), na=False)]

    if entities.get("supplier") and "supplier" in result.columns:
        supplier = _ai_normalize(entities["supplier"])
        result = result[result["supplier"].fillna("").astype(str).str.lower().str.contains(re.escape(supplier), na=False)]

    if entities.get("category") and "category" in result.columns:
        category = _ai_normalize(entities["category"])
        result = result[result["category"].fillna("").astype(str).str.lower().str.contains(re.escape(category), na=False)]

    if entities.get("batch") and "batch_number" in result.columns:
        batch = _ai_normalize(entities["batch"])
        result = result[result["batch_number"].fillna("").astype(str).str.lower().str.contains(re.escape(batch), na=False)]

    return result


def _ai_keyword_search(df, question):
    """Fallback keyword search across important columns."""
    if df is None or df.empty:
        return pd.DataFrame()

    tokens = _ai_extract_keywords(question)
    if not tokens:
        return pd.DataFrame()

    search_cols = [
        c for c in ["product_name", "batch_number", "supplier", "category", "active_ingredient", "storage_location", "remarks"]
        if c in df.columns
    ]
    mask = pd.Series(False, index=df.index)

    for token in tokens:
        token_mask = pd.Series(False, index=df.index)
        for col in search_cols:
            token_mask = token_mask | df[col].fillna("").astype(str).str.lower().str.contains(re.escape(token), na=False)
        mask = mask | token_mask

    return df[mask].copy()


def _ai_key_findings(result_df):
    if result_df is None or result_df.empty:
        return [
            "No matching records were found.",
            "Try using a product name, supplier name, category, or batch number.",
        ]

    findings = []
    qty = pd.to_numeric(result_df.get("current_inventory", 0), errors="coerce").fillna(0).sum() if "current_inventory" in result_df.columns else None
    if qty is not None:
        findings.append(f"Total current quantity in matched records: {clean_number(qty)}")

    if "product_name" in result_df.columns:
        findings.append(f"Unique products: {clean_number(result_df['product_name'].fillna('').astype(str).str.upper().nunique())}")

    if "batch_number" in result_df.columns:
        non_blank_batches = result_df["batch_number"].fillna("").astype(str).str.strip()
        findings.append(f"Batch records: {clean_number(len(result_df))}")

    if "supplier" in result_df.columns and result_df["supplier"].notna().any():
        top_suppliers = result_df["supplier"].fillna("No Supplier").value_counts().head(3)
        if not top_suppliers.empty:
            findings.append("Top supplier(s): " + ", ".join([f"{idx} ({val})" for idx, val in top_suppliers.items()]))

    if "expiry_date_dt" in result_df.columns:
        valid_exp = result_df[pd.notna(result_df["expiry_date_dt"])].copy()
        if not valid_exp.empty:
            nearest = valid_exp.sort_values("expiry_date_dt").iloc[0]
            findings.append(f"Nearest expiry: {_ai_text(nearest.get('expiry_date'))} — {_ai_text(nearest.get('product_name'))}")

    return findings


def _ai_suggest_action(intent, result_df, entities=None):
    entities = entities or {}
    if result_df is None or result_df.empty:
        return "Review your spelling or search using a known product, batch number, supplier, or category."

    if intent in ["expired", "expired_with_stock"]:
        with_stock = result_df[pd.to_numeric(result_df.get("current_inventory", 0), errors="coerce").fillna(0) > 0] if "current_inventory" in result_df.columns else result_df
        if len(with_stock) > 0:
            return "Review expired records with remaining stock for disposal, documentation, or management approval. Do not release expired samples without approval."
        return "Keep the records for documentation, but no stock movement is needed if current quantity is already zero."

    if intent == "expiring":
        return "Prioritize these products for FEFO release, trial use, or customer sample planning before they expire."

    if intent == "available":
        return "Use FEFO when releasing these samples. Check nearest expiry first before selecting a batch."

    if intent == "out_of_stock":
        return "Review if these products need replenishment, supplier follow-up, or cleanup from active monitoring lists."

    if intent == "fefo":
        return "Release or use the earliest valid expiry date first. Avoid expired batches unless approved for documentation-only purposes."

    if intent in ["received_this_month", "released_this_month"]:
        return "Use this movement summary for monthly reporting and reconciliation."

    if intent in ["supplier_summary", "category_summary"]:
        return "Review groups with high expired or near-expiry counts for focused action."

    return "Review the related records and use the filters/search table for more detailed checking."


def _ai_format_answer(title, direct_answer, findings, action, detected=None):
    detected = detected or {}
    findings_html = "".join([f"<li>{f}</li>" for f in findings]) if findings else "<li>No key findings available.</li>"
    detected_items = []
    for k, v in detected.items():
        if v:
            detected_items.append(f"<span class='ai-detected-pill'>{k.title()}: <b>{v}</b></span>")
    detected_html = "".join(detected_items)

    if detected_html:
        detected_html = f"<div class='ai-detected-row'>{detected_html}</div>"

    return f"""
    <div class="ai-answer-card">
        <div class="ai-answer-title">{title}</div>
        {detected_html}
        <div class="ai-answer-section">
            <div class="ai-answer-label">Direct Answer</div>
            <div class="ai-answer-text">{direct_answer}</div>
        </div>
        <div class="ai-answer-section">
            <div class="ai-answer-label">Key Findings</div>
            <ul class="ai-answer-list">{findings_html}</ul>
        </div>
        <div class="ai-action-box">
            <b>Suggested Action:</b> {action}
        </div>
    </div>
    """


def smart_ai_inventory_answer(question):
    """Return title, HTML answer, related dataframe, and detected intent/entities."""
    q_original = str(question or "").strip()
    df = products_base_df()

    if df is None or df.empty:
        answer_html = _ai_format_answer(
            "Inventory Assistant",
            "I could not find product records yet. Please add product samples first.",
            ["No product/batch records are available."],
            "Add product sample records, then ask the assistant again.",
            {}
        )
        return "Inventory Assistant", answer_html, pd.DataFrame(), {"intent": "no_data"}

    intent = _ai_detect_intent(q_original)
    entities = _ai_detect_entities(q_original, df)
    filtered = _ai_apply_entities(df, entities)
    if filtered.empty and any(k in entities for k in ["product", "supplier", "category", "batch"]):
        filtered = pd.DataFrame()

    working_df = filtered if not filtered.empty else df.copy()
    detected = {"intent": intent}
    detected.update(entities)

    title = "Inventory Assistant"
    direct = ""
    result_df = pd.DataFrame()

    if intent == "management_summary":
        summary = _ai_status_summary(df)
        expiring_df = df[df["final_status"].isin(["Expiring Soon", "Urgent Expiry"])].sort_values("days_to_expiry")
        expired_with_stock = df[(df["expiry_status"] == "Expired") & (pd.to_numeric(df["current_inventory"], errors="coerce").fillna(0) > 0)]
        result_df = pd.concat([expiring_df.head(40), expired_with_stock.head(40)], ignore_index=True).drop_duplicates() if (not expiring_df.empty or not expired_with_stock.empty) else df.head(80)
        direct = (
            f"As of today, there are {clean_number(summary['total_records'])} product/batch records, "
            f"{clean_number(summary['usable'])} usable records, "
            f"{clean_number(summary['expiring'] + summary['urgent'])} records needing expiry attention, "
            f"{clean_number(summary['expired'])} expired records, and "
            f"{clean_number(summary['out'])} out-of-stock records."
        )
        findings = [
            f"Total current quantity: {clean_number(summary['total_qty'])}",
            f"Expired records with remaining stock: {clean_number(summary['expired_with_stock'])}",
            f"Suppliers represented: {clean_number(summary['suppliers'])}",
            f"Categories represented: {clean_number(summary['categories'])}",
        ]
        action = _ai_suggest_action(intent, result_df, entities)
        return title, _ai_format_answer(title, direct, findings, action, detected), _ai_detail_table(result_df), detected

    if intent == "fefo":
        result_df = working_df[
            (pd.to_numeric(working_df["current_inventory"], errors="coerce").fillna(0) > 0)
            & (working_df["expiry_status"] != "Expired")
        ].sort_values("expiry_date_dt").head(50)
        direct = f"I listed {clean_number(len(result_df))} FEFO guidance records based on available stock and nearest valid expiry date."
        findings = _ai_key_findings(result_df)
        action = _ai_suggest_action(intent, result_df, entities)
        return "FEFO Priority Recommendation", _ai_format_answer("FEFO Priority Recommendation", direct, findings, action, detected), _ai_detail_table(result_df), detected

    if intent == "received_this_month":
        result_df = _ai_stock_in_month_df()
        total_qty = pd.to_numeric(result_df.get("Quantity", 0), errors="coerce").fillna(0).sum() if not result_df.empty else 0
        direct = f"There are {clean_number(len(result_df))} Sample-In transactions this month with total received quantity of {clean_number(total_qty)}."
        findings = [
            f"Receiving transactions this month: {clean_number(len(result_df))}",
            f"Total quantity received: {clean_number(total_qty)}",
        ]
        if not result_df.empty and "Supplier" in result_df.columns:
            top_supplier = result_df["Supplier"].fillna("Unknown Supplier").value_counts().head(1)
            if not top_supplier.empty:
                findings.append(f"Top supplier by transaction count: {top_supplier.index[0]} ({top_supplier.iloc[0]})")
        action = _ai_suggest_action(intent, result_df, entities)
        return "Sample-In This Month", _ai_format_answer("Sample-In This Month", direct, findings, action, detected), result_df, detected

    if intent == "released_this_month":
        result_df = _ai_stock_out_month_df()
        total_qty = pd.to_numeric(result_df.get("Quantity", 0), errors="coerce").fillna(0).sum() if not result_df.empty else 0
        direct = f"There are {clean_number(len(result_df))} Sample-Out transactions this month with total released quantity of {clean_number(total_qty)}."
        findings = [
            f"Released transactions this month: {clean_number(len(result_df))}",
            f"Total quantity released: {clean_number(total_qty)}",
        ]
        action = _ai_suggest_action(intent, result_df, entities)
        return "Sample-Out This Month", _ai_format_answer("Sample-Out This Month", direct, findings, action, detected), result_df, detected

    if intent == "expired_with_stock":
        result_df = working_df[
            (working_df["expiry_status"] == "Expired")
            & (pd.to_numeric(working_df["current_inventory"], errors="coerce").fillna(0) > 0)
        ].sort_values("expiry_date_dt")
        direct = f"I found {clean_number(len(result_df))} expired product/batch records that still have remaining stock."
        title = "Expired Products with Stock"
    elif intent == "expired":
        result_df = working_df[working_df["expiry_status"] == "Expired"].sort_values("expiry_date_dt")
        direct = f"I found {clean_number(len(result_df))} expired product/batch records."
        title = "Expired Product Records"
    elif intent == "expiring":
        result_df = working_df[
            (pd.to_numeric(working_df["current_inventory"], errors="coerce").fillna(0) > 0)
            & (working_df["final_status"].isin(["Expiring Soon", "Urgent Expiry"]))
        ].sort_values("days_to_expiry")
        direct = f"I found {clean_number(len(result_df))} expiring or urgent-expiry records with available stock."
        title = "Expiring Soon Products"
    elif intent == "out_of_stock":
        result_df = working_df[pd.to_numeric(working_df["current_inventory"], errors="coerce").fillna(0) <= 0].sort_values(["product_name", "batch_number"])
        direct = f"I found {clean_number(len(result_df))} out-of-stock product/batch records."
        title = "Out-of-Stock Products"
    elif intent == "available":
        result_df = working_df[
            (pd.to_numeric(working_df["current_inventory"], errors="coerce").fillna(0) > 0)
            & (working_df["expiry_status"] != "Expired")
        ].sort_values(["product_name", "expiry_date_dt"])
        direct = f"I found {clean_number(len(result_df))} usable product/batch records with stock and not expired."
        title = "Available Products"
    elif intent == "supplier_summary":
        if not filtered.empty:
            result_df = filtered.sort_values(["supplier", "product_name", "batch_number"])
            direct = f"I found {clean_number(len(result_df))} records matching the detected supplier/product/category filters."
            title = "Supplier/Product Search Result"
        else:
            result_df = (
                df.groupby("supplier", dropna=False)
                .agg(
                    Records=("product_name", "count"),
                    Available_Records=("current_inventory", lambda s: int((pd.to_numeric(s, errors="coerce").fillna(0) > 0).sum())),
                    Total_Current_Qty=("current_inventory", "sum"),
                )
                .reset_index()
                .rename(columns={"supplier": "Supplier"})
                .sort_values("Records", ascending=False)
            )
            direct = f"I summarized inventory by supplier. There are {clean_number(len(result_df))} suppliers in the current records."
            findings = [
                f"Supplier groups: {clean_number(len(result_df))}",
                f"Total records summarized: {clean_number(result_df['Records'].sum()) if not result_df.empty else '0'}",
            ]
            action = _ai_suggest_action(intent, result_df, entities)
            return "Supplier Summary", _ai_format_answer("Supplier Summary", direct, findings, action, detected), result_df, detected
    elif intent == "category_summary":
        if not filtered.empty:
            result_df = filtered.sort_values(["category", "product_name", "batch_number"])
            direct = f"I found {clean_number(len(result_df))} records matching the detected category/product/supplier filters."
            title = "Category/Product Search Result"
        else:
            result_df = (
                df.groupby("category", dropna=False)
                .agg(
                    Records=("product_name", "count"),
                    Available_Records=("current_inventory", lambda s: int((pd.to_numeric(s, errors="coerce").fillna(0) > 0).sum())),
                    Total_Current_Qty=("current_inventory", "sum"),
                )
                .reset_index()
                .rename(columns={"category": "Category"})
                .sort_values("Records", ascending=False)
            )
            direct = f"I summarized inventory by category. There are {clean_number(len(result_df))} categories in the current records."
            findings = [
                f"Category groups: {clean_number(len(result_df))}",
                f"Total records summarized: {clean_number(result_df['Records'].sum()) if not result_df.empty else '0'}",
            ]
            action = _ai_suggest_action(intent, result_df, entities)
            return "Category Summary", _ai_format_answer("Category Summary", direct, findings, action, detected), result_df, detected
    else:
        result_df = filtered if not filtered.empty else _ai_keyword_search(df, q_original)
        if result_df.empty:
            direct = "I could not find matching records for that question."
            findings = ["Try using a product name, supplier name, category, batch number, or a status such as available, expired, or expiring soon."]
            action = "Example: “Show available Fytofert products” or “What products are expiring soon?”"
            return "Search Result", _ai_format_answer("Search Result", direct, findings, action, detected), pd.DataFrame(), detected
        result_df = result_df.sort_values(["product_name", "batch_number", "expiry_date_dt"])
        direct = f"I found {clean_number(len(result_df))} product/batch records related to your question."
        title = "Search Result"

    findings = _ai_key_findings(result_df)
    action = _ai_suggest_action(intent, result_df, entities)
    return title, _ai_format_answer(title, direct, findings, action, detected), _ai_detail_table(result_df), detected


def ai_assistant_page():
    page_header("AI Assistant", "Ask inventory questions, get key findings, and receive action suggestions without using an external API.")
    module_hero(
        "Smart Offline Inventory Assistant",
        "Detects products, suppliers, categories, batches, status intent, expiry risks, FEFO priorities, and monthly sample movements.",
        "🤖",
        "Version 1.2.1"
    )

    st.markdown(
        """
        <style>
        .ai-answer-card {
            background: linear-gradient(180deg, rgba(255,255,255,.96), rgba(250,253,251,.96));
            border: 1px solid #dceee4;
            border-radius: 24px;
            padding: 1.2rem 1.25rem;
            box-shadow: 0 14px 34px rgba(16,24,40,.07);
            margin: .85rem 0 1rem 0;
        }
        .ai-answer-title {
            color: #17312a;
            font-size: 1.25rem;
            font-weight: 950;
            margin-bottom: .7rem;
        }
        .ai-answer-label {
            color: #166534;
            font-size: .8rem;
            font-weight: 900;
            text-transform: uppercase;
            letter-spacing: .04em;
            margin-bottom: .25rem;
        }
        .ai-answer-text {
            color: #475467;
            line-height: 1.55;
            font-weight: 560;
        }
        .ai-answer-section {
            margin-top: .8rem;
        }
        .ai-answer-list {
            color: #475467;
            margin-top: .2rem;
            margin-bottom: .2rem;
        }
        .ai-action-box {
            margin-top: 1rem;
            padding: .85rem .95rem;
            border-radius: 16px;
            background: #f0fdf4;
            border: 1px solid #bbf7d0;
            color: #14532d;
            line-height: 1.45;
        }
        .ai-detected-row {
            display: flex;
            gap: .4rem;
            flex-wrap: wrap;
            margin-bottom: .75rem;
        }
        .ai-detected-pill {
            display: inline-flex;
            border-radius: 999px;
            padding: .32rem .55rem;
            background: #f4fbf6;
            border: 1px solid #dceee4;
            color: #166534;
            font-size: .78rem;
            font-weight: 750;
        }
        </style>
        """,
        unsafe_allow_html=True,
    )

    section_shell(
        "Ask the Assistant",
        "This smarter offline assistant uses your local app database only. No API key is required and no company data is sent outside the app."
    )

    suggested_questions = [
        "What products are expiring soon?",
        "Show expired products with remaining stock.",
        "Which products are available?",
        "Show out of stock products.",
        "Show CERADIS fertilizer products with stock.",
        "Which Fytofert batches are still available?",
        "What samples were received this month?",
        "What samples were released this month?",
        "Summarize inventory by supplier.",
        "Summarize inventory by category.",
    ]

    c1, c2 = st.columns([1.05, 2])
    with c1:
        selected_prompt = st.selectbox("Suggested questions", [""] + suggested_questions)
    with c2:
        question = st.text_input(
            "Ask a question",
            value=selected_prompt,
            placeholder="Example: Show CERADIS fertilizer products with stock",
        )

    ask_clicked = st.button("Ask Assistant", use_container_width=True)

    final_question = question

    if ask_clicked and str(final_question or "").strip():
        title, answer_html, related_df, detected = smart_ai_inventory_answer(final_question)

        st.markdown(answer_html, unsafe_allow_html=True)

        if related_df is not None and not related_df.empty:
            st.markdown("<div class='section-title'>Related Records</div>", unsafe_allow_html=True)
            display_df(related_df, height=430, table_key="smart_ai_related_records")
            st.download_button(
                "Download Related Records as CSV",
                to_csv_bytes(related_df),
                file_name="smart_ai_related_records.csv",
                mime="text/csv",
                use_container_width=True,
            )
        else:
            st.info("No related records table for this answer.")

        st.session_state.setdefault("ai_assistant_history", [])
        plain_answer = re.sub("<[^<]+?>", " ", answer_html)
        plain_answer = re.sub(r"\s+", " ", plain_answer).strip()
        st.session_state["ai_assistant_history"].insert(0, {"question": final_question, "answer": plain_answer[:380]})
        st.session_state["ai_assistant_history"] = st.session_state["ai_assistant_history"][:8]

    if st.session_state.get("ai_assistant_history"):
        with st.expander("Recent Assistant Questions", expanded=False):
            for item in st.session_state["ai_assistant_history"]:
                st.markdown(f"**Q:** {item['question']}")
                st.caption(item["answer"])

    end_section_shell()

    section_shell("What the Smart Assistant Can Understand", "Version 1.2 recognizes intent and common inventory entities.")
    col_a, col_b, col_c = st.columns(3)
    with col_a:
        st.markdown("""
        **Detected Entities**
        - Product names
        - Supplier names
        - Categories
        - Batch numbers
        """)
    with col_b:
        st.markdown("""
        **Detected Intent**
        - Available stock
        - Expiring soon
        - Expired with stock
        - FEFO guidance
        """)
    with col_c:
        st.markdown("""
        **Output**
        - Direct answer
        - Key findings
        - Suggested action
        - Related records table
        """)
    end_section_shell()


# -----------------------------
# AI ASSISTANT ACCURACY OVERRIDE - VERSION 1.2.2
# This overrides the earlier AI Assistant functions with more accurate filters.
# -----------------------------

def _ai22_norm(value):
    text = "" if value is None else str(value)
    return re.sub(r"\s+", " ", text.strip().lower())


def _ai22_contains(series, needle):
    needle = _ai22_norm(needle)
    if not needle:
        return pd.Series(False, index=series.index)
    return series.fillna("").astype(str).str.lower().str.contains(re.escape(needle), na=False)


def _ai22_status_intent(question):
    q = _ai22_norm(question)

    if any(x in q for x in ["received this month", "sample in this month", "stock in this month", "received samples this month"]):
        return "received_this_month"

    if any(x in q for x in ["released this month", "sample out this month", "stock out this month", "released samples this month"]):
        return "released_this_month"

    if any(x in q for x in ["expired with stock", "expired still have stock", "expired products still", "expired batches still", "expired with remaining"]):
        return "expired_with_stock"

    if any(x in q for x in ["expire soon", "expiring soon", "near expiry", "near expiration", "urgent expiry", "expiry alert"]):
        return "expiring"

    if any(x in q for x in ["expired", "expiration"]):
        return "expired"

    if any(x in q for x in ["out of stock", "zero stock", "no stock", "without stock"]):
        return "out_of_stock"

    if any(x in q for x in ["available", "usable", "with stock", "still have stock", "has stock", "current stock", "in stock"]):
        return "available"

    if any(x in q for x in ["fefo", "release first", "use first", "prioritize", "priority"]):
        return "fefo"

    if any(x in q for x in ["supplier summary", "by supplier", "summarize by supplier"]):
        return "supplier_summary"

    if any(x in q for x in ["category summary", "by category", "summarize by category"]):
        return "category_summary"

    if any(x in q for x in ["summary", "overview", "inventory status", "status report"]):
        return "summary"

    return "search"


def _ai22_known_values(df, column):
    if df is None or df.empty or column not in df.columns:
        return []
    values = []
    for value in df[column].dropna().astype(str).unique().tolist():
        value = value.strip()
        if value:
            values.append(value)
    return sorted(values, key=len, reverse=True)


def _ai22_exact_or_partial_match(question, values, min_len=3):
    """Return a known database value found in the question.

    Uses exact phrase first, then meaningful token overlap.
    """
    q = _ai22_norm(question)
    if not q:
        return None

    # Exact phrase match first.
    for value in values:
        norm = _ai22_norm(value)
        if len(norm) >= min_len and norm in q:
            return value

    # Token overlap fallback.
    q_tokens = set(re.findall(r"[a-zA-Z0-9]+", q))
    best = None
    best_score = 0
    for value in values:
        norm = _ai22_norm(value)
        tokens = [t for t in re.findall(r"[a-zA-Z0-9]+", norm) if len(t) >= min_len]
        if not tokens:
            continue

        overlap = [t for t in tokens if t in q_tokens]
        if overlap:
            # Favor shorter exact-looking matches for supplier/category, longer for product.
            score = len(overlap) * 100 + sum(len(t) for t in overlap)
            if score > best_score:
                best = value
                best_score = score

    return best


def _ai22_extract_possible_product_terms(question, df, known_filters):
    """Extract search words likely referring to product, active ingredient, or batch.

    This avoids using generic words like available, expired, products, stock, supplier.
    """
    q = _ai22_norm(question)
    stop_words = {
        "what", "which", "show", "give", "list", "products", "product", "batch", "batches",
        "sample", "samples", "stock", "available", "expired", "expiring", "soon", "details",
        "of", "the", "are", "is", "in", "from", "with", "qty", "quantity", "current", "how",
        "many", "supplier", "suppliers", "records", "record", "status", "inventory", "please",
        "can", "you", "me", "for", "and", "or", "still", "have", "has", "this", "that",
        "summarize", "summary", "report", "release", "released", "received", "month",
        "near", "expiry", "expiration", "urgent", "usable", "zero", "without", "first",
        "using", "priority", "prioritize", "category", "categories"
    }

    # Remove already detected supplier/category/batch phrases to avoid double filtering.
    remove_terms = []
    for key in ["supplier", "category", "batch"]:
        if known_filters.get(key):
            remove_terms.append(_ai22_norm(known_filters[key]))
    q_clean = q
    for term in remove_terms:
        if term:
            q_clean = q_clean.replace(term, " ")

    tokens = [
        t for t in re.findall(r"[a-zA-Z0-9\-]+", q_clean)
        if len(t) >= 3 and t not in stop_words
    ]

    # Keep only tokens that appear in product names, active ingredients, or batch numbers.
    if df is None or df.empty or not tokens:
        return []

    searchable_text = ""
    for col in ["product_name", "active_ingredient", "batch_number"]:
        if col in df.columns:
            searchable_text += " " + " ".join(df[col].fillna("").astype(str).str.lower().tolist())

    useful = []
    for token in tokens:
        if token.lower() in searchable_text:
            useful.append(token)

    return useful


def _ai22_detect_filters(question, df):
    filters = {}

    supplier = _ai22_exact_or_partial_match(question, _ai22_known_values(df, "supplier"))
    category = _ai22_exact_or_partial_match(question, _ai22_known_values(df, "category"))
    batch = _ai22_exact_or_partial_match(question, _ai22_known_values(df, "batch_number"), min_len=3)

    if supplier:
        filters["supplier"] = supplier
    if category:
        filters["category"] = category
    if batch:
        filters["batch"] = batch

    # Product phrase exact/partial matching.
    product = _ai22_exact_or_partial_match(question, _ai22_known_values(df, "product_name"))
    if product:
        filters["product"] = product

    # Product terms allow broader search, e.g. "Fytofert" matches both FYTOFERT PHOSCO and FYTOFERT second source.
    product_terms = _ai22_extract_possible_product_terms(question, df, filters)
    if product_terms:
        filters["product_terms"] = product_terms

    return filters


def _ai22_apply_filters(df, filters):
    if df is None or df.empty:
        return pd.DataFrame()

    result = df.copy()

    if filters.get("supplier") and "supplier" in result.columns:
        result = result[_ai22_contains(result["supplier"], filters["supplier"])]

    if filters.get("category") and "category" in result.columns:
        result = result[_ai22_contains(result["category"], filters["category"])]

    if filters.get("batch") and "batch_number" in result.columns:
        result = result[_ai22_contains(result["batch_number"], filters["batch"])]

    # Product exact match is used only if no broad product terms are detected.
    # This prevents "Fytofert" from being narrowed too much to only FYTOFERT PHOSCO.
    if filters.get("product_terms"):
        mask = pd.Series(False, index=result.index)
        for term in filters["product_terms"]:
            for col in ["product_name", "active_ingredient", "batch_number"]:
                if col in result.columns:
                    mask = mask | _ai22_contains(result[col], term)
        result = result[mask]
    elif filters.get("product") and "product_name" in result.columns:
        result = result[_ai22_contains(result["product_name"], filters["product"])]

    return result


def _ai22_apply_intent_filter(df, intent):
    if df is None or df.empty:
        return pd.DataFrame()

    result = df.copy()
    qty = pd.to_numeric(result.get("current_inventory", result.get("quantity", 0)), errors="coerce").fillna(0)

    if intent == "available":
        return result[(qty > 0) & (result["expiry_status"] != "Expired")].copy()

    if intent == "expiring":
        return result[(qty > 0) & (result["final_status"].isin(["Expiring Soon", "Urgent Expiry"]))].sort_values("days_to_expiry").copy()

    if intent == "expired":
        return result[result["expiry_status"] == "Expired"].sort_values("expiry_date_dt").copy()

    if intent == "expired_with_stock":
        return result[(result["expiry_status"] == "Expired") & (qty > 0)].sort_values("expiry_date_dt").copy()

    if intent == "out_of_stock":
        return result[qty <= 0].sort_values(["product_name", "batch_number"]).copy()

    if intent == "fefo":
        return result[(qty > 0) & (result["expiry_status"] != "Expired")].sort_values("expiry_date_dt").copy()

    return result.copy()


def _ai22_detail_table(df, max_rows=150):
    if df is None or df.empty:
        return pd.DataFrame()

    temp = df.copy()
    preferred = [
        "product_name", "batch_number", "supplier", "category", "active_ingredient",
        "current_inventory", "pack_size", "pack_type", "expiry_date", "days_to_expiry",
        "expiry_status", "final_status", "storage_location", "remarks"
    ]
    cols = [c for c in preferred if c in temp.columns]
    temp = temp[cols].head(max_rows).copy()

    rename_map = {
        "product_name": "Product",
        "batch_number": "Batch No.",
        "supplier": "Supplier",
        "category": "Category",
        "active_ingredient": "Active Ingredient",
        "current_inventory": "Current Qty",
        "pack_size": "Pack Size",
        "pack_type": "Pack Type",
        "expiry_date": "Expiry Date",
        "days_to_expiry": "Days to Expiry",
        "expiry_status": "Expiry Status",
        "final_status": "Status",
        "storage_location": "Storage Location",
        "remarks": "Remarks",
    }
    return temp.rename(columns=rename_map)


def _ai22_month_in():
    today = date.today()
    month_start = date(today.year, today.month, 1)
    return query_df(
        """
        SELECT
            si.date_received AS "Date Received",
            si.reference_no AS "Reference No.",
            COALESCE(p.product_name, '') AS "Product",
            COALESCE(p.batch_number, si.batch_number, '') AS "Batch No.",
            COALESCE(s.supplier_name, 'Unknown Supplier') AS "Supplier",
            si.quantity AS "Quantity",
            CASE
                WHEN p.pack_size_value IS NOT NULL AND p.pack_size_unit IS NOT NULL
                THEN printf('%g %s', p.pack_size_value, p.pack_size_unit)
                ELSE ''
            END AS "Pack Size",
            COALESCE(p.pack_type, '') AS "Pack Type",
            COALESCE(si.received_by, '') AS "Received By",
            COALESCE(si.remarks, '') AS "Remarks"
        FROM stock_in si
        LEFT JOIN products p ON si.product_id = p.product_id
        LEFT JOIN suppliers s ON si.supplier_id = s.supplier_id
        WHERE date(si.date_received) >= date(?)
        ORDER BY date(si.date_received) DESC, si.created_at DESC
        LIMIT 150
        """,
        (str(month_start),),
    )


def _ai22_month_out():
    today = date.today()
    month_start = date(today.year, today.month, 1)
    return query_df(
        """
        SELECT
            so.date_released AS "Date Released",
            so.reference_no AS "Reference No.",
            COALESCE(p.product_name, '') AS "Product",
            COALESCE(p.batch_number, '') AS "Batch No.",
            COALESCE(s.supplier_name, 'Unknown Supplier') AS "Supplier",
            so.quantity AS "Quantity",
            COALESCE(p.pack_type, '') AS "Pack Type",
            COALESCE(so.purpose, '') AS "Purpose",
            COALESCE(so.requested_by, '') AS "Requested By",
            COALESCE(so.released_by, '') AS "Released By",
            COALESCE(so.status, '') AS "Status",
            COALESCE(so.remarks, '') AS "Remarks"
        FROM stock_out so
        LEFT JOIN products p ON so.product_id = p.product_id
        LEFT JOIN suppliers s ON p.supplier_id = s.supplier_id
        WHERE date(so.date_released) >= date(?)
        ORDER BY date(so.date_released) DESC, so.created_at DESC
        LIMIT 150
        """,
        (str(month_start),),
    )


def _ai22_findings(df):
    if df is None or df.empty:
        return ["No matching records found."]

    findings = []
    if "current_inventory" in df.columns:
        qty = pd.to_numeric(df["current_inventory"], errors="coerce").fillna(0).sum()
        findings.append(f"Total current quantity: {clean_number(qty)}")

    if "product_name" in df.columns:
        findings.append(f"Unique products: {clean_number(df['product_name'].fillna('').astype(str).str.upper().nunique())}")

    findings.append(f"Batch/record count: {clean_number(len(df))}")

    if "supplier" in df.columns and not df["supplier"].dropna().empty:
        top = df["supplier"].fillna("No Supplier").value_counts().head(3)
        findings.append("Top supplier(s): " + ", ".join([f"{idx} ({val})" for idx, val in top.items()]))

    if "expiry_date_dt" in df.columns:
        valid = df[pd.notna(df["expiry_date_dt"])].copy()
        if not valid.empty:
            nearest = valid.sort_values("expiry_date_dt").iloc[0]
            findings.append(f"Nearest expiry: {nearest.get('expiry_date', '')} — {nearest.get('product_name', '')}")

    return findings


def _ai22_action(intent, df):
    if df is None or df.empty:
        return "Try checking the spelling or use a known product name, supplier, category, or batch number."

    if intent == "available":
        return "Use FEFO when releasing these samples. Check the nearest expiry date first."
    if intent == "expiring":
        return "Prioritize these products for trial use, customer sample release, or internal use before they expire."
    if intent in ["expired", "expired_with_stock"]:
        return "Review these records for disposal, documentation, or management approval. Do not release expired samples without approval."
    if intent == "out_of_stock":
        return "Check if these products should be replenished or removed from active monitoring."
    if intent == "fefo":
        return "Release or use the earliest valid expiry date first."
    if intent in ["received_this_month", "released_this_month"]:
        return "Use this table for monthly movement reconciliation and reporting."
    return "Review the related records and use the table search/export for follow-up."


def _ai22_filter_summary(filters, intent):
    chips = [f"Intent: {intent}"]
    for k, v in filters.items():
        if k == "product_terms":
            chips.append("Product keywords: " + ", ".join(v))
        else:
            chips.append(f"{k.title()}: {v}")
    return " | ".join(chips)


def _ai22_answer_html(title, direct, findings, action, filter_summary):
    findings_html = "".join([f"<li>{f}</li>" for f in findings])
    return f"""
    <div class="ai-answer-card">
        <div class="ai-answer-title">{title}</div>
        <div class="ai-detected-row">
            <span class="ai-detected-pill">{filter_summary}</span>
        </div>
        <div class="ai-answer-section">
            <div class="ai-answer-label">Direct Answer</div>
            <div class="ai-answer-text">{direct}</div>
        </div>
        <div class="ai-answer-section">
            <div class="ai-answer-label">Key Findings</div>
            <ul class="ai-answer-list">{findings_html}</ul>
        </div>
        <div class="ai-action-box"><b>Suggested Action:</b> {action}</div>
    </div>
    """


def smart_ai_inventory_answer(question):
    """Improved v1.2.2 offline assistant.

    More accurate because it:
    - detects intent first
    - detects supplier/category/batch/product terms separately
    - applies filters in a deterministic order
    - shows detected filters for verification
    """
    question = str(question or "").strip()
    df = products_base_df()

    if df is None or df.empty:
        html = _ai22_answer_html(
            "AI Assistant",
            "No product records are available yet.",
            ["Please add product/sample records first."],
            "Add inventory records, then ask again.",
            "No data"
        )
        return "AI Assistant", html, pd.DataFrame(), {"intent": "no_data"}

    intent = _ai22_status_intent(question)
    filters = _ai22_detect_filters(question, df)

    # Movement queries use transaction tables.
    if intent == "received_this_month":
        result = _ai22_month_in()
        total_qty = pd.to_numeric(result.get("Quantity", 0), errors="coerce").fillna(0).sum() if not result.empty else 0
        direct = f"I found {clean_number(len(result))} Sample-In transactions this month with total received quantity of {clean_number(total_qty)}."
        html = _ai22_answer_html(
            "Samples Received This Month",
            direct,
            [f"Transaction count: {clean_number(len(result))}", f"Total quantity received: {clean_number(total_qty)}"],
            _ai22_action(intent, result),
            _ai22_filter_summary(filters, intent)
        )
        return "Samples Received This Month", html, result, {"intent": intent, **filters}

    if intent == "released_this_month":
        result = _ai22_month_out()
        total_qty = pd.to_numeric(result.get("Quantity", 0), errors="coerce").fillna(0).sum() if not result.empty else 0
        direct = f"I found {clean_number(len(result))} Sample-Out transactions this month with total released quantity of {clean_number(total_qty)}."
        html = _ai22_answer_html(
            "Samples Released This Month",
            direct,
            [f"Transaction count: {clean_number(len(result))}", f"Total quantity released: {clean_number(total_qty)}"],
            _ai22_action(intent, result),
            _ai22_filter_summary(filters, intent)
        )
        return "Samples Released This Month", html, result, {"intent": intent, **filters}

    # Summary/grouping queries.
    if intent == "summary":
        qty = pd.to_numeric(df.get("current_inventory", 0), errors="coerce").fillna(0)
        expired_with_stock = df[(df["expiry_status"] == "Expired") & (qty > 0)]
        expiring = df[(qty > 0) & (df["final_status"].isin(["Expiring Soon", "Urgent Expiry"]))]
        usable = df[(qty > 0) & (df["expiry_status"] != "Expired")]
        out = df[qty <= 0]
        direct = (
            f"Current inventory has {clean_number(len(df))} product/batch records, "
            f"{clean_number(len(usable))} usable records, "
            f"{clean_number(len(expiring))} expiring/urgent records, "
            f"{clean_number(len(expired_with_stock))} expired records with stock, and "
            f"{clean_number(len(out))} out-of-stock records."
        )
        result = pd.concat([expiring.head(40), expired_with_stock.head(40)], ignore_index=True).drop_duplicates()
        html = _ai22_answer_html(
            "Inventory Summary",
            direct,
            _ai22_findings(df),
            "Focus first on expired records with stock and expiring records that can still be used or released through FEFO.",
            _ai22_filter_summary(filters, intent)
        )
        return "Inventory Summary", html, _ai22_detail_table(result), {"intent": intent, **filters}

    if intent == "supplier_summary":
        result = (
            df.groupby("supplier", dropna=False)
            .agg(
                Records=("product_name", "count"),
                Available_Records=("current_inventory", lambda s: int((pd.to_numeric(s, errors="coerce").fillna(0) > 0).sum())),
                Total_Current_Qty=("current_inventory", "sum"),
            )
            .reset_index()
            .rename(columns={"supplier": "Supplier"})
            .sort_values("Records", ascending=False)
        )
        direct = f"I summarized inventory by supplier. There are {clean_number(len(result))} supplier groups."
        html = _ai22_answer_html(
            "Supplier Summary",
            direct,
            [f"Supplier groups: {clean_number(len(result))}", f"Total records summarized: {clean_number(result['Records'].sum()) if not result.empty else '0'}"],
            "Review suppliers with high expired or near-expiry records for follow-up.",
            _ai22_filter_summary(filters, intent)
        )
        return "Supplier Summary", html, result, {"intent": intent, **filters}

    if intent == "category_summary":
        result = (
            df.groupby("category", dropna=False)
            .agg(
                Records=("product_name", "count"),
                Available_Records=("current_inventory", lambda s: int((pd.to_numeric(s, errors="coerce").fillna(0) > 0).sum())),
                Total_Current_Qty=("current_inventory", "sum"),
            )
            .reset_index()
            .rename(columns={"category": "Category"})
            .sort_values("Records", ascending=False)
        )
        direct = f"I summarized inventory by category. There are {clean_number(len(result))} category groups."
        html = _ai22_answer_html(
            "Category Summary",
            direct,
            [f"Category groups: {clean_number(len(result))}", f"Total records summarized: {clean_number(result['Records'].sum()) if not result.empty else '0'}"],
            "Review categories with high expired or near-expiry records for action planning.",
            _ai22_filter_summary(filters, intent)
        )
        return "Category Summary", html, result, {"intent": intent, **filters}

    # Product/supplier/category/batch filters + status intent.
    working = _ai22_apply_filters(df, filters)
    if working.empty and filters:
        # If the structured filter is too strict, try keyword fallback.
        working = _ai22_apply_filters(df, {k: v for k, v in filters.items() if k == "product_terms"})

    if working.empty and not filters:
        working = df.copy()

    result = _ai22_apply_intent_filter(working, intent)

    # If no status intent and no structured results, fallback keyword search.
    if result.empty and intent == "search":
        result = _ai22_apply_filters(df, filters)
        if result.empty:
            # broad keyword fallback
            tokens = _ai22_extract_possible_product_terms(question, df, {})
            if tokens:
                result = _ai22_apply_filters(df, {"product_terms": tokens})
                filters["product_terms"] = tokens

    title_map = {
        "available": "Available Products",
        "expiring": "Expiring Products",
        "expired": "Expired Products",
        "expired_with_stock": "Expired Products with Stock",
        "out_of_stock": "Out-of-Stock Products",
        "fefo": "FEFO Priority",
        "search": "Search Result",
    }
    title = title_map.get(intent, "AI Assistant")

    if result.empty:
        direct = "I could not find matching records for that question."
    else:
        direct = f"I found {clean_number(len(result))} matching product/batch records."

    # If a status intent was not stated but filters were detected, make the answer clearer.
    if intent == "search" and filters and not result.empty:
        direct = f"I found {clean_number(len(result))} records matching the detected product/supplier/category/batch filters."

    html = _ai22_answer_html(
        title,
        direct,
        _ai22_findings(result),
        _ai22_action(intent, result),
        _ai22_filter_summary(filters, intent)
    )

    return title, html, _ai22_detail_table(result), {"intent": intent, **filters}


def ai_assistant_page():
    page_header("AI Assistant", "Ask inventory questions, get key findings, and receive action suggestions without using an external API.")
    module_hero(
        "Smarter Offline Inventory Assistant",
        "More accurate intent and product/supplier/category/batch detection using your local inventory database only.",
        "🤖",
        "Version 1.2.2"
    )

    st.markdown(
        """
        <style>
        .ai-answer-card {
            background: linear-gradient(180deg, rgba(255,255,255,.96), rgba(250,253,251,.96));
            border: 1px solid #dceee4;
            border-radius: 24px;
            padding: 1.2rem 1.25rem;
            box-shadow: 0 14px 34px rgba(16,24,40,.07);
            margin: .85rem 0 1rem 0;
        }
        .ai-answer-title {
            color: #17312a;
            font-size: 1.25rem;
            font-weight: 950;
            margin-bottom: .7rem;
        }
        .ai-answer-label {
            color: #166534;
            font-size: .8rem;
            font-weight: 900;
            text-transform: uppercase;
            letter-spacing: .04em;
            margin-bottom: .25rem;
        }
        .ai-answer-text {
            color: #475467;
            line-height: 1.55;
            font-weight: 560;
        }
        .ai-answer-section {
            margin-top: .8rem;
        }
        .ai-answer-list {
            color: #475467;
            margin-top: .2rem;
            margin-bottom: .2rem;
        }
        .ai-action-box {
            margin-top: 1rem;
            padding: .85rem .95rem;
            border-radius: 16px;
            background: #f0fdf4;
            border: 1px solid #bbf7d0;
            color: #14532d;
            line-height: 1.45;
        }
        .ai-detected-row {
            display: flex;
            gap: .4rem;
            flex-wrap: wrap;
            margin-bottom: .75rem;
        }
        .ai-detected-pill {
            display: inline-flex;
            border-radius: 999px;
            padding: .32rem .55rem;
            background: #f4fbf6;
            border: 1px solid #dceee4;
            color: #166534;
            font-size: .78rem;
            font-weight: 750;
        }
        </style>
        """,
        unsafe_allow_html=True,
    )

    section_shell(
        "Ask the Assistant",
        "This offline assistant uses deterministic database filters. It shows the detected intent and filters so you can verify the answer."
    )

    suggested_questions = [
        "What products are expiring soon?",
        "Show expired products with remaining stock.",
        "Which products are available?",
        "Show out of stock products.",
        "Show CERADIS fertilizer products with stock.",
        "Which Fytofert batches are still available?",
        "What samples were received this month?",
        "What samples were released this month?",
        "Summarize inventory by supplier.",
        "Summarize inventory by category.",
    ]

    c1, c2 = st.columns([1.05, 2])
    with c1:
        selected_prompt = st.selectbox("Suggested questions", [""] + suggested_questions)
    with c2:
        question = st.text_input(
            "Ask a question",
            value=selected_prompt,
            placeholder="Example: Show CERADIS fertilizer products with stock",
        )

    ask_clicked = st.button("Ask Assistant", use_container_width=True)

    if ask_clicked and str(question or "").strip():
        title, answer_html, related_df, detected = smart_ai_inventory_answer(question)

        st.markdown(answer_html, unsafe_allow_html=True)

        if related_df is not None and not related_df.empty:
            st.markdown("<div class='section-title'>Related Records</div>", unsafe_allow_html=True)
            display_df(related_df, height=430, table_key="smart_ai_related_records")
            st.download_button(
                "Download Related Records as CSV",
                to_csv_bytes(related_df),
                file_name="smart_ai_related_records.csv",
                mime="text/csv",
                use_container_width=True,
            )
        else:
            st.info("No related records table for this answer.")

        st.session_state.setdefault("ai_assistant_history", [])
        plain_answer = re.sub("<[^<]+?>", " ", answer_html)
        plain_answer = re.sub(r"\s+", " ", plain_answer).strip()
        st.session_state["ai_assistant_history"].insert(0, {"question": question, "answer": plain_answer[:380]})
        st.session_state["ai_assistant_history"] = st.session_state["ai_assistant_history"][:8]

    if st.session_state.get("ai_assistant_history"):
        with st.expander("Recent Assistant Questions", expanded=False):
            for item in st.session_state["ai_assistant_history"]:
                st.markdown(f"**Q:** {item['question']}")
                st.caption(item["answer"])

    end_section_shell()

    section_shell("Accuracy Notes", "How to get the most accurate answers.")
    st.markdown("""
    The offline assistant is most accurate when your question includes one or more of these:
    - Product name, for example **Fytofert**
    - Supplier name, for example **CERADIS**
    - Category, for example **Fertilizer**
    - Batch number
    - Status, for example **available**, **expired**, **expiring soon**, or **out of stock**
    """)
    end_section_shell()

def auto_backup_database():
    """Create one automatic SQLite backup per day when the app starts."""
    try:
        backup_dir = APP_DIR / "auto_backups"
        backup_dir.mkdir(exist_ok=True)
        today_key = datetime.now().strftime("%Y%m%d")
        existing = list(backup_dir.glob(f"inventory_app_auto_backup_{today_key}_*.db"))
        if DB_PATH.exists() and not existing:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            backup_path = backup_dir / f"inventory_app_auto_backup_{timestamp}.db"
            backup_path.write_bytes(DB_PATH.read_bytes())
        # Keep the most recent 30 auto-backups to avoid unlimited growth.
        backups = sorted(backup_dir.glob("inventory_app_auto_backup_*.db"), key=lambda x: x.stat().st_mtime, reverse=True)
        for old_backup in backups[30:]:
            try:
                old_backup.unlink()
            except Exception:
                pass
    except Exception as exc:
        # Do not stop the app if backup fails; show a compact warning after login.
        st.session_state["auto_backup_warning"] = str(exc)


def get_latest_backup_label():
    backup_dir = APP_DIR / "auto_backups"
    backups = sorted(backup_dir.glob("inventory_app_auto_backup_*.db"), key=lambda x: x.stat().st_mtime, reverse=True)
    if not backups:
        return "No auto-backup yet"
    latest = backups[0]
    return datetime.fromtimestamp(latest.stat().st_mtime).strftime("%Y-%m-%d %I:%M %p")


def login_page():
    """Premium centered login page for Inventory App Version 1.0."""
    st.markdown(
        """
        <style>
        /* ================================
           VERSION 1.0 PREMIUM LOGIN PAGE
           ================================ */

        .stApp {
            background:
                radial-gradient(circle at 14% 10%, rgba(34,197,94,.10) 0 9%, transparent 10%),
                radial-gradient(circle at 88% 18%, rgba(16,185,129,.13) 0 13%, transparent 14%),
                linear-gradient(180deg, #f7fbf8 0%, #eef6f2 100%) !important;
        }

        [class*="st-key-login_card"] {
            margin-top: 5.5vh;
            padding: 0 !important;
        }

        [class*="st-key-login_card"] [data-testid="stVerticalBlock"] {
            gap: 0.72rem !important;
        }

        .login-premium-card {
            width: 100%;
            background:
                radial-gradient(circle at 92% 88%, rgba(34,197,94,.10) 0 18%, transparent 19%),
                linear-gradient(180deg, rgba(255,255,255,.94) 0%, rgba(255,255,255,.86) 100%);
            border: 1px solid rgba(207, 221, 214, .95);
            border-radius: 30px;
            padding: 2.35rem 2.1rem 1.45rem 2.1rem;
            box-shadow:
                0 24px 56px rgba(16,24,40,.10),
                0 2px 8px rgba(16,24,40,.04);
            backdrop-filter: blur(12px);
            position: relative;
            overflow: hidden;
            margin-bottom: 1rem;
        }

        .login-premium-card::before {
            content: "";
            position: absolute;
            inset: 0 0 auto 0;
            height: 7px;
            background: linear-gradient(90deg, #16a34a, #22c55e, #86efac);
        }

        .login-top-row {
            display: flex;
            align-items: center;
            gap: 1.05rem;
            margin-bottom: 1.2rem;
        }

        .login-logo-badge {
            width: 78px;
            height: 78px;
            border-radius: 24px;
            display: grid;
            place-items: center;
            font-size: 2.25rem;
            background:
                radial-gradient(circle at 30% 20%, rgba(255,255,255,.95) 0 22%, transparent 23%),
                linear-gradient(180deg, #effaf3 0%, #ddf3e5 100%);
            border: 1px solid #d6eadb;
            box-shadow: 0 14px 32px rgba(34,197,94,.16);
            flex: 0 0 auto;
        }

        .login-title {
            font-size: 2.2rem;
            line-height: 1.05;
            font-weight: 950;
            color: #17312a;
            letter-spacing: -0.045em;
            margin: 0;
        }

        .login-subtitle {
            color: #667085;
            font-size: 1.02rem;
            line-height: 1.45;
            margin-top: .45rem;
            font-weight: 560;
        }

        .login-mini-row {
            display: flex;
            flex-wrap: wrap;
            gap: .55rem;
            margin-top: 1.05rem;
        }

        .login-pill {
            display: inline-flex;
            align-items: center;
            gap: .35rem;
            padding: .42rem .68rem;
            border-radius: 999px;
            border: 1px solid #dceee4;
            background: #f3fbf6;
            color: #166534;
            font-size: .78rem;
            font-weight: 800;
        }

        [class*="st-key-login_card"] [data-testid="stForm"] {
            background: rgba(255,255,255,.82) !important;
            border: 1px solid rgba(218,226,232,.95) !important;
            border-radius: 26px !important;
            padding: 1.25rem 1.3rem 1.35rem 1.3rem !important;
            box-shadow: 0 16px 36px rgba(16,24,40,.07) !important;
        }

        [class*="st-key-login_card"] label,
        [class*="st-key-login_card"] [data-testid="stWidgetLabel"] {
            color: #344054 !important;
            font-weight: 850 !important;
            font-size: .93rem !important;
        }

        [class*="st-key-login_card"] input {
            min-height: 3.25rem !important;
            border-radius: 16px !important;
            border: 1px solid #d8e1e6 !important;
            background: #fbfdfc !important;
            color: #17312a !important;
            padding-left: 1rem !important;
            font-weight: 650 !important;
            box-shadow: inset 0 1px 0 rgba(255,255,255,.65) !important;
        }

        [class*="st-key-login_card"] input:focus {
            border-color: #22c55e !important;
            box-shadow: 0 0 0 4px rgba(34,197,94,.10) !important;
        }

        [class*="st-key-login_card"] div[data-testid="stFormSubmitButton"] button {
            min-height: 3.25rem !important;
            border-radius: 16px !important;
            border: none !important;
            color: white !important;
            font-weight: 900 !important;
            font-size: 1.02rem !important;
            background: linear-gradient(180deg, #22c55e 0%, #159947 100%) !important;
            box-shadow: 0 16px 30px rgba(34,197,94,.22) !important;
            transition: transform .18s ease, box-shadow .18s ease, filter .18s ease !important;
        }

        [class*="st-key-login_card"] div[data-testid="stFormSubmitButton"] button:hover {
            transform: translateY(-2px) !important;
            filter: brightness(1.03) !important;
            box-shadow: 0 20px 38px rgba(34,197,94,.28) !important;
        }

        [class*="st-key-login_card"] div[data-testid="stFormSubmitButton"] button:active {
            transform: translateY(0px) scale(.992) !important;
        }

        .login-helper {
            margin-top: .85rem;
            color: #71847d;
            font-size: .84rem;
            line-height: 1.45;
            background: rgba(244,250,246,.9);
            border: 1px solid #dceee4;
            border-radius: 16px;
            padding: .85rem .95rem;
        }

        .login-version {
            text-align: center;
            color: #587169;
            font-size: .82rem;
            font-weight: 800;
            margin-top: .9rem;
        }

        @media (max-width: 768px) {
            [class*="st-key-login_card"] {
                margin-top: 1.25rem;
            }
            .login-premium-card {
                border-radius: 22px;
                padding: 1.6rem 1.25rem 1.15rem 1.25rem;
            }
            .login-top-row {
                align-items: flex-start;
                gap: .85rem;
            }
            .login-logo-badge {
                width: 62px;
                height: 62px;
                border-radius: 18px;
                font-size: 1.85rem;
            }
            .login-title {
                font-size: 1.55rem;
            }
            .login-subtitle {
                font-size: .9rem;
            }
        }
        </style>
        """,
        unsafe_allow_html=True,
    )

    left_col, center_col, right_col = st.columns([0.85, 1.18, 0.85])

    with center_col:
        try:
            login_container = st.container(key="login_card")
        except TypeError:
            login_container = st.container()

        with login_container:
            st.markdown(
                f"""
                <div class="login-premium-card">
                    <div class="login-top-row">
                        <div class="login-logo-badge">🌿</div>
                        <div>
                            <div class="login-title">Sample Inventory App</div>
                            <div class="login-subtitle">
                                Secure access to product sample monitoring, stock balances, expiry alerts, and movement records.
                            </div>
                        </div>
                    </div>
                    <div class="login-mini-row">
                        <span class="login-pill">📦 Inventory Monitoring</span>
                        <span class="login-pill">🔐 Secure Login</span>
                        <span class="login-pill">{APP_VERSION}</span>
                    </div>
                </div>
                """,
                unsafe_allow_html=True,
            )

            with st.form("login_form"):
                username = st.text_input("Username", placeholder="Enter admin or viewer")
                password = st.text_input("Password", type="password", placeholder="Enter password")
                submitted = st.form_submit_button("Login", use_container_width=True)

            st.markdown(
                """
                <div class="login-version">Inventory App • Version {APP_VERSION}</div>
                """,
                unsafe_allow_html=True,
            )

    if submitted:
        account = USER_ACCOUNTS.get(username.strip().lower())
        if account and password == account["password"]:
            st.session_state["is_logged_in"] = True
            st.session_state["username"] = username.strip().lower()
            st.session_state["user_role"] = account["role"]
            st.session_state["display_name"] = account["display_name"]
            st.rerun()
        else:
            st.error("Invalid username or password.")


def require_login():
    if not st.session_state.get("is_logged_in"):
        login_page()
        st.stop()


def is_admin():
    return st.session_state.get("user_role") == "Admin"

# APP START
# -----------------------------

init_db()
auto_backup_database()
require_login()

app_name_display = get_setting("app_name", "Sample Inventory App")
company_name_display = get_setting("company_name", "Farmfix")
user_role_display = st.session_state.get("user_role", "Viewer")
user_name_display = st.session_state.get("display_name", "User")
app_logo_uri = image_to_data_uri(APP_LOGO_PATH) or image_to_data_uri(LOGO_PATH)

with st.sidebar:
    st.markdown(
        f"""
        <div class="sidebar-brand">
            <div class="sidebar-logo"><img class="sidebar-logo-img" src="{app_logo_uri}" alt="FarmFix logo" /></div>
            <div>
                <div class="sidebar-title">{app_name_display}</div>
                <div class="sidebar-subtitle">{company_name_display} • Product sample monitoring</div>
                <div class="sidebar-subtitle">{APP_VERSION} • {user_role_display}</div>
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )
    st.markdown("---")
    st.markdown("<div class='sidebar-menu-title'>Navigation</div>", unsafe_allow_html=True)
    if is_admin():
        menu_items = [
            "📊 Dashboard",
            "🤖 AI Assistant",
            "📦 Products",
            "⬇️ Sample In",
            "⬆️ Sample Out",
            "📋 Stock Balance",
            "🏢 Suppliers",
            "📑 Reports",
            "⚙️ Settings",
        ]
    else:
        menu_items = [
            "📊 Dashboard",
            "🤖 AI Assistant",
            "📋 Stock Balance",
            "📑 Reports",
        ]
    page = st.radio("Menu", menu_items, label_visibility="collapsed")
    st.markdown("<div class='sidebar-farm'></div>", unsafe_allow_html=True)
    current_user_sidebar = f"{user_name_display} ({user_role_display})"
    st.markdown(
        f'''<div class="sidebar-info-card">
            <div class="sidebar-info-title">Quick Tip</div>
            <div class="sidebar-info-text">Use Settings to customize chart colors, alert days, dropdown lists, and backup tools. Current user: <b>{current_user_sidebar}</b>.</div>
        </div>''',
        unsafe_allow_html=True,
    )
    st.markdown("---")
    st.caption(f"App version: {APP_VERSION}")
    st.caption(f"Latest auto-backup: {get_latest_backup_label()}")
    if st.session_state.get("auto_backup_warning"):
        st.warning("Auto-backup warning: " + st.session_state.get("auto_backup_warning"))
    if st.button("Logout", use_container_width=True):
        for key in ["is_logged_in", "username", "user_role", "display_name"]:
            st.session_state.pop(key, None)
        st.rerun()
    st.caption("Compact navigation design.")

if page == "📊 Dashboard":
    dashboard_page()
elif page == "🤖 AI Assistant":
    ai_assistant_page()
elif page == "📦 Products":
    if is_admin():
        products_page()
    else:
        st.error("Viewer access is read-only. Please login as Admin to manage products.")
elif page == "⬇️ Sample In":
    if is_admin():
        stock_in_page()
    else:
        st.error("Viewer access is read-only. Please login as Admin to record Sample In transactions.")
elif page == "⬆️ Sample Out":
    if is_admin():
        stock_out_page()
    else:
        st.error("Viewer access is read-only. Please login as Admin to record Sample Out transactions.")
elif page == "📋 Stock Balance":
    stock_balance_page()
elif page == "🏢 Suppliers":
    if is_admin():
        suppliers_page()
    else:
        st.error("Viewer access is read-only. Please login as Admin to manage suppliers.")
elif page == "📑 Reports":
    reports_page()
elif page == "⚙️ Settings":
    if is_admin():
        settings_page()
    else:
        st.error("Only Admin can access Settings.")


st.markdown("""
<style>
/* Collapse-friendly top spacing: keeps content high without forcing the sidebar open */
[data-testid="stAppViewBlockContainer"],
.main .block-container,
.block-container {
    padding-top: 0.55rem !important;
    margin-top: 0rem !important;
}
.dashboard-topbar,
.app-header {
    margin-top: 0rem !important;
}
/* Keep the native Streamlit sidebar show/expand button visible after the sidebar is hidden */
[data-testid="stSidebarCollapsedControl"],
[data-testid="collapsedControl"] {
    display: flex !important;
    visibility: visible !important;
    opacity: 1 !important;
    position: fixed !important;
    top: 0.55rem !important;
    left: 0.65rem !important;
    z-index: 999999 !important;
    background: rgba(255,255,255,0.96) !important;
    border: 1px solid #d0d5dd !important;
    border-radius: 999px !important;
    box-shadow: 0 8px 24px rgba(16,24,40,.14) !important;
    width: 2.15rem !important;
    height: 2.15rem !important;
    align-items: center !important;
    justify-content: center !important;
}
[data-testid="stSidebarCollapsedControl"] button,
[data-testid="collapsedControl"] button {
    display: flex !important;
    visibility: visible !important;
    opacity: 1 !important;
    color: #12352b !important;
}
@media (max-width: 900px) {
    [data-testid="stAppViewBlockContainer"],
    .main .block-container,
    .block-container {
        padding-top: 0.7rem !important;
        padding-left: 0.65rem !important;
        padding-right: 0.65rem !important;
    }
}
</style>
""", unsafe_allow_html=True)


st.markdown("""
<style>
/* Sidebar reopen fix: keep the native Streamlit collapsed-sidebar control accessible. */
header[data-testid="stHeader"],
div[data-testid="stToolbar"] {
    display: flex !important;
    visibility: visible !important;
    opacity: 1 !important;
    pointer-events: auto !important;
    overflow: visible !important;
    z-index: 999999 !important;
}
[data-testid="stSidebarCollapsedControl"],
[data-testid="collapsedControl"],
button[kind="header"] {
    visibility: visible !important;
    opacity: 1 !important;
    pointer-events: auto !important;
}
[data-testid="stSidebarCollapsedControl"],
[data-testid="collapsedControl"] {
    position: fixed !important;
    top: .55rem !important;
    left: .65rem !important;
    z-index: 999999 !important;
}
</style>
""", unsafe_allow_html=True)
